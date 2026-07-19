#!/usr/bin/env python3
"""
BIB reviewer feedback -> advisor-discussion deck (Traditional Chinese).

Purpose: turn the reviewer's comments on the BIB "Problem-solving Protocol"
manuscript into decision-oriented slides to discuss with the advisor. Each
decision slide states: what the reviewer flagged, the options (with trade-offs),
and a recommendation + what the advisor needs to decide.

Reuses the design system from make_paper_designed.py (validated dataviz palette,
dark cover/dividers, card layout, CJK fonts). Chinese only, per request.

Output: bib_review_discussion_zh.pptx
"""

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from matplotlib.patches import FancyBboxPatch, Rectangle
import numpy as np
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DOCS = Path(__file__).parent
TMP = DOCS / "bib_review_figs"
TMP.mkdir(exist_ok=True)
OUT = DOCS / "bib_review_discussion_zh.pptx"

# ---- design tokens (from make_paper_designed.py) ----
INK, INK_2, MUTED = "#0B0B0B", "#52514E", "#898781"
SURFACE, PAGE, HAIRLINE = "#FCFCFB", "#F7F7F5", "#E1E0D9"
DARK_BG, DARK_INK2 = "#0D1B2A", "#AEB8C4"
BLUE, AQUA, YELLOW, GREEN = "#2A78D6", "#1BAF7A", "#EDA100", "#008300"
VIOLET, RED, MAGENTA, ORANGE = "#4A3AA7", "#E34948", "#E87BA4", "#EB6834"
GOOD, WARNING, SERIOUS, CRITICAL = "#0CA30C", "#FAB219", "#EC835A", "#D03B3B"

CAT_A = CRITICAL   # 科學歸因
CAT_B = ORANGE     # 資料定義與嚴謹
CAT_C = BLUE       # 定位與規範

CJK_BOLD_PATH = "/work/ymj1123ntu/NTU-Thesis-Writing-Template/fonts/chinese/cwTeX_Hei.ttf"
CJK_REG_PATH = "/work/ymj1123ntu/NTU-Thesis-Writing-Template/fonts/chinese/cwTeX_Yuan.ttf"
fp_bold = fm.FontProperties(fname=CJK_BOLD_PATH)
fp_reg = fm.FontProperties(fname=CJK_REG_PATH)
FONT = "Microsoft JhengHei"

plt.rcParams.update({
    "axes.spines.top": False, "axes.spines.right": False, "axes.spines.left": False,
    "figure.facecolor": SURFACE, "savefig.dpi": 200, "savefig.bbox": "tight",
    "savefig.pad_inches": 0.15, "font.family": "DejaVu Sans",
})

SW, SH = Inches(13.333), Inches(7.5)
FOOTER_TITLE = "BIB Review · 投稿前決策討論 · 2026-07-14"


def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ═══════════════════════════════════════════════════════════════════════
# pptx primitives
# ═══════════════════════════════════════════════════════════════════════

def set_no_line(s): s.line.fill.background()
def fill_solid(s, c): s.fill.solid(); s.fill.fore_color.rgb = rgb(c)


def add_bg(sl, color):
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    fill_solid(bg, color); set_no_line(bg); bg.shadow.inherit = False
    return bg


def txt(sl, text, left, top, width, height, size=14, color=INK, bold=False,
        align=PP_ALIGN.LEFT, italic=False, anchor=None, line_spacing=None):
    tb = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = rgb(color); r.font.name = FONT
    return tb


def accent_rule(sl, left, top, width=1.0, color=BLUE, thickness=0.05):
    r = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(thickness))
    fill_solid(r, color); set_no_line(r); r.shadow.inherit = False
    return r


def card(sl, left, top, width, height, fill=SURFACE, border=HAIRLINE, border_w=1.0, radius=True):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    c = sl.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    fill_solid(c, fill)
    if border is not None:
        c.line.color.rgb = rgb(border); c.line.width = Pt(border_w)
    else:
        set_no_line(c)
    c.shadow.inherit = False
    return c


def footer(sl, page_no, total, color=MUTED):
    txt(sl, FOOTER_TITLE, 0.6, 7.12, 8, 0.3, size=9, color=color)
    txt(sl, f"{page_no:02d} / {total:02d}", 12.2, 7.12, 0.6, 0.3, size=9, color=color, align=PP_ALIGN.RIGHT)


def content_header(sl, eyebrow_text, title_text, accent, tag=None):
    add_bg(sl, PAGE)
    txt(sl, eyebrow_text.upper(), 0.6, 0.5, 11.5, 0.35, size=12, color=accent, bold=True)
    txt(sl, title_text, 0.6, 0.82, 12.1, 0.85, size=25, color=INK, bold=True)
    accent_rule(sl, 0.62, 1.58, width=1.0, color=accent)
    if tag:
        pill = card(sl, 11.15, 0.55, 1.6, 0.42, fill=accent, border=None)
        txt(sl, tag, 11.15, 0.55, 1.6, 0.42, size=12, color="#FFFFFF", bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════════
# decision-slide builder
# ═══════════════════════════════════════════════════════════════════════

def reviewer_note(sl, text, top, accent):
    """Full-width note strip: what the reviewer flagged."""
    h = 1.15
    card(sl, 0.6, top, 12.13, h, fill="#FDF3F0" if accent == CAT_A else "#FBF6EF" if accent == CAT_B else "#F0F5FC",
         border=accent, border_w=1.0)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(top), Inches(0.09), Inches(h))
    fill_solid(bar, accent); set_no_line(bar); bar.shadow.inherit = False
    txt(sl, "REVIEWER 指出", 0.85, top + 0.1, 3, 0.3, size=11, color=accent, bold=True)
    txt(sl, text, 0.85, top + 0.4, 11.7, h - 0.45, size=12.5, color=INK, line_spacing=1.1)
    return top + h + 0.18


def option_cards(sl, options, top, height):
    """options: list of (header, tone_color, [body_lines]). 1-3 cards side by side."""
    n = len(options)
    total_w = 12.13
    gap = 0.3
    w = (total_w - gap * (n - 1)) / n
    left = 0.6
    for header, tone, lines in options:
        card(sl, left, top, w, height, fill=SURFACE, border=HAIRLINE, border_w=1.0)
        pill = card(sl, left, top, w, 0.5, fill=tone, border=None)
        # square top corners visually by overlaying — acceptable as rounded pill header
        txt(sl, header, left + 0.15, top + 0.03, w - 0.3, 0.44, size=13.5, color="#FFFFFF",
            bold=True, anchor=MSO_ANCHOR.MIDDLE)
        yy = top + 0.62
        for ln in lines:
            txt(sl, ln, left + 0.22, yy, w - 0.4, 1.2, size=11, color=INK_2, line_spacing=1.05)
            # rough line count (conservative: assume narrow CJK cells)
            cpl = max(6, int((w - 0.4) * 4.6))
            nlines = max(1, -(-len(ln) // cpl))
            yy += nlines * (11 / 72 * 1.3) + 0.14
        if yy > top + height + 0.02:
            print(f"    !! option card overflow: needs {yy - top:.2f}in of {height:.2f}in "
                  f"(header {header!r})")
        left += w + gap
    return top + height + 0.18


def reco_bar(sl, reco, advisor, top, height=1.5):
    card(sl, 0.6, top, 12.13, height, fill=DARK_BG, border=None)
    txt(sl, "建議", 0.85, top + 0.14, 2, 0.3, size=12, color=AQUA, bold=True)
    txt(sl, reco, 0.85, top + 0.44, 11.6, height - 0.5, size=12.5, color="#FFFFFF", line_spacing=1.12)
    if advisor:
        # advisor line pinned near bottom
        txt(sl, "需老師拍板：" + advisor, 0.85, top + height - 0.5, 11.6, 0.42,
            size=12, color=YELLOW, bold=True, line_spacing=1.05)
    return top + height


# ═══════════════════════════════════════════════════════════════════════
# figures
# ═══════════════════════════════════════════════════════════════════════

def fig_scorecard():
    rows = [
        ("BIB 主題與讀者適配度", "良好", GOOD),
        ("Problem-solving Protocol 定位", "大致成立，但仍有 original research 過重的風險", WARNING),
        ("形式規範完整度", "尚未達標", CRITICAL),
        ("核心發現的新穎性", "良好", GOOD),
        ("實驗因果歸因的嚴謹度", "有幾個重大問題需修正", CRITICAL),
        ("真實資料驗證", "已有進步，但仍屬 preliminary", WARNING),
        ("可重現性文件", "尚未完成", CRITICAL),
        ("今日直接投稿的建議", "不建議", CRITICAL),
    ]
    fig, ax = plt.subplots(figsize=(12.4, 5.0))
    ax.set_xlim(0, 12); ax.set_ylim(0, len(rows)); ax.axis("off")
    for i, (dim, verd, color) in enumerate(rows):
        y = len(rows) - 1 - i
        ax.add_patch(Rectangle((0, y + 0.08), 12, 0.84, facecolor=SURFACE if i % 2 == 0 else "#F1F0EC",
                               edgecolor="none"))
        ax.add_patch(Rectangle((0, y + 0.08), 0.09, 0.84, facecolor=color, edgecolor="none"))
        ax.text(0.3, y + 0.5, dim, va="center", ha="left", fontsize=13, color=INK, fontproperties=fp_bold)
        ax.text(5.4, y + 0.5, verd, va="center", ha="left", fontsize=12, color=color, fontproperties=fp_reg)
    plt.savefig(TMP / "scorecard.png"); plt.close()


def fig_decision_map():
    fig, ax = plt.subplots(figsize=(12.6, 5.3))
    ax.set_xlim(0, 12); ax.set_ylim(0, 6); ax.axis("off")
    cols = [
        ("A · 科學歸因（最重要）", CAT_A, [
            "D1  Pre-training 效果歸因混淆",
            "D2  Tokenization「固定容量」表述",
            "D3  Train-fit 天花板（Fig 5）驗證",
            "D4  Out-of-reference robustness claim",
        ]),
        ("B · 資料定義與嚴謹度", CAT_B, [
            "D5  1,535「species」定義",
            "D6  Scaling 資料 balancing 是否固定",
            "D7  +0.22pp 是否在 seed noise 內",
            "D8  真實 mock 結論保守化 + 公平性",
        ]),
        ("C · 定位與規範", CAT_C, [
            "D9   字數壓縮 7,500 → 5,000 字以下",
            "D10  文章定位再往 protocol 靠",
            "—    規範 / logistics checklist",
            "—    修改優先順序（三層）",
        ]),
    ]
    w = 3.7; gap = 0.35; x = 0.35
    for title, color, items in cols:
        ax.add_patch(FancyBboxPatch((x, 0.4), w, 5.0, boxstyle="round,pad=0.05",
                                    facecolor=SURFACE, edgecolor=color, linewidth=1.6))
        ax.add_patch(FancyBboxPatch((x, 4.55), w, 0.85, boxstyle="round,pad=0.05",
                                    facecolor=color, edgecolor="none"))
        ax.text(x + w/2, 4.97, title, ha="center", va="center", fontsize=12.5, color="white",
                fontproperties=fp_bold)
        yy = 4.0
        for it in items:
            ax.text(x + 0.25, yy, it, ha="left", va="center", fontsize=11.5, color=INK,
                    fontproperties=fp_reg)
            yy -= 0.72
        x += w + gap
    plt.savefig(TMP / "decision_map.png"); plt.close()


def fig_priority():
    fig, ax = plt.subplots(figsize=(12.6, 5.6))
    ax.set_xlim(0, 12); ax.set_ylim(0, 6); ax.axis("off")
    tiers = [
        ("投稿前不可缺少", CRITICAL, [
            "1 字數 5,000 字以下　2 清除所有 placeholder　3 建 repo/Zenodo DOI　4 作者資訊/bio/ORCID/CRediT",
            "5 cover letter 揭露碩論　6 修正 AI disclosure 用語　7 釐清 1,535 是否真 species",
            "8 釐清 scaling 各檔 balancing/distribution　9 +13.2pp 改非純 pre-training 或補 same-arch control",
            "10 重算並確認 Fig 5 train accuracy　11 移除無 genome-disjoint 支持的 out-of-reference claim",
            "12 補齊 Supplementary 的 exact commands 與 configs",
        ]),
        ("強烈建議", WARNING, [
            "13 修 Bracken 135/150bp mismatch　14 real mock 同報 full-community + in-set metrics",
            "15 解決 Clostridioides/Clostridium taxonomy　16 報告 seed limitation 或補 replicate",
            "17 加 total model / embedding 參數量帳　18 Fig 4 改成兩個獨立比較（非 additive 分解）",
        ]),
        ("有資源再補", GOOD, [
            "19 第二個 real/mock community　20 Centrifuge 或另一 production baseline",
            "21 genome-disjoint（或至少 locus-disjoint）實驗　22 MT 與 NT-v2 統一 RC-TTA protocol",
        ]),
    ]
    y = 5.55
    for title, color, lines in tiers:
        h = 0.5 + len(lines) * 0.42 + 0.2
        ax.add_patch(FancyBboxPatch((0.2, y - h), 11.6, h, boxstyle="round,pad=0.03",
                                    facecolor=SURFACE, edgecolor=color, linewidth=1.4))
        ax.add_patch(Rectangle((0.2, y - h), 0.1, h, facecolor=color, edgecolor="none"))
        ax.text(0.45, y - 0.32, title, ha="left", va="center", fontsize=13, color=color,
                fontproperties=fp_bold)
        yy = y - 0.72
        for ln in lines:
            ax.text(0.55, yy, ln, ha="left", va="center", fontsize=10.5, color=INK, fontproperties=fp_reg)
            yy -= 0.42
        y -= h + 0.22
    plt.savefig(TMP / "priority.png"); plt.close()


# ═══════════════════════════════════════════════════════════════════════
# build
# ═══════════════════════════════════════════════════════════════════════

def build():
    print("figures...")
    fig_scorecard(); fig_decision_map(); fig_priority()

    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    L6 = prs.slide_layouts[6]
    pg = [0]
    TOTAL = 16

    def blank():
        return prs.slides.add_slide(L6)

    def page(sl, dark=False):
        pg[0] += 1
        footer(sl, pg[0], TOTAL, color=DARK_INK2 if dark else MUTED)

    # 1 — cover
    sl = blank(); add_bg(sl, DARK_BG)
    accent_rule(sl, 1.0, 1.55, width=0.7, color=RED, thickness=0.05)
    txt(sl, "BIB REVIEWER 意見 · 投稿前決策討論", 1.0, 1.75, 11, 0.4, size=13, color=RED, bold=True)
    txt(sl, "依 Review 意見修改 Paper", 1.0, 2.2, 11.3, 1.0, size=40, color="#FFFFFF", bold=True)
    txt(sl, "待與老師討論的決策點 · 修改優先順序 · 分工", 1.0, 3.25, 11.3, 0.5, size=16, color="#C7D2E0")
    box = card(sl, 1.0, 4.1, 11.3, 1.15, fill="#132A40", border=RED, border_w=1.0)
    txt(sl, "Reviewer 總判斷：具備投 BIB 的合理基礎，但今天不建議直接投——\n"
            "需先處理「字數、claim 強度、實驗歸因」三大問題", 1.3, 4.25, 10.7, 0.9,
        size=15, color="#FFFFFF", bold=True, line_spacing=1.15)
    txt(sl, "楊明儒 · 台大生醫電資所 · 2026-07-14", 1.0, 6.6, 11.3, 0.4, size=11, color=DARK_INK2)
    page(sl, dark=True)

    # 2 — scorecard
    sl = blank(); content_header(sl, "Reviewer 總判斷", "八個面向的評價", INK)
    add_image = lambda p, l, t, w, h: sl.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))
    add_image(TMP / "scorecard.png", 0.5, 1.5, 12.3, 4.55)
    txt(sl, "最大優點：已不再是「提出一個 NT-v2 分類器」，而是 controlled benchmark + tokenizer-dependent scaling "
            "+ 因子拆解 + Kraken2/Bracken + real mock + Box1/Box2 —— 方向很符合 BIB。",
        0.6, 6.15, 12.1, 0.7, size=11.5, color=INK_2, line_spacing=1.1)
    page(sl)

    # 3 — decision map
    sl = blank(); content_header(sl, "決策地圖", "今天要跟老師過的十個決策 + 兩份清單", INK)
    sl.shapes.add_picture(str(TMP / "decision_map.png"), Inches(0.35), Inches(1.5), Inches(12.6), Inches(5.3))
    page(sl)

    # 4 — D1 pre-training attribution
    sl = blank(); content_header(sl, "決策 1 · 科學歸因（最重要）", "「+13.2pp pre-training 效果」不是純 pre-training", CAT_A, tag="A")
    y = reviewer_note(sl,
        "已證實：+13.2pp 比的是「29 層 / 498M / pretrained NT-v2」對「1 層、d_model=128 的 random-init shallow transformer」"
        "（v11_shallow config）——同時改變了 pre-training、深度、容量、架構。configs 裡沒有任何 same-architecture 的 29 層 random-init 對照。"
        "Abstract / Key Points / Fig 4 / cover letter 都把它當成 isolated pre-training effect。",
        1.42, CAT_A)
    y = option_cards(sl, [
        ("選項 A · 補實驗（最有說服力）", BLUE, [
            "加一個「相同 NT-v2 架構」的 random-init control：同 29 層、同 hidden dim、同 tokenizer、同資料。",
            "全模型從頭訓成本太高 → 退一步用 random NT-v2 backbone + LoRA/head，界定為「PEFT regime 下的 pre-training effect」。",
            "代價：需要算力。",
        ]),
        ("選項 B · 降級 claim（投稿底線）", ORANGE, [
            "改成：相同資料與 6-mer 下，pretrained NT-v2 system 勝過小型 shallow scratch baseline 13.2pp；因架構與容量也不同，不可解讀為純 pre-training。",
            "Fig 4 不再畫成 additive 分解（+13.2 / +20.4 不是同一 factorization 可相加）。",
        ]),
    ], y, 2.75)
    reco_bar(sl, "B 一定要做（投稿前的底線改法）；若算力允許，補 A 的 LoRA-regime control 最能站得住腳。",
             "是否值得花算力補 same-architecture control？", y, 1.3)
    page(sl)

    # 5 — D2 tokenization fixed capacity
    sl = blank(); content_header(sl, "決策 2 · 科學歸因", "「Tokenization, at fixed capacity」也不完全正確", CAT_A, tag="A")
    y = reviewer_note(sl,
        "6-mer → 13-mer 同時改了 vocabulary、embedding 參數量、序列長度、memory、optimization。你自己算過 overlapping 13-mer "
        "realized vocab ≈ 33.5M × 64-dim ≈ 2.1B embedding 參數 —— 所以不能宣稱「tokenization, not capacity」。",
        1.42, CAT_A)
    y = option_cards(sl, [
        ("要做的修改", BLUE, [
            "段落標題「Tokenization, at fixed capacity, helps more」→「Long-k tokenization within a fixed encoder design yields the largest gain」。",
            "主表加一欄：realized vocabulary / encoder params / embedding params / total trainable params。",
            "naive Bayes control 支持的是「高階 k-mer 表徵更 discriminative」，不是「已排除 capacity」。",
        ]),
    ], y, 2.75)
    reco_bar(sl, "這一項主要是 framing + 補參數量表，較明確，不需補實驗。我可先草擬新標題與參數量表。",
             "framing 是否接受？參數量表放主文還是 Supplementary？", y, 1.3)
    page(sl)

    # 6 — D3 train-fit ceiling
    sl = blank(); content_header(sl, "決策 3 · 科學歸因", "Train-fit 天花板（Fig 5）論證太強、且需驗證", CAT_A, tag="A")
    y = reviewer_note(sl,
        "「train 上不去 66% → 一定是 representation，不是 capacity/optimization」推論太強。更關鍵：train acc < val acc 很可疑——"
        "必須確認 train accuracy 是否在 eval mode（關 dropout/augmentation、deterministic）算的。若在 train mode 算，就不能跟 deterministic val 直接比。",
        1.42, CAT_A)
    y = option_cards(sl, [
        ("查證結果（已證實，問題成立）", CRITICAL, [
            "train acc 用 train mode（dropout + RC-aug 開著）算，非 eval-mode。",
            "且 train < val（250M scratch 63.2<64.1、warm 65.8<66.6）—— 被低估的紅旗。",
            "→「6-mer 連 train 都 fit 不上」目前站不住。",
        ]),
        ("要做的修改", BLUE, [
            "用 model.eval()、關 dropout/aug、同 metric 重算 train acc。",
            "Supp 補：duplicate rate、衝突 label 數、OOV、per-genome acc。",
            "措辭改「與 representation ceiling 一致，但不排除 capacity/opt.」。",
        ]),
    ], y, 2.75)
    reco_bar(sl, "已確認原本是 train-mode 且 train<val → Fig 5 必須用 eval-mode 重算，否則刪掉「fit 不了 training」這條論據、改用其他證據。",
             "重算 Fig 5，還是直接拿掉 train-fit 論證？", y, 1.3)
    page(sl)

    # 7 — D4 out-of-reference claim
    sl = blank(); content_header(sl, "決策 4 · 科學歸因", "Out-of-reference robustness claim 目前無證據支持", CAT_A, tag="A")
    y = reviewer_note(sl,
        "已證實：「219 species Kraken2 得 0%」其實是那些 species 不在 Kraken2 自建 DB 內，但它們都在 neural 的訓練集內 —— 這是"
        "「部署覆蓋不對稱」，不是 neural 對未見物種的泛化。（團隊內部 5/30–5/31 已改口，但論文 §6 文字仍寫 out-of-database、retain non-zero posterior。）"
        "且 coverage-matched pool read acc：Kraken2 77.7% > NT-v2 6-mer 68.9%，只有 MT 13-mer（非 GFM）94.9%。",
        1.42, CAT_A)
    y = option_cards(sl, [
        ("選項 A · 改成正確 framing（投稿底線）", ORANGE, [
            "改成「部署覆蓋不對稱」：Kraken2 要 FASTA 才能建 DB、neural 不用。",
            "刪「distinct value = read-level accuracy」（最高的是 MT 13-mer、非 GFM）。",
            "closed-set 98.7% 不算泛化；non-zero posterior ≠ generalization。",
        ]),
        ("選項 B · 做 P0 實驗", BLUE, [
            "做 §8 已列的 P0：genome-disjoint / out-of-genome generalization，才能真正宣稱 out-of-reference robustness。",
            "這也是全文最重要的後續工作。",
        ]),
    ], y, 2.75)
    reco_bar(sl, "先做 A（改成「覆蓋不對稱」的正確 framing，投稿底線）；P0 列為投稿後或平行進行的最重要實驗。",
             "P0 實驗的排程與是否納入這次投稿範圍？", y, 1.3)
    page(sl)

    # 8 — D5 species definition
    sl = blank(); content_header(sl, "決策 5 · 資料定義", "1,535「species」的定義要對齊", CAT_B, tag="B")
    y = reviewer_note(sl,
        "碩論寫明：source taxonomy 只到 genus，因此每個 genome 被當成一個 distinct「species」、以 genome ID 為 label。"
        "但 BIB 稿多次直接寫「1,535 species」——若資料定義沒變，這在 taxonomy 上不精確。",
        1.42, CAT_B)
    y = option_cards(sl, [
        ("查證結果（已證實）", CRITICAL, [
            "碩論明講：每個 genome 當一個「species」、用 genome ID 當 label。",
            "label 檔全是 genome ID、非學名；genus（120）才是唯一真階層。",
        ]),
        ("選項 A · 重新映射", VIOLET, [
            "用 GTDB/NCBI taxonomy 重新映射成真正 species labels。工作量大。",
        ]),
        ("選項 B · 改稱呼（推薦）", ORANGE, [
            "全文改成 1,535 genome-level classes / species proxies / genome identifiers within 120 genera。誠實、低成本。",
        ]),
    ], y, 2.75)
    reco_bar(sl, "推薦 B（誠實、幾乎零實驗成本）；若時間允許再考慮 A 的正式重映射。",
             "採用 B（改稱呼）還是 A（重映射）？", y, 1.3)
    page(sl)

    # 9 — D6 scaling balancing
    sl = blank(); content_header(sl, "決策 6 · 資料定義", "Scaling 各檔的 balancing / distribution 是否固定", CAT_B, tag="B")
    y = reviewer_note(sl,
        "已證實（最關鍵）：500K 是 imbalanced/natural，5M / 50M / 250M 都是 species-balanced，且各 scale 由 258M 源池「獨立抽樣」（seed 42）、"
        "並非 nested subset。所以 500K→5M 的 +7.76pp 同時混了「資料量 + natural→balanced 切換」。好消息：50M→250M 兩邊都 balanced，是乾淨比較。",
        1.42, CAT_B)
    y = option_cards(sl, [
        ("要做的修改（多半只需表述）", BLUE, [
            "加 dataset table：nominal / unique reads（源池 42.64M/258.67M；50M→17.8M、250M→41.9M unique）、sampling scheme、seed、是否 nested。",
            "明講 500K 是 natural、5M+ 是 species-balanced；把飽和論點的證據明確落在「50M→250M（皆 balanced）」這段乾淨比較上。",
            "「natural-distribution pool」→「catalogue-proportional pool」（來自目錄各屬 genome 數，不是真實 microbiome abundance）。",
        ]),
    ], y, 2.75)
    reco_bar(sl, "核心飽和論點（50M→250M）本身乾淨、不用補跑；主要是把 500K→5M 的 balancing 切換誠實講清楚 + 加 dataset table。",
             "500K→5M 這段：只重新表述即可，還是要補一個一致 distribution 的對照點？", y, 1.3)
    page(sl)

    # 10 — D7 seed variation
    sl = blank(); content_header(sl, "決策 7 · 嚴謹度", "+0.22pp 可能落在 training-seed 變異內", CAT_B, tag="B")
    y = reviewer_note(sl,
        "已證實：全部 17 個 config 都只有 seed 42、單次執行，沒有 multi-seed。而 250M scratch vs warm 光是 init 差異就有 ~2.5pp（64.8 vs 67.3）——"
        "遠大於 +0.22pp。所以 50M→250M 的 +0.22pp 沒有 noise 估計，不宜當精確實驗效果。（13-mer 的 +11pp 太大，不受影響。）",
        1.42, CAT_B)
    y = option_cards(sl, [
        ("選項 A · 補 seeds", BLUE, [
            "50M / 250M 各跑 2–3 seeds，或報 checkpoint-to-checkpoint variability，或事先設 practical equivalence margin（±0.5pp）。",
        ]),
        ("選項 B · 降級表述（推薦）", ORANGE, [
            "改成「在本次訓練條件下，50M→250M 沒有實務上有意義的提升」，並明確聲明 full-scale 是 single-run。",
            "核心「飽和」論點仍成立。",
        ]),
    ], y, 2.75)
    reco_bar(sl, "推薦 B（低成本、論點不受損）；有算力再補 A 的 multi-seed 讓數字更硬。",
             "是否要為 6-mer 50M/250M 補 seeds？", y, 1.3)
    page(sl)

    # 11 — D8 real mock conservatism
    sl = blank(); content_header(sl, "決策 8 · 嚴謹度", "真實 mock 結論要更保守 + 三個公平性問題", CAT_B, tag="B")
    y = reviewer_note(sl,
        "D6331 的加入是重要進步，但：① 已證實 Bracken 用 -r 150、只建 150-mer distribution，真實 reads median 135bp（mismatch）；"
        "② in-set renormalization 掩蓋 open-set failure，100% 「Reads clf.」只是沒 abstention；③ Clostridioides/Clostridium taxonomy 還寫「may / if older names」。",
        1.42, CAT_B)
    y = option_cards(sl, [
        ("公平性修正", BLUE, [
            "Bracken：trim 到 135/150 或建對應 distribution，至少做 sensitivity。",
            "同時報 full-community + in-set metrics；「Reads clf.」→「Assigned-read fraction」。",
            "查清 taxonomy，把 C. difficile 定 in/out-of-set，刪「may / if」。",
        ]),
        ("措辭保守化", ORANGE, [
            "標題「in-database dominance does not transfer」→「A preliminary real-mock probe reveals a large simulation-to-real gap」。",
            "結論改「simulated in-database rankings may not transfer unchanged to real communities」。",
        ]),
    ], y, 2.75)
    reco_bar(sl, "措辭保守化一定要做；Bracken 重跑與 open-set metrics 強烈建議補（會讓「classical 優勢消失」的話站得更穩）。",
             "是否重跑 Bracken（135bp）、補哪些 open-set metrics？", y, 1.3)
    page(sl)

    # 12 — D9 word count
    sl = blank(); content_header(sl, "決策 9 · 定位與規範", "字數壓縮：7,500 → ≤5,000（目標 4,500–4,800）", CAT_C, tag="C")
    y = reviewer_note(sl,
        "BIB Problem-solving Protocol 建議 2,000–5,000 words，超過上限 1,000+ words 常被 editorial office 直接退回。"
        "目前 references 前約 7,500 words，幾乎確定超標。",
        1.42, CAT_C)
    y = option_cards(sl, [
        ("移到 Supplementary", MUTED, [
            "RC-TTA 完整 Fig 3 + 機制、genus/species-balanced 完整結果、species hierarchical routing 細節、",
            "unique-key/majority/naive Bayes 完整推導、full compute-cost、per-replicate mock residuals、完整 backbone comparison。",
        ]),
        ("主文優先保留", BLUE, [
            "protocol、50M→250M scaling、正確表述的 pretrain/tokenizer 比較、",
            "Kraken2+Bracken、real mock、practical recommendations、limitations。",
            "Box1/Box2/Table5 重複 → 併成一 protocol box + 一 decision table。",
        ]),
    ], y, 2.75)
    reco_bar(sl, "我可以先出一版壓到 ~4,800 字的主文草案（把上述細節移到 Supplementary、合併 Box）。",
             "主文的紅線：哪些一定要留在正文？", y, 1.3)
    page(sl)

    # 13 — D10 positioning
    sl = blank(); content_header(sl, "決策 10 · 定位與規範", "文章定位要再往「protocol」中心靠", CAT_C, tag="C")
    y = reviewer_note(sl,
        "BIB 是 review journal、不收 pure original research。目前 11 頁大多是自己的新實驗與結果 → editor 可能視為「original research + 一個 protocol box」，有 desk-reject 風險。",
        1.42, CAT_C)
    y = option_cards(sl, [
        ("結構性重組（降低 desk-reject 風險）", BLUE, [
            "把 Box 1 提前到 benchmark design 之後；Methods 以「benchmark should be constructed as follows」組織。",
            "每節先提出 protocol principle，再用你的結果支持；減少實驗歷程式敘述；增加與既有 benchmark practices 的連結。",
            "Conclusion 以「讀者該如何評估」為中心，而不是「我們的 13-mer 最好」。",
        ]),
    ], y, 2.75)
    reco_bar(sl, "方向正確、值得做，但屬較大的結構工程；建議與字數壓縮（D9）一起做，一次到位。",
             "是否接受這種結構性重組？要做到多徹底？", y, 1.3)
    page(sl)

    # 14 — logistics checklist
    sl = blank(); content_header(sl, "規範 / Logistics", "投稿前必補（多數是行政、非科學）", CAT_C, tag="C")
    items = [
        ("字數 ≤5,000（見 D9）", INK_2),
        ("清除所有 placeholder：Second Author、funding、repository、DOI、X.X./Y.Y. CRediT、colleagues/reviewers、supp「confirm scheduler」", INK_2),
        ("Data availability 不能等 acceptance：先建 Zenodo release + 固定 git commit + split manifests + read hashes", INK_2),
        ("＋ taxonomy mapping、per-read predictions、ART params/seed/accession、Kraken/Bracken build scripts、persistent DOI", MUTED),
        ("Title page：作者/單位、通訊電話+email、每位作者 ~30 words bio、submitting author 連 ORCID", INK_2),
        ("Cover letter 揭露碩士論文（相同/相近資料）+ 與投稿稿差異 + 附 thesis 副本", INK_2),
        ("AI disclosure：把「help draft」改成準確的「language editing / formatting / code assistance / consistency checking」（依實際使用）", INK_2),
        ("　※ 若曾用 prompt 產生大段正文 → 與 BIB 政策衝突，需作者重寫並如實揭露，必要時先問 editorial office", CRITICAL),
    ]
    yy = 1.55
    cbox = card(sl, 0.6, 1.5, 12.13, 4.55, fill=SURFACE, border=HAIRLINE)
    yy = 1.72
    for text, color in items:
        m = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(yy + 0.06), Inches(0.09), Inches(0.09))
        fill_solid(m, CAT_C if color == INK_2 else color); set_no_line(m); m.shadow.inherit = False
        txt(sl, text, 1.1, yy, 11.4, 0.6, size=12.5, color=color, line_spacing=1.05)
        cpl = 40
        yy += max(1, -(-len(text) // cpl)) * (12.5 / 72 * 1.3) + 0.16
    txt(sl, "需老師 / 行政提供：funding 資訊、致謝名單、第二作者、ORCID、通訊電話。",
        0.6, 6.2, 12.1, 0.5, size=12.5, color=INK, bold=True)
    page(sl)

    # 15 — priority tiers
    sl = blank(); content_header(sl, "修改優先順序", "Reviewer 的三層清單", INK)
    sl.shapes.add_picture(str(TMP / "priority.png"), Inches(0.35), Inches(1.5), Inches(12.6), Inches(5.5))
    page(sl)

    # 16 — next steps / division of labor
    sl = blank(); content_header(sl, "建議下一步", "分工：今天拍板 / 我能立刻做 / 需實驗端", GREEN)
    option_cards(sl, [
        ("今天會議要拍板", CAT_A, [
            "pre-training control 要不要補（算力）",
            "1,535 label 用哪種稱呼",
            "scaling 若 confound 怎麼處理",
            "P0 排程、Bracken 是否重跑",
            "字數紅線 / 是否結構重組",
        ]),
        ("我可立刻做（pull 後）", BLUE, [
            "純措辭降級：pre-train / tokenization / information-ceiling / out-of-ref claim",
            "加 Kraken2+Bracken row、Fig S1「multi-billion-parameter」",
            "natural → catalogue-proportional、合併 Box",
            "字數壓縮草案、dataset table 骨架",
        ]),
        ("需你 / 實驗端", ORANGE, [
            "Fig 5 train acc 重算",
            "balancing 查證後定稿",
            "seed、Bracken 重跑",
            "taxonomy 查清、Zenodo",
            "作者資訊 / ORCID",
        ]),
    ], 1.55, 4.0)
    card(sl, 0.6, 5.75, 12.13, 0.95, fill=DARK_BG, border=None)
    txt(sl, "底線：先把「必補」清單清完再投——今天最該解決的不是再加 headline，而是讓每個 claim 都與實驗能支持的範圍一致。",
        0.85, 5.9, 11.6, 0.7, size=13.5, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    page(sl)

    prs.save(OUT)
    print(f"saved: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()

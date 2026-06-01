#!/usr/bin/env python3
"""Generate weekly progress slides for 2026-05-19."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

TEMPLATE = Path("/work/ymj1123ntu/token_level_gfm_classifier/docs/0512.pptx")
OUT      = Path("/work/ymj1123ntu/token_level_gfm_classifier/docs/0519.pptx")
FIG      = Path("/work/ymj1123ntu/token_level_gfm_classifier/docs/figures_0519")
THESIS_FIG = Path("/work/ymj1123ntu/thesis/figures")

BLUE   = RGBColor(0x00, 0x70, 0xC0)
BLACK  = RGBColor(0x00, 0x00, 0x00)
GRAY   = RGBColor(0x40, 0x40, 0x40)
GREEN  = RGBColor(0x00, 0x70, 0x00)
RED    = RGBColor(0xC0, 0x00, 0x00)
ORANGE = RGBColor(0xE0, 0x70, 0x00)

prs = Presentation(TEMPLATE)
xml_slides = prs.slides._sldIdLst
for sld in list(xml_slides):
    xml_slides.remove(sld)

LAY_TITLE   = prs.slide_layouts[0]
LAY_CONTENT = prs.slide_layouts[1]
LAY_SECTION = prs.slide_layouts[19]

slide_w = prs.slide_width
slide_h = prs.slide_height


def add_section(text):
    slide = prs.slides.add_slide(LAY_SECTION)
    txBox = slide.shapes.add_textbox(
        Emu(590550), Emu(1530900), Emu(7962900), Emu(2081700))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(48)
    run.font.color.rgb = BLUE
    run.font.name = "Calibri"


def add_content(title, bullets):
    slide = prs.slides.add_slide(LAY_CONTENT)
    title_ph = slide.placeholders[0]
    title_ph.text = title
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(24)
    body_ph = slide.placeholders[1]
    tf = body_ph.text_frame
    tf.clear()
    first = True
    for (level, text, bold, color) in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.size = Pt(16) if level == 0 else Pt(14)
    return slide


def add_figure_slide(title, img_path, caption=""):
    slide = prs.slides.add_slide(LAY_CONTENT)
    title_ph = slide.placeholders[0]
    title_ph.text = title
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(24)
    body_ph = slide.placeholders[1]
    body_ph.text = ""
    img_top  = Emu(int(slide_h * 0.16))
    img_left = Emu(int(slide_w * 0.04))
    img_w    = Emu(int(slide_w * 0.92))
    img_h    = Emu(int(slide_h * 0.70))
    slide.shapes.add_picture(str(img_path), img_left, img_top, img_w, img_h)
    if caption:
        txBox = slide.shapes.add_textbox(
            Emu(int(slide_w * 0.05)),
            Emu(int(slide_h * 0.88)),
            Emu(int(slide_w * 0.90)),
            Emu(int(slide_h * 0.10)),
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = caption
        run.font.size = Pt(12)
        run.font.color.rgb = GRAY
    return slide


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(LAY_TITLE)
    slide.placeholders[0].text = title
    for para in slide.placeholders[0].text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.bold = True
    sub_ph = slide.placeholders[1]
    sub_ph.text = subtitle
    for para in sub_ph.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(18)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: Title
# ─────────────────────────────────────────────────────────────────────────────
add_title_slide(
    "Weekly Progress Update",
    "0519   |   楊明儒 (資工所碩二)"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: This Week's Highlights
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "This Week's Highlights",
    [
        (0, "1  Genus-level: all 3 tokenization variants evaluated at sample level", True, BLUE),
        (1, "MT 13-mer genus: Pearson r = 0.9993, BC = 0.028  ← near-perfect", True, GREEN),
        (1, "NT-Genus: Pearson r = 0.9932, BC = 0.094  |  MT 6-mer: r = 0.984, BC = 0.167", False, BLACK),
        (1, "Pearson r saturates ≥ 0.98 for ALL; Bray-Curtis is the discriminating metric at genus level", True, ORANGE),
        (0, "2  Species-level: MT 13-mer crosses the practical utility threshold", True, BLUE),
        (1, "MT 13-mer flat: Pearson r = 0.466, ROC AUC = 0.966  |  hier: 0.478, 0.967  ← useful", True, GREEN),
        (1, "NT-Species flat r = 0.135, hier r = 0.106  |  MT 6-mer: 0.034–0.065  ← all below threshold", True, RED),
        (0, "3  Hierarchical masking direction is router-dependent (monotonic in router accuracy)", True, GREEN),
        (1, "MT 13-mer router 87.5% → +0.012 r, +0.0015 ROC, +1.12 pp acc (consistent gain)", True, GREEN),
        (1, "NT-Genus router 66% → −0.029 r, +0.002 ROC, −2.0 pp acc (4 of 5 metrics hurt)", True, RED),
        (1, "MT 6-mer router 48.9% → −0.031 r, −0.023 ROC, −2.8 pp acc (largest degradation)", True, RED),
        (0, "4  Tokenization effect holds at sample level — not just read accuracy", True, GREEN),
        (1, "13-mer BC 0.028 (genus) / 0.369 (species) vs 6-mer 0.167 / 0.686 — systematic gap", True, GREEN),
        (0, "5  Next: MT per-genus classifiers (Taiwana2) — expected to improve species detection", True, ORANGE),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3: Section
# ─────────────────────────────────────────────────────────────────────────────
add_section("Genus-Level: Tokenization Stratifies Abundance Quality")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4: Genus 3-model results table
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "Genus Sample-Level Evaluation — 3 Models  (5M pool, 120 genera, 50K reads/sample)",
    [
        (0, "Model               Read Acc   Pearson r    Bray-Curtis   Spearman r   ROC AUC", True, BLUE),
        (0, "─────────────────────────────────────────────────────────────────────────────", False, GRAY),
        (0, "MT 13-mer genus       87.4%     0.9993       0.0279         0.9778       0.900  ✓", True, GREEN),
        (0, "NT-Genus        67.1%     0.9932       0.0937         0.9312       0.705", False, BLACK),
        (0, "MT  6-mer genus       48.8%     0.9840       0.1666         0.8537       0.569", False, BLACK),
        (0, "─────────────────────────────────────────────────────────────────────────────", False, GRAY),
        (0, "Observations:", True, BLUE),
        (1, "Pearson r ≥ 0.984 for ALL three — genus-level is forgiving (only 120 classes)", True, GREEN),
        (1, "Bray-Curtis clearly separates: 0.028 → 0.094 → 0.167 — BC is the better discriminator", True, ORANGE),
        (1, "ROC AUC reflects class-level detection; low for 6-mer (many genera with few predictions)", False, GRAY),
        (0, "Why genus is forgiving:", True, BLUE),
        (1, "120 genera: expected abundance per genus ≈ 1/120 = 0.83%  → much more reads per class", False, BLACK),
        (1, "Even 50% classifier correctly places most reads near the right genus", False, BLACK),
        (0, "Key finding: tokenization still stratifies quality — better tokenizer → lower BC dissimilarity", True, GREEN),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5: Genus 3-model bar figure
# ─────────────────────────────────────────────────────────────────────────────
add_figure_slide(
    "Genus-Level: Effect of Tokenization on Abundance Quality  (5M pool, 120 genera)",
    FIG / "genus_comparison_3model.png",
    "All 3 models achieve Pearson r ≥ 0.984; Bray-Curtis clearly separates (0.028 / 0.094 / 0.167)"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5b: Genus scatter — accuracy vs abundance metrics
# ─────────────────────────────────────────────────────────────────────────────
add_figure_slide(
    "Genus: Read Accuracy vs Abundance Quality  (BC degrades, Pearson r saturates)",
    FIG / "genus_acc_vs_abundance.png",
    "Pearson r near-saturated across all models; Bray-Curtis is the sensitive metric at genus level"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6: Section
# ─────────────────────────────────────────────────────────────────────────────
add_section("Species-Level: MT 13-mer Enables Practical Detection")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7: Species results summary table
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "Species Sample-Level Evaluation  (1535 species, reads_per_sample=1000)",
    [
        (0, "Model                Read Acc   Pearson r   Bray-Curtis   ROC AUC", True, BLUE),
        (0, "─────────────────────────────────────────────────────────────────", False, GRAY),
        (0, "NT-Species  flat     17.8%      0.135       0.589         0.794", False, BLACK),
        (0, "NT-Species  hier.    15.8%      0.106       0.608         0.796", True, RED),
        (0, "─────────────────────────────────────────────────────────────────", False, GRAY),
        (0, "MT  6-mer  flat        9.2%      0.065       0.654         0.690", False, BLACK),
        (0, "MT  6-mer  hierarchical  6.4%    0.034       0.686         0.667", True, RED),
        (0, "─────────────────────────────────────────────────────────────────", False, GRAY),
        (0, "MT 13-mer  flat       49.7%      0.466       0.378         0.966  ✓", True, GREEN),
        (0, "MT 13-mer  hierarchical 50.9%    0.478       0.369         0.967  ✓ best", True, GREEN),
        (0, "─────────────────────────────────────────────────────────────────", False, GRAY),
        (0, "All models: 100K independent test set, 1K reads/sample", False, GRAY),
        (0, "Key observations:", True, BLUE),
        (1, "1. NT-Species 17.8% → Pearson r = 0.135 — noise floor constrains r; 3.4× below MT 13-mer", True, RED),
        (1, "2. MT 13-mer achieves Pearson r 0.47 and ROC AUC 0.967 — practically useful", True, GREEN),
        (1, "3. MT 6-mer: higher read acc (9.2%) than hier. (6.4%) but both unusable (r < 0.07)", True, RED),
        (1, "4. Hier masking effect monotonic in router quality: 87.5%→help, 66%→hurt, 48.9%→hurt more", True, BLUE),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8: Figure — species summary 4-panel
# ─────────────────────────────────────────────────────────────────────────────
add_figure_slide(
    "Species-Level Metrics: All 4 Dimensions  (100K test, 1535 species)",
    FIG / "species_summary_4panel.png",
    "Read accuracy, Pearson r, 1 − Bray-Curtis, ROC AUC — all favor MT 13-mer"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9: Figure — abundance estimation bar
# ─────────────────────────────────────────────────────────────────────────────
add_figure_slide(
    "Abundance Correlation and Dissimilarity Across All Models",
    FIG / "abundance_comparison_bar.png",
    "All models evaluated on 100K independent test set  (1535 species, reads_per_sample=1K)"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10: Figure — ROC AUC bar
# ─────────────────────────────────────────────────────────────────────────────
add_figure_slide(
    "Species Detection Performance: ROC AUC Comparison",
    FIG / "roc_auc_comparison.png",
    "MT 13-mer 0.966–0.967  vs  NT-Species flat 0.794 / hier 0.796 (only metric not hurt)  vs  MT 6-mer 0.667–0.690"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11: Figure — operating points
# ─────────────────────────────────────────────────────────────────────────────
add_figure_slide(
    "Species Detection: Sensitivity at Fixed Specificity  (ROC Operating Points)",
    FIG / "roc_operating_points.png",
    "At 90% specificity: MT 13-mer hier. sens=92.5%  vs  NT-Species sens=55.8%  vs  MT 6-mer flat sens=41.2%"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12: Figure — detection threshold curves
# ─────────────────────────────────────────────────────────────────────────────
add_figure_slide(
    "Species Detection: Sensitivity & Specificity vs Read Count Threshold",
    FIG / "detection_threshold_curves.png",
    "MT 13-mer flat: sens=78.4% / spec=99.6% at ≥5 reads  |  MT 6-mer: 20.3% / 98.8% — poor recall"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13: MT 13-mer flat vs hier detail
# ─────────────────────────────────────────────────────────────────────────────
add_figure_slide(
    "MT 13-mer: Flat vs Hierarchical — Abundance Threshold Detection",
    FIG / "mt13_flat_vs_hier_detection.png",
    "Hierarchical consistently better at all thresholds — gains most at ≥1% abundance (sens 56.7% → 58.4%)"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14: Section
# ─────────────────────────────────────────────────────────────────────────────
add_section("Why Read Accuracy → Abundance Quality")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15: The threshold finding
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "Connecting Read Accuracy to Sample-Level Utility",
    [
        (0, "Rule of thumb from the data:", True, BLUE),
        (1, "< ~20% read acc  →  Pearson r low (noise floor limited)  (NT-Species: 17.8% → 0.135, 6-mer: 9.2% → 0.065)", True, RED),
        (1, "~50% read acc   →  Pearson r ≈ 0.47  (MT 13-mer flat: 49.7% → 0.466)", True, GREEN),
        (1, "~67% read acc   →  Pearson r ≈ 0.99  (NT-v2 genus: 67.1% → 0.993) — but only 120 classes", True, GREEN),
        (0, "The non-linear jump (≈ 0.466 at 50%) makes sense:", True, BLUE),
        (1, "1535 species: true abundance per species = 1/1535 ≈ 0.065%  → each misclassification heavily dilutes", False, BLACK),
        (1, "A 50% classifier correctly assigns half the reads → abundance vector is partially recoverable", False, BLACK),
        (1, "A 9% classifier assigns 91% of reads to wrong species → abundance vector is dominated by noise", False, BLACK),
        (0, "Practical threshold for metagenomics: ~40–50% top-1 read accuracy needed", True, ORANGE),
        (1, "NT-v2 architecture cannot reach this with 6-mer tokenization (ceiling: 15.9% NT-Species)", True, RED),
        (1, "MT 13-mer already crosses it (49.7%); overlapping tokenization is the key design choice", True, GREEN),
        (0, "Conclusion: tokenization design is not just a classification metric — it determines practical utility", True, GREEN),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16: Section
# ─────────────────────────────────────────────────────────────────────────────
add_section("Next Steps")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17: Pending experiments
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "Remaining Experiments & Timeline",
    [
        (0, "On TWCC (done / running):", True, GREEN),
        (1, "✓  NT-Genus sample-level evaluation  (Pearson r = 0.993)", False, GREEN),
        (1, "✓  NT-Species flat sample-level evaluation", False, GREEN),
        (1, "✓  MT 13-mer / 6-mer species (flat + hierarchical) sample-level evaluation", False, GREEN),
        (0, "On Taiwana2 (in progress):", True, ORANGE),
        (1, "Step 2: Inspect MT 6-mer training config for per-genus training command", False, BLACK),
        (1, "Step 3: Split 45M reads by genus (split_by_genus_mt.py)", False, BLACK),
        (1, "Step 4: Train 81 MT 6-mer per-genus species classifiers", False, BLACK),
        (1, "Step 5: Extract per-genus predictions on reads_remapped.fa → scp to TWCC", False, BLACK),
        (0, "After Taiwana2 per-genus results arrive:", True, BLUE),
        (1, "TWCC: run per-genus ensemble + sample-level eval (Exp B/C/F)", False, BLACK),
        (1, "Expected: oracle per-genus routing may unlock >50% species detection", False, GRAY),
        (0, "After all experiments complete: §4.13 sample-level section in thesis", True, BLUE),
        (1, "Write-up: abundance estimation, ROC AUC, detection thresholds, practical utility section", False, BLACK),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 18: Updated complete results
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "Complete Results — Updated Numeric Reference  (as of 2026-05-17)",
    [
        (0, "Genus (120 classes, RC TTA)", True, BLUE),
        (1, "NT-Genus (50M) → 67.07%  |  NT-Genus-5M (5M) → 63.05%  |  Shallow-Genus → 53.88%", False, BLACK),
        (0, "Tokenization ablation (MT, 50M, genus)", True, BLUE),
        (1, "13-mer s=1 → 87.42%  |  12-mer s=1 → 78.83%  |  6-mer s=1 → 48.87%  |  6-mer s=6 → 38.47%", False, BLACK),
        (0, "Species (1535 classes, 100K test set)", True, BLUE),
        (1, "MT 13-mer flat → 49.74%  |  MT 13-mer hierarchical → 50.86%", True, GREEN),
        (1, "NT-Species flat → 15.51%  |  MT 6-mer flat → 9.22%  |  MT 6-mer hier. → 6.41%", False, BLACK),
        (0, "Sample-level — genus (5M pool, 120 genera)", True, BLUE),
        (1, "MT 13-mer: r = 0.9993, BC = 0.028  |  NT-Genus: r = 0.9932, BC = 0.094  |  MT 6-mer: r = 0.984, BC = 0.167", True, GREEN),
        (0, "Sample-level — species (100K pool, 1535 classes)", True, BLUE),
        (1, "MT 13-mer flat: Pearson r = 0.466, BC = 0.378, ROC AUC = 0.966", True, GREEN),
        (1, "MT 13-mer hier: Pearson r = 0.478, BC = 0.369, ROC AUC = 0.967  ← best species", True, GREEN),
        (1, "NT-Species flat: Pearson r = 0.135, BC = 0.589, ROC AUC = 0.794  (100K test set)", True, RED),
        (1, "MT 6-mer flat/hier: Pearson r = 0.034–0.065  ← unusable", True, RED),
        (0, "Hierarchical (subgenus routing): 13.2–13.6%  ← FAILED (sil. = 0.03)", True, RED),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 19: Section — backup
# ─────────────────────────────────────────────────────────────────────────────
add_section("Backup / Reference")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 20: Sample eval methodology
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "Sample-Level Evaluation Methodology",
    [
        (0, "Two sampling strategies (evaluate_sample.py):", True, BLUE),
        (1, "1. Random partition samples: non-overlapping subsets of test pool", False, BLACK),
        (1, "2. Sparse community samples: replace=False from present-genera pool  (simulates low-diversity env)", False, BLACK),
        (0, "Abundance vector: predicted class count histogram / reads_per_sample", False, BLACK),
        (0, "Metrics per sample, then mean ± std across 100 samples:", False, BLACK),
        (1, "Pearson r: linear correlation between true and predicted abundance vectors", False, BLACK),
        (1, "Bray-Curtis dissimilarity: ecological β-diversity metric (0=identical, 1=no overlap)", False, BLACK),
        (1, "ROC AUC: binary detection at species level, threshold on predicted abundance", False, BLACK),
        (0, "Parameters (species, 100K pool):", True, ORANGE),
        (1, "reads_per_sample=1000, n_partition_samples=100, n_sparse_samples=200, genera_present=50", False, BLACK),
        (1, "Constrained by pool size: 100K reads / 1535 species ≈ 65 reads/species", False, GRAY),
        (0, "Parameters (genus, 5M pool):", True, BLUE),
        (1, "reads_per_sample=50000, n_partition_samples=100, n_sparse_samples=200, genera_present=60", False, BLACK),
        (0, "Caveat: genus/species metrics are NOT directly comparable (different pool sizes & reads_per_sample)", True, RED),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 21: Data scaling figure
# ─────────────────────────────────────────────────────────────────────────────
add_figure_slide(
    "Data Scaling — Genus RC TTA Accuracy vs Training Volume",
    THESIS_FIG / "data_scaling.png",
    "500K→5M: +7.76 pp  |  5M→50M: +4.02 pp"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 22–24: ROC and scatter figures from individual eval runs
# ─────────────────────────────────────────────────────────────────────────────
RES = Path("/work/ymj1123ntu/token_level_gfm_classifier/results")

add_figure_slide(
    "NT-Genus — ROC Detection Curve  (species presence/absence at genus level)",
    RES / "nt_token_genus_lora_v9_50M/eval_sample_level/roc_detection.png",
    "AUC = 0.705 — low because all 120 genera present in every sample (detection trivially high)"
)

add_figure_slide(
    "MT 13-mer Flat — ROC Detection Curve  (100K test, 1535 species)",
    RES / "mt_species_flat/eval_sample_level_100K/roc_detection.png",
    "AUC = 0.966 — high discriminative power for species presence/absence"
)

add_figure_slide(
    "MT 13-mer Hierarchical — Abundance Scatter  (100K test, 1535 species)",
    RES / "mt_hierarchical/eval_sample_level_100K/abundance_scatter.png",
    "Pearson r = 0.478 — predicted abundance correlates with true abundance at species level"
)

prs.save(OUT)
print(f"Saved: {OUT}")

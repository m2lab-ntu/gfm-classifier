#!/usr/bin/env python3
"""
Discussion / next-steps section to APPEND to the paper walkthrough deck
(paper_designed_zh). Framing: the STUDENT presenting their own forward plan to
the advisor for the first time — what to strengthen before submission, which
experiments need extra compute (possible Nano4/5 top-up), and decisions where
the advisor's steer is wanted. NOT a "respond-to-reviewer" deck.

Same design system as make_paper_designed.py so it imports seamlessly.
Output: paper_discussion_addon_zh.pptx  (import via Google Slides > Import slides)
"""

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from matplotlib.patches import Rectangle
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DOCS = Path(__file__).parent
TMP = DOCS / "discussion_addon_figs"
TMP.mkdir(exist_ok=True)
OUT = DOCS / "paper_discussion_addon_zh.pptx"

# ---- design tokens (identical to make_paper_designed.py) ----
INK, INK_2, MUTED = "#0B0B0B", "#52514E", "#898781"
SURFACE, PAGE, HAIRLINE = "#FCFCFB", "#F7F7F5", "#E1E0D9"
DARK_BG, DARK_INK2 = "#0D1B2A", "#AEB8C4"
BLUE, AQUA, YELLOW, GREEN = "#2A78D6", "#1BAF7A", "#EDA100", "#008300"
VIOLET, RED, ORANGE = "#4A3AA7", "#E34948", "#EB6834"
GOOD, WARNING, CRITICAL = "#0CA30C", "#FAB219", "#D03B3B"

CJK_BOLD = "/work/ymj1123ntu/NTU-Thesis-Writing-Template/fonts/chinese/cwTeX_Hei.ttf"
CJK_REG = "/work/ymj1123ntu/NTU-Thesis-Writing-Template/fonts/chinese/cwTeX_Yuan.ttf"
fp_bold = fm.FontProperties(fname=CJK_BOLD)
fp_reg = fm.FontProperties(fname=CJK_REG)
FONT = "Microsoft JhengHei"

plt.rcParams.update({
    "axes.spines.top": False, "axes.spines.right": False, "axes.spines.left": False,
    "figure.facecolor": SURFACE, "savefig.dpi": 200, "savefig.bbox": "tight",
    "savefig.pad_inches": 0.12, "font.family": "DejaVu Sans",
})

SW, SH = Inches(13.333), Inches(7.5)
FOOTER_TITLE = "GFM Metagenomic Benchmark · BIB Problem-solving Protocol"


def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_no_line(s): s.line.fill.background()
def fill_solid(s, c): s.fill.solid(); s.fill.fore_color.rgb = rgb(c)


def add_bg(sl, color):
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    fill_solid(bg, color); set_no_line(bg); bg.shadow.inherit = False
    return bg


def txt(sl, text, left, top, width, height, size=14, color=INK, bold=False,
        align=PP_ALIGN.LEFT, italic=False, anchor=None, line_spacing=None):
    tb = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = rgb(color); r.font.name = FONT
    return tb


def accent_rule(sl, left, top, width=1.0, color=BLUE, thickness=0.05):
    r = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(thickness))
    fill_solid(r, color); set_no_line(r); r.shadow.inherit = False


def card(sl, left, top, width, height, fill=SURFACE, border=HAIRLINE, border_w=1.0):
    c = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    fill_solid(c, fill)
    if border is not None:
        c.line.color.rgb = rgb(border); c.line.width = Pt(border_w)
    else:
        set_no_line(c)
    c.shadow.inherit = False
    return c


def footer(sl, color=MUTED):
    txt(sl, FOOTER_TITLE, 0.6, 7.12, 9, 0.3, size=9, color=color)


def content_header(sl, eyebrow_text, title_text, accent):
    add_bg(sl, PAGE)
    txt(sl, eyebrow_text.upper(), 0.6, 0.55, 11.8, 0.35, size=12, color=accent, bold=True)
    txt(sl, title_text, 0.6, 0.86, 12.1, 0.85, size=26, color=INK, bold=True)
    accent_rule(sl, 0.62, 1.62, width=1.0, color=accent)


def bullet_card(sl, items, left, top, width, height, base_size=15, gap=0.14):
    """items: (text, color, level[, bold]). marker square for level-0."""
    card(sl, left, top, width, height, fill=SURFACE, border=HAIRLINE)
    yy = top + 0.28
    for item in items:
        if len(item) == 4:
            text, color, level, bold = item
        else:
            text, color, level = item; bold = (level == 0)
        indent = 0.32 + level * 0.35
        if level == 0:
            m = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.28), Inches(yy + 0.09),
                                    Inches(0.1), Inches(0.1))
            fill_solid(m, color); set_no_line(m); m.shadow.inherit = False
        size = base_size if level == 0 else base_size - 2
        txt(sl, text, left + indent, yy, width - indent - 0.3, 1.4,
            size=size, color=(INK if level == 0 else INK_2), bold=(bold and level == 0), line_spacing=1.12)
        cpl = max(8, int((width - indent - 0.3) * 4.5 * (15 / size)))
        nlines = max(1, -(-len(text) // cpl))
        yy += nlines * (size / 72 * 1.3) + gap
    if yy > top + height + 0.03:
        print(f"    !! overflow: {yy-top:.2f}/{height:.2f}  ({items[0][0][:20]!r})")


def key_row(sl, title, body, color, left, top, width, height):
    card(sl, left, top, width, height)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.09), Inches(height))
    fill_solid(bar, color); set_no_line(bar); bar.shadow.inherit = False
    txt(sl, title, left + 0.35, top + 0.12, width - 0.6, 0.5, size=15, color=INK, bold=True)
    txt(sl, body, left + 0.35, top + 0.56, width - 0.6, height - 0.62, size=12.5, color=INK_2, line_spacing=1.1)


def divider(prs, num_label, title, sub, color):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, DARK_BG)
    accent_rule(sl, 0.9, 2.55, width=0.55, color=color, thickness=0.06)
    txt(sl, num_label, 0.9, 2.15, 8, 0.4, size=14, color=color, bold=True)
    txt(sl, title, 0.85, 2.75, 11.6, 1.2, size=40, color="#FFFFFF", bold=True)
    txt(sl, sub, 0.9, 3.9, 11.0, 0.9, size=15, color=DARK_INK2, line_spacing=1.2)
    footer(sl, color=DARK_INK2)
    return sl


# ═══════════════════════════════════════════════════════════════════════
# compute table figure
# ═══════════════════════════════════════════════════════════════════════

def fig_compute_table():
    rows = [
        ("Out-of-genome 泛化（genome-disjoint 重訓）", "回應 closed-set 98.7% 最關鍵的一步", "大", "T2 / Nano（需儲值）", "P0", CRITICAL),
        ("相同架構 pre-training 對照", "乾淨分離 pre-training 的貢獻", "中", "T2 / Nano（需儲值）", "高", ORANGE),
        ("50M（必）+ 250M（可選）multi-seed", "支撐 50M→250M 飽和不是 seed 抖動", "50M 小 · 250M 大", "同上，50M 先跑", "中", BLUE),
        ("Bracken 135bp 重估 + open-set 指標", "讓真實 mock 三方比較更公平", "小（推論為主）", "本地 4090", "中", BLUE),
        ("第二個 mock / Centrifuge baseline（可選）", "真實驗證更廣、baseline 更完整", "中", "本地 4090 + 建 DB", "低", MUTED),
    ]
    n = len(rows)
    fig, ax = plt.subplots(figsize=(12.6, 4.7))
    ax.set_xlim(0, 100); ax.set_ylim(0, n + 1.0); ax.axis("off")
    # columns: left block (name+purpose) | 量級 | 建議跑法 | 優先度
    X_NAME, X_SIZE, X_WHERE, X_PRIO = 1.5, 54, 68, 92
    ytop = n + 0.55
    ax.text(X_NAME, ytop, "實驗 / 目的", fontsize=12, fontweight="bold", color=INK, va="center", fontproperties=fp_bold)
    ax.text(X_SIZE, ytop, "量級", fontsize=12, fontweight="bold", color=INK, va="center", fontproperties=fp_bold)
    ax.text(X_WHERE, ytop, "建議跑法", fontsize=12, fontweight="bold", color=INK, va="center", fontproperties=fp_bold)
    ax.text(X_PRIO, ytop, "優先度", fontsize=12, fontweight="bold", color=INK, va="center", ha="center", fontproperties=fp_bold)
    ax.plot([0, 100], [ytop - 0.42, ytop - 0.42], color=HAIRLINE, lw=1.2)
    for i, (name, why, size, where, prio, color) in enumerate(rows):
        yc = n - 1 - i + 0.5   # row center
        if i % 2 == 0:
            ax.add_patch(Rectangle((0, yc - 0.5), 100, 1.0, facecolor="#F1F0EC", edgecolor="none"))
        ax.text(X_NAME, yc + 0.16, name, fontsize=11, color=INK, va="center", fontproperties=fp_bold)
        ax.text(X_NAME, yc - 0.22, why, fontsize=9.5, color=MUTED, va="center", fontproperties=fp_reg)
        ax.text(X_SIZE, yc, size, fontsize=10, color=INK_2, va="center", fontproperties=fp_reg)
        ax.text(X_WHERE, yc, where, fontsize=9.5, color=INK_2, va="center", fontproperties=fp_reg)
        ax.add_patch(Rectangle((X_PRIO - 3.3, yc - 0.22), 6.6, 0.44, facecolor=color, edgecolor="none"))
        ax.text(X_PRIO, yc, prio, fontsize=10.5, color="white", va="center", ha="center",
                fontweight="bold", fontproperties=fp_bold)
    plt.savefig(TMP / "compute_table.png"); plt.close()


# ═══════════════════════════════════════════════════════════════════════
# build
# ═══════════════════════════════════════════════════════════════════════

def build():
    print("figures...")
    fig_compute_table()
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    L6 = prs.slide_layouts[6]

    # 1 — divider
    divider(prs, "後續規劃", "後續規劃與討論",
            "投稿前想再補強的方向 · 需要的運算資源 · 想請老師給方向的決策", GREEN)

    # 2 — self-review: what to strengthen
    sl = prs.slides.add_slide(L6)
    content_header(sl, "自我檢視", "投稿前，我想再補強的幾個地方", BLUE)
    txt(sl, "重新檢視實驗細節後，我認為在投出去前，下面幾點可以讓論證更站得住腳：",
        0.6, 1.35, 12.1, 0.4, size=13.5, color=INK_2)
    bullet_card(sl, [
        ("因果歸因再收緊：目前「pre-training +13.2pp」是拿 29 層 pretrained 對 1 層 random-init，容量也一起變了；"
         "想補一個相同架構的對照，才能乾淨分離 pre-training 的貢獻。", BLUE, 0),
        ("Tokenization 說法精準化：13-mer 也帶進較大的 embedding 容量；想改成「固定 encoder 下、長 k-mer 帶來最大增益」，並附參數量帳。", BLUE, 0),
        ("從 closed-set 走向真實泛化：98.7% 是同一目錄的 closed-set；想補 out-of-genome 驗證，才能談對未見基因體的表現。", BLUE, 0),
        ("Scaling 敘述講清楚：500K 是自然分布、5M 以上是 species-balanced；飽和結論主要靠乾淨的 50M→250M 這段，想在文中明確交代。", BLUE, 0),
        ("用詞與範圍：1,535「species」其實是 genome-level classes（來源 taxonomy 只到 genus）；豐度優勢的說法也限縮在正確範圍。", BLUE, 0),
    ], 0.6, 1.85, 12.13, 4.9, base_size=14)
    footer(sl)

    # 3 — compute / resources
    sl = prs.slides.add_slide(L6)
    content_header(sl, "需要討論 · 運算資源", "想再跑的實驗與所需資源", ORANGE)
    txt(sl, "上面有些只需改寫，有些需要再跑訓練 / 測試。想跟老師確認要投入哪些——"
            "部分可能需要在 Nano4 / 5 儲值運算額度。",
        0.6, 1.35, 12.1, 0.5, size=13, color=INK_2, line_spacing=1.1)
    sl.shapes.add_picture(str(TMP / "compute_table.png"), Inches(0.35), Inches(1.95), Inches(12.6), Inches(4.7))
    txt(sl, "想法：儲值前先跟老師確認優先序，把預算集中在 P0（out-of-genome）最有價值；50M multi-seed 與 Bracken 重估相對便宜、可先做。",
        0.6, 6.7, 12.1, 0.5, size=11.5, color=ORANGE, bold=True)
    footer(sl)

    # 4 — decisions wanting advisor's steer
    sl = prs.slides.add_slide(L6)
    content_header(sl, "需要討論", "想請老師給方向的幾個決策", VIOLET)
    y = 1.5; hh = 1.15
    for title, body, color in [
        ("投稿定位要多往「評估準則 / protocol」靠嗎？",
         "BIB 偏 review / protocol；越往準則靠越貼合，但改寫幅度也越大。想聽老師覺得該走到哪個程度。", VIOLET),
        ("out-of-genome 這次投稿就補，還是先投、後續補？",
         "這是最關鍵的實驗，但也最花時間 / 運算。要納入這次投稿範圍，還是當作第一輪 revision 再補？", CRITICAL),
        ("若運算有限，先投入哪一兩個實驗？",
         "我的排序是 out-of-genome > 相同架構 pre-training 對照 > multi-seed；想確認是否符合老師的期待。", ORANGE),
        ("目標期刊與投稿時程？",
         "是否仍以 Briefings in Bioinformatics 為主要目標？希望抓一個目標投稿時間，好回推實驗排程。", BLUE),
    ]:
        key_row(sl, title, body, color, 0.6, y, 12.13, hh)
        y += hh + 0.12
    footer(sl)

    # 5 — submission logistics (light)
    sl = prs.slides.add_slide(L6)
    content_header(sl, "投稿前準備", "投稿前的行政待辦（我可以先弄草稿）", GREEN)
    bullet_card(sl, [
        ("資料 / 程式公開：建 Zenodo release + 固定 git commit + split manifest（BIB 要求 data availability 附 DOI）。", GREEN, 0),
        ("作者資訊：第二作者、ORCID、CRediT 分工、通訊資訊、每位作者約 30 字簡介。", GREEN, 0),
        ("字數：目前約 7,500 字，BIB Problem-solving Protocol 上限 5,000；把細節（RC-TTA、balancing、routing、計算成本）移到 Supplementary。", GREEN, 0),
        ("cover letter：主動揭露本文改寫自碩論、並說明差異（BIB 規定）。", GREEN, 0),
        ("以上多數我可以先做出骨架 / 草稿，老師確認方向後我再補完。", MUTED, 0),
    ], 0.6, 1.5, 12.13, 4.5, base_size=14)
    footer(sl)

    # 6 — what I hope to get from today
    sl = prs.slides.add_slide(L6)
    content_header(sl, "今天的目標", "今天想跟老師確認的三件事", INK)
    y = 1.65; hh = 1.35
    for title, body, color in [
        ("① 整體架構與定位", "這個 controlled-benchmark + tokenizer 主線的方向，老師覺得 OK 嗎？有沒有想調整的重點？", BLUE),
        ("② 要投入哪些額外實驗", "特別是 out-of-genome（P0）要不要這次就做；若要，是否在 Nano4 / 5 儲值運算額度。", ORANGE),
        ("③ 目標投稿時程", "抓一個目標時間，好回推實驗與改寫的排程。", GREEN),
    ]:
        key_row(sl, title, body, color, 0.6, y, 12.13, hh)
        y += hh + 0.18
    footer(sl)

    prs.save(OUT)
    print(f"saved: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()

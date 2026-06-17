#!/usr/bin/env python3
"""
Weekly meeting deck · 2026-06-19
Covers progress 2026-06-09 → 2026-06-19.
Key story: Scaling NT-v2 genus 50M → 258M did NOT improve accuracy —
           it dropped (66.6% → ~44.5%). A code+data audit found TWO
           compounding bugs:
           (1) WRONG DATASET — v10/v12/v13 used the HPC 3,507-species
               file instead of the correct CrucialX9 1,535-species version.
           (2) READER BUG — train_ddp.py _read_seq reads only the first
               line of each record; the 258M FASTA is 60-col wrapped, so
               every 150 bp read was truncated to 60 bp.
           The 66.6% v9 baseline is also inflated (test/train leakage).
           Fix: switch to correct CrucialX9 data (on Nano4 since 6/17) +
           patch _read_seq + rebuild clean test set.
"""

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle, FancyArrowPatch
import numpy as np
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DOCS = Path(__file__).parent
TMP  = DOCS / "meeting_0619_figs"
TMP.mkdir(exist_ok=True)
OUT  = DOCS / "meeting_2026_0619.pptx"

NAVY   = "#0A2540"
TEAL   = "#1E6091"
GREEN  = "#2E7D32"
RED    = "#C62828"
ORANGE = "#E65100"
GOLD   = "#B8860B"
GRAY   = "#5A6068"
LIGHT  = "#F4F6F8"
WHITE  = "#FFFFFF"
PURPLE = "#6A1B9A"

def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

plt.rcParams.update({"font.family":"DejaVu Sans","font.size":11,
                     "axes.spines.top":False,"axes.spines.right":False,
                     "figure.facecolor":"white","savefig.dpi":200,
                     "savefig.bbox":"tight","savefig.pad_inches":0.15})


# ═══════════════════════════════════════════════════════════════════════════════
# FIGURES
# ═══════════════════════════════════════════════════════════════════════════════

def fig_timeline():
    fig, ax = plt.subplots(figsize=(13, 5.0))
    ax.set_xlim(0, 13); ax.set_ylim(0, 5.5); ax.axis("off")
    ax.add_patch(Rectangle((0.3, 2.5), 12.4, 0.13, facecolor=NAVY))

    # (x, date, label, color, up=True means label above the line)
    events = [
        (0.8,  "6/9",  "Benchmark done\n(7 models)",        TEAL,   True),
        (2.4,  "6/11", "v10 258M fired\n(warm, weights=T)",  TEAL,   False),
        (4.0,  "6/13", "v10 COLLAPSE\ntest 40.1%",           RED,    True),
        (5.6,  "6/14", "1st guess:\nclass_weights",          ORANGE, False),
        (7.2,  "6/14", "v12 + v13 fired\n(weights=F)",       TEAL,   True),
        (8.8,  "6/15", "v13 done\n44.5%",                    GREEN,  False),
        (10.4, "6/18", "AUDIT: wrong\ndataset + 60bp\nreader bug",      RED,    True),
        (12.0, "6/19", "TODAY\nfix + re-test",               GOLD,   False),
    ]
    for x, date, label, color, up in events:
        ax.scatter(x, 2.565, s=200, color=color, zorder=4, edgecolor="white", linewidth=2)
        if up:
            ax.text(x, 3.05, date,  ha="center", fontsize=10, fontweight="bold", color=color)
            ax.text(x, 3.45, label, ha="center", fontsize=9,  color=NAVY, va="bottom")
            ax.plot([x, x], [2.62, 3.0], color=color, lw=1, alpha=0.5)
        else:
            ax.text(x, 2.10, date,  ha="center", fontsize=10, fontweight="bold", color=color)
            ax.text(x, 1.95, label, ha="center", fontsize=9,  color=NAVY, va="top")
            ax.plot([x, x], [2.50, 2.15], color=color, lw=1, alpha=0.5)

    ax.text(6.5, 5.2, "Progress Timeline · 6/9 → 6/19 (today)",
            ha="center", fontsize=13, fontweight="bold", color=NAVY)
    plt.savefig(TMP/"timeline.png", bbox_inches="tight"); plt.close()


def fig_paradox():
    """The scaling paradox: expected vs actual."""
    fig, axes = plt.subplots(1, 2, figsize=(13, 5))
    fig.patch.set_facecolor("white")

    # Left: what we expected
    ax = axes[0]
    ax.set_title("Expectation: 5× data → +2–5 pp", fontsize=12, fontweight="bold", color=NAVY)
    x = [50, 258]
    y_exp = [66.6, 70.0]
    ax.plot(x, y_exp, "o--", color=GREEN, lw=2.2, ms=11)
    ax.annotate("", xy=(258, 70.0), xytext=(50, 66.6),
                arrowprops=dict(arrowstyle="->", color=GREEN, lw=2, alpha=0.5))
    ax.text(150, 69.2, "+3.4 pp\n(hoped)", color=GREEN, fontsize=10, ha="center", fontweight="bold")
    ax.text(50, 65.2, "v9\n66.6%", ha="center", fontsize=9.5, color=NAVY)
    ax.text(258, 71.0, "?", ha="center", fontsize=13, color=GREEN, fontweight="bold")
    ax.set_xlim(0, 300); ax.set_ylim(35, 75)
    ax.set_xlabel("Training reads (M)"); ax.set_ylabel("Genus accuracy (%)")
    ax.set_xticks([50, 258])

    # Right: what actually happened
    ax2 = axes[1]
    ax2.set_title("Reality: 5× data → −22 pp", fontsize=12, fontweight="bold", color=RED)
    ax2.plot([50], [66.6], "o", color=GREEN, ms=12)
    ax2.text(50, 68.0, "v9 (50M balanced)\n66.6%", ha="center", fontsize=9.5, color=GREEN, fontweight="bold")
    # 258M points
    ax2.plot([258], [40.1], "v", color=RED, ms=13)
    ax2.text(258, 37.0, "v10 (weights=T)\n40.1% test", ha="center", fontsize=9, color=RED, fontweight="bold")
    ax2.plot([258], [45.4], "s", color=ORANGE, ms=11)
    ax2.text(268, 47.0, "v12 warm\n45.4%", ha="left", fontsize=8.5, color=ORANGE)
    ax2.plot([258], [44.5], "D", color=TEAL, ms=10)
    ax2.text(200, 42.0, "v13 scratch\n44.5%", ha="right", fontsize=8.5, color=TEAL)
    ax2.annotate("", xy=(258, 44.5), xytext=(50, 66.6),
                 arrowprops=dict(arrowstyle="->", color=RED, lw=2, alpha=0.45, linestyle="--"))
    ax2.set_xlim(0, 300); ax2.set_ylim(35, 75)
    ax2.set_xlabel("Training reads (M)"); ax2.set_ylabel("Genus accuracy (%)")
    ax2.set_xticks([50, 258])

    plt.tight_layout()
    plt.savefig(TMP/"paradox.png"); plt.close()


def fig_v10_collapse():
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.8))
    fig.patch.set_facecolor("white")

    # Left: v10 val_acc curve, flat ~21%
    ax = axes[0]
    epochs = list(range(1, 13))
    val = [19.8, 21.16, 20.4, 20.9, 21.0, 20.7, 21.1, 20.9, 21.18, 20.80, 21.18, 21.01]
    ax.plot(epochs, val, "o-", color=RED, lw=2, ms=6)
    ax.axhline(66.6, color=GREEN, linestyle="--", lw=1.6, alpha=0.7)
    ax.text(6, 64.0, "v9 baseline 66.6%", color=GREEN, fontsize=10, ha="center")
    ax.axhline(40.14, color=ORANGE, linestyle=":", lw=1.6, alpha=0.8)
    ax.text(9, 42.0, "v10 100K test = 40.1%", color=ORANGE, fontsize=9, ha="center")
    ax.set_ylim(15, 72); ax.set_xlim(0.5, 12.5)
    ax.set_xlabel("Epoch"); ax.set_ylabel("Val accuracy (%)")
    ax.set_title("v10: val_acc frozen ~21% (weights + 60 bp reads)",
                 fontsize=11.5, fontweight="bold", color=RED)

    # Right: root-cause box
    ax2 = axes[1]; ax2.axis("off"); ax2.set_xlim(0,1); ax2.set_ylim(0,1)
    ax2.text(0.5, 0.97, "Real Root Cause (code+data audit, 6/18)", ha="center",
             fontsize=13, fontweight="bold", color=NAVY, va="top")

    rows = [
        (RED,    "Wrong dataset — primary cause",
                 "v10/v12/v13 used HPC 3,507-species file; correct is\nCrucialX9 1,535-species (arrived on Nano4 2026-06-17)."),
        (ORANGE, "Reads truncated 150 → 60 bp — compounds the damage",
                 "train_ddp.py/_read_seq reads only 1st line; 258M FASTA\nis 60-col wrapped → model saw 40% of each read."),
        (TEAL,   "v9 used a different, correct path",
                 "train.py/_parse_fasta joins all lines (150 bp) on a\nsingle-line file. Same model, different reader."),
        (GREEN,  "Baseline 66.6% is leakage-inflated",
                 "100K 'test' seq_ids are 100% inside the 50M train set\n→ v9's number is not a clean generalization metric."),
    ]
    y = 0.84
    for color, title, body in rows:
        ax2.add_patch(FancyBboxPatch((0.03, y-0.115), 0.94, 0.135,
                                     boxstyle="round,pad=0.01",
                                     facecolor=color+"18", edgecolor=color, linewidth=1.3))
        ax2.text(0.07, y+0.005, title, fontsize=10.5, color=color, fontweight="bold", va="top")
        ax2.text(0.07, y-0.04, body,  fontsize=9.3, color=NAVY, va="top")
        y -= 0.165

    plt.tight_layout()
    plt.savefig(TMP/"v10_collapse.png"); plt.close()


def fig_experiments_table():
    fig, ax = plt.subplots(figsize=(13, 5.6))
    ax.axis("off"); ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.text(0.5, 0.97, "Four Experiments on NT-v2 Genus (LoRA, 120 genera)",
            ha="center", fontsize=14, fontweight="bold", color=NAVY, va="top")

    # columns: Run | Data | Init | class_w | Epochs | Result | Note
    cols   = [0.04, 0.14, 0.255, 0.42, 0.55, 0.645, 0.78]
    widths = ["Run", "Data", "Init", "class_w", "Epochs", "Result", "Verdict"]
    header_y = 0.86
    for cx, h in zip(cols, widths):
        ax.text(cx, header_y, h, fontsize=10.5, color=TEAL, fontweight="bold")
    ax.add_patch(Rectangle((0.02, header_y-0.03), 0.96, 0.004, facecolor=TEAL))

    rows = [
        ("v9",  "50M bal.\n(1535 sp.)",  "scratch", "false", "30", "66.6%", "baseline ✓", GREEN),
        ("v10", "258M†\n(HPC 3507sp)", "v9 warm", "TRUE",  "12*", "40.1%", "collapsed ✗", RED),
        ("v12", "258M†\n(HPC 3507sp)", "v9 warm", "false", "3**", "45.4%", "plateau", ORANGE),
        ("v13", "258M†\n(HPC 3507sp)", "scratch", "false", "25", "44.5%", "plateau", ORANGE),
    ]
    y = header_y - 0.075
    for run, data, init, cw, ep, res, verdict, color in rows:
        ax.add_patch(FancyBboxPatch((0.02, y-0.052), 0.96, 0.075,
                                    boxstyle="round,pad=0.004",
                                    facecolor=color+"12", edgecolor=color+"66", linewidth=1))
        vals = [run, data, init, cw, ep, res, verdict]
        for cx, v in zip(cols, vals):
            fw = "bold" if cx in (cols[0], cols[5], cols[6]) else "normal"
            c  = color if cx in (cols[5], cols[6]) else NAVY
            ax.text(cx, y, v, fontsize=10.5, color=c, fontweight=fw, va="center")
        y -= 0.105

    ax.text(0.04, 0.20, "*  v10 early-stopped at epoch 12/15 (val_acc plateau).",
            fontsize=8.8, color=GRAY)
    ax.text(0.04, 0.165, "** v12 only reached 3 epochs — cold NFS = 476 min/epoch hit the 23 h wall.",
            fontsize=8.8, color=GRAY)
    ax.text(0.04, 0.128, "†  All three 258M runs used the HPC old 3,507-species file —",
            fontsize=8.8, color=RED)
    ax.text(0.04, 0.093, "   not the correct CrucialX9 1,535-species data (arrived on Nano4 2026-06-17).",
            fontsize=8.8, color=RED)

    ax.add_patch(FancyBboxPatch((0.02, 0.01), 0.96, 0.065,
                                boxstyle="round,pad=0.01",
                                facecolor="#FFF8E1", edgecolor=GOLD, linewidth=1.6))
    ax.text(0.5, 0.055, "All three 258M runs: wrong dataset AND reader bug (60 bp) — both active simultaneously.",
            ha="center", fontsize=10.5, color=NAVY, fontweight="bold")
    ax.text(0.5, 0.025, "Explains the 22 pp gap and why init/weights made no difference.",
            ha="center", fontsize=9.5, color=GOLD)

    plt.savefig(TMP/"experiments_table.png", bbox_inches="tight"); plt.close()


def fig_key_finding():
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.8))
    fig.patch.set_facecolor("white")

    # Left: two compounding bugs as stacked boxes
    ax = axes[0]
    ax.axis("off"); ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.text(0.5, 0.97, "Two Bugs, Both Active in All 258M Runs",
            ha="center", fontsize=12, fontweight="bold", color=NAVY, va="top")

    # Bug 1
    ax.add_patch(FancyBboxPatch((0.03, 0.60), 0.94, 0.32,
                                 boxstyle="round,pad=0.02",
                                 facecolor=RED+"1A", edgecolor=RED, linewidth=2))
    ax.text(0.08, 0.895, "Bug 1  —  Wrong Dataset  (primary)",
            fontsize=11.5, fontweight="bold", color=RED, va="center")
    ax.text(0.08, 0.790, "Used:    HPC 3,507-species file",
            fontsize=10.5, color=RED, va="center")
    ax.text(0.08, 0.685, "Correct: CrucialX9 1,535-species  (on Nano4 since 6/17)",
            fontsize=10, color=NAVY, va="center")

    ax.text(0.5, 0.565, "+", ha="center", fontsize=20, color=GRAY, fontweight="bold")

    # Bug 2
    ax.add_patch(FancyBboxPatch((0.03, 0.20), 0.94, 0.32,
                                 boxstyle="round,pad=0.02",
                                 facecolor=ORANGE+"1A", edgecolor=ORANGE, linewidth=2))
    ax.text(0.08, 0.495, "Bug 2  —  Reader Truncation  (compounding)",
            fontsize=11.5, fontweight="bold", color=ORANGE, va="center")
    ax.text(0.08, 0.390, "train_ddp.py _read_seq reads only the 1st line",
            fontsize=10, color=NAVY, va="center")
    ax.text(0.08, 0.285, "258M FASTA is 60-col wrapped  →  150 bp fed as 60 bp to model",
            fontsize=10, color=ORANGE, va="center")

    # Bottom result box
    ax.add_patch(FancyBboxPatch((0.03, 0.02), 0.94, 0.135,
                                 boxstyle="round,pad=0.01",
                                 facecolor="#FFF3E0", edgecolor=GOLD, linewidth=1.5))
    ax.text(0.5, 0.115, "Both bugs active in ALL v10/v12/v13 → explains the full 22 pp gap",
            ha="center", fontsize=10.5, color=NAVY, fontweight="bold")
    ax.text(0.5, 0.060, "CrucialX9 is also 60-col wrapped — must fix reader after switching data",
            ha="center", fontsize=9.5, color=GRAY)

    # Right: takeaway text
    ax2 = axes[1]; ax2.axis("off"); ax2.set_xlim(0, 1); ax2.set_ylim(0, 1)
    ax2.text(0.5, 0.97, "What This Means", ha="center",
             fontsize=13, fontweight="bold", color=NAVY, va="top")

    points = [
        (RED,    "Wrong dataset invalidates all 258M runs",
                 "3,507 sp. HPC vs 1,535 sp. test — different genus\ndistribution; model never saw the right label space."),
        (ORANGE, "Reader bug compounds the damage",
                 "_read_seq reads 1 line; CrucialX9 is also 60-col\nwrapped → must fix regardless of dataset swap."),
        (TEAL,   "The 66.6% baseline is leaky",
                 "100K test ⊂ 50M train (100% overlap) → v9 is\ninflated; honest baseline will be lower."),
        (GREEN,  "Correct data is on Nano4 — fix is ready",
                 "CrucialX9 1,535 sp. at …/labeled_multi_level_1535sp/\nPatch _read_seq + rebuild test set → re-run."),
    ]
    y = 0.85
    for color, title, body in points:
        ax2.add_patch(FancyBboxPatch((0.03, y-0.115), 0.94, 0.135,
                                     boxstyle="round,pad=0.01",
                                     facecolor=color+"18", edgecolor=color, linewidth=1.3))
        ax2.text(0.07, y+0.005, title, fontsize=10.5, color=color, fontweight="bold", va="top")
        ax2.text(0.07, y-0.04, body,  fontsize=9.3, color=NAVY, va="top")
        y -= 0.168

    plt.tight_layout()
    plt.savefig(TMP/"key_finding.png"); plt.close()


def fig_infra():
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.8))
    fig.patch.set_facecolor("white")

    # Left: NFS cold vs warm bar (the 10x lesson)
    ax = axes[0]
    labels = ["Cold NFS\n(v12 nodes)", "Warm NFS\n(v13 nodes)"]
    vals   = [476, 43]
    colors = [RED, GREEN]
    bars = ax.bar(labels, vals, color=colors, width=0.55)
    for b, v in zip(bars, vals):
        ax.text(b.get_x()+b.get_width()/2, v+8, f"{v} min/epoch", ha="center",
                fontsize=11, fontweight="bold", color=colors[0] if v>100 else GREEN)
    ax.set_ylim(0, 540); ax.set_ylabel("Minutes per epoch")
    ax.set_title("NFS cache → 11× epoch-time swing", fontsize=12, fontweight="bold", color=NAVY)
    ax.text(0.5, 260, "11×  faster", ha="center", fontsize=18, color=GRAY, fontweight="bold")

    # Right: infra facts
    ax2 = axes[1]; ax2.axis("off"); ax2.set_xlim(0,1); ax2.set_ylim(0,1)
    ax2.text(0.5, 0.97, "DDP Infrastructure Notes", ha="center",
             fontsize=13, fontweight="bold", color=NAVY, va="top")
    rows = [
        ("Scale",     "8 nodes × 8 H200 = 64 GPUs"),
        ("Batch",     "128 / GPU × 64 = 8192 effective"),
        ("Bottleneck","Lazy FASTA index over NFS (47 GB, 258M reads)"),
        ("Cold vs warm","476 vs 43 min/epoch — depends on node landing"),
        ("Quota",     "MST114414 + MST114550 juggled (MaxGRES/account)"),
        ("Wall limit","23 h soft cap → auto-resume from last.pt"),
    ]
    y = 0.84
    for k, v in rows:
        ax2.text(0.05, y, k, fontsize=10.3, color=TEAL, fontweight="bold", va="center")
        ax2.text(0.34, y, v, fontsize=10.0, color=NAVY, va="center")
        y -= 0.115
    ax2.add_patch(FancyBboxPatch((0.03, 0.02), 0.94, 0.085,
                                 boxstyle="round,pad=0.01",
                                 facecolor="#E8F5E9", edgecolor=GREEN, linewidth=1.5))
    ax2.text(0.5, 0.073, "Lesson: pre-warm / pin the FASTA index before large DDP runs.",
             ha="center", fontsize=9.8, color=GREEN, fontweight="bold")
    ax2.text(0.5, 0.04, "Cold-cache runs waste 10× the GPU-hours for the same epoch.",
             ha="center", fontsize=9, color=NAVY)

    plt.tight_layout()
    plt.savefig(TMP/"infra.png"); plt.close()


def fig_next_steps():
    fig, ax = plt.subplots(figsize=(13, 5))
    ax.axis("off"); ax.set_xlim(0,1); ax.set_ylim(0,1)
    ax.text(0.5, 0.97, "Next Steps · Fix the reader, then re-test", ha="center",
            fontsize=14, fontweight="bold", color=NAVY, va="top")

    steps = [
        ("Immediate fixes (two bugs + leaky baseline)",  RED, [
            ("Switch dataset",   "Use CrucialX9 1,535 sp. (/work/…/labeled_multi_level_1535sp/) — already on Nano4 since 6/17"),
            ("Patch _read_seq",  "Read all lines until next '>' — CrucialX9 is also 60-col wrapped; verify reads = 150 bp"),
            ("Clean test set",   "Rebuild 100K test with ZERO seq_id overlap with training; re-baseline v9 honestly"),
        ]),
        ("Re-run with correct data + full 150 bp reads", ORANGE, [
            ("258M full-length", "Re-train correct 258M (weights=false, fixed reader) — first real test of VOLUME"),
            ("v14 250M balanced","subsample_balanced.py (job 118820, single-line output) as controlled BALANCE comparison"),
        ]),
        ("Only then conclude", GREEN, [
            ("Volume",          "Does full-read 258M beat the clean v9 baseline? (finally a controlled answer)"),
            ("Balance",         "v14 vs full-read 258M isolates balance — meaningful only after the reader fix"),
        ]),
    ]
    y = 0.86
    for section, color, items in steps:
        ax.text(0.02, y, section, fontsize=11.5, color=color, fontweight="bold")
        y -= 0.05
        for label, detail in items:
            ax.add_patch(FancyBboxPatch((0.02, y-0.065), 0.96, 0.07,
                                        boxstyle="round,pad=0.01",
                                        facecolor=color+"15", edgecolor=color+"55", linewidth=1))
            ax.text(0.05, y-0.015, label, fontsize=10, color=color, fontweight="bold")
            ax.text(0.27, y-0.015, detail, fontsize=9.3, color=NAVY)
            y -= 0.085
        y -= 0.025

    plt.tight_layout()
    plt.savefig(TMP/"next_steps.png"); plt.close()


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def add_slide(prs, layout=6):
    return prs.slides.add_slide(prs.slide_layouts[layout])

def title_bar(slide, title, subtitle=""):
    W, H = Inches(13.33), Inches(7.5)
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(NAVY)
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = rgb(WHITE)
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(14); r2.font.color.rgb = rgb("#90CAF9")

def add_image(slide, path, left, top, width, height):
    slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                              Inches(width), Inches(height))

def add_text_box(slide, text, left, top, width, height,
                 fontsize=12, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fontsize); r.font.bold = bold; r.font.color.rgb = rgb(color)


def build():
    print("Generating figures...")
    fig_timeline()
    fig_paradox()
    fig_v10_collapse()
    fig_experiments_table()
    fig_key_finding()
    fig_infra()
    fig_next_steps()

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    sl = add_slide(prs)
    bg = sl.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = rgb(NAVY); bg.line.fill.background()

    add_text_box(sl, "Weekly Progress Update", 1, 1.5, 11, 1,
                 fontsize=38, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, "Scaling NT-v2 Genus 50M → 258M · debugging a 22 pp regression",
                 1, 2.7, 11, 0.6, fontsize=16, color="#90CAF9", align=PP_ALIGN.CENTER)
    add_text_box(sl, "June 19, 2026", 1, 3.4, 11, 0.5,
                 fontsize=14, color="#BBDEFB", align=PP_ALIGN.CENTER)

    bullets = "5× data → −22 pp  ·  wrong dataset (3507 sp.) + reader bug (60 bp)  ·  baseline was leaky"
    add_text_box(sl, bullets, 1, 4.2, 11, 0.6,
                 fontsize=13, color=GOLD, align=PP_ALIGN.CENTER)

    add_text_box(sl, "YMJ · M2 Lab · NTU CSIE", 1, 6.6, 11, 0.5,
                 fontsize=11, color="#78909C", align=PP_ALIGN.CENTER)

    # ── Slide 2: Timeline ─────────────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "Progress Timeline · 6/9 → 6/19")
    add_image(sl, TMP/"timeline.png", 0.2, 1.1, 12.9, 3.8)
    rows = [
        ("6/11", "Fired v10 — 258M, warm-start from v9, class_weights=TRUE"),
        ("6/13", "v10 collapsed: 100K test 40.1% (vs v9 66.6%)"),
        ("6/14", "First guess: class_weights → fired v12 (warm) + v13 (scratch), both weights=false"),
        ("6/15–18", "Both plateau ~44.5% → audit found TWO bugs: wrong dataset (HPC 3507sp) + reads truncated 150→60 bp"),
    ]
    y = 5.05
    for date, text in rows:
        add_text_box(sl, f"  {date}   {text}", 0.3, y, 12.7, 0.35, fontsize=11, color=NAVY)
        y += 0.37

    # ── Slide 3: The Paradox ──────────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "The Scaling Paradox", "We expected more data to help. It hurt.")
    add_image(sl, TMP/"paradox.png", 0.2, 1.2, 12.9, 5.2)

    # ── Slide 4: v10 Collapse + Root Cause ────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "v10 Collapse → Real Root Cause", "Two bugs: wrong dataset (HPC 3507 sp.) + reads truncated 150→60 bp")
    add_image(sl, TMP/"v10_collapse.png", 0.2, 1.2, 12.9, 5.2)

    # ── Slide 5: Experiments Table ────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "Four Experiments, One Conclusion")
    add_image(sl, TMP/"experiments_table.png", 0.2, 1.1, 12.9, 5.9)

    # ── Slide 6: Key Finding ──────────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "Key Finding: Two Compounding Bugs", "Wrong dataset (HPC 3507 sp.) + reader truncation — fix ready before any volume/balance claim")
    add_image(sl, TMP/"key_finding.png", 0.2, 1.2, 12.9, 5.2)

    # ── Slide 7: Infrastructure ───────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "Infrastructure Lessons (64-GPU DDP)")
    add_image(sl, TMP/"infra.png", 0.2, 1.2, 12.9, 5.2)

    # ── Slide 8: Next Steps ───────────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "Next Steps · Fix the Reader, then Re-test", "Patch _read_seq + clean test set before any volume/balance claim")
    add_image(sl, TMP/"next_steps.png", 0.2, 1.1, 12.9, 5.9)

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()

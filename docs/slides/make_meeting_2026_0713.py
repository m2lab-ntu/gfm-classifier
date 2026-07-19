#!/usr/bin/env python3
"""
Weekly meeting deck · 2026-07-13
Covers progress 2026-07-03 (BIB submission direction confirmed) -> 2026-07-13 (today).

Story arc:
  1) Confirmed direction: reposition thesis as a BIB Problem-solving Protocol;
     identified the two most dangerous reviewer objections in advance
     (incomplete Kraken2/Bracken baseline; no real-community validation).
  2) Drafted manuscript + fixed build/citation issues; added a 6-mer Naive
     Bayes control to isolate k-mer length as a confound.
  3) HONEST CORRECTION: added Kraken2+Bracken baseline -> the "neural beats
     Kraken2 on abundance" claim was an artefact of an incomplete baseline
     (raw Kraken2 r=0.823 -> Bracken r=0.997). Retracted and reframed Sec 6.2.
  4) Real-world trial: ZymoBIOMICS D6331 mock community shows ALL methods
     collapse from simulated r~0.99 to r=0.42-0.58 with no clear winner --
     the classical pipeline's simulated dominance does not transfer.
  5) Narrative alignment + final visual polish (today): Fig 1 redesign,
     compute figure, Box 1 enrichment, first-author info filled in.
"""

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from matplotlib.patches import Rectangle, FancyBboxPatch
import numpy as np
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DOCS = Path(__file__).parent
TMP = DOCS / "meeting_0713_figs"
TMP.mkdir(exist_ok=True)
OUT = DOCS / "meeting_2026_0713.pptx"

NAVY = "#0A2540"
TEAL = "#1E6091"
GREEN = "#2E7D32"
RED = "#C62828"
ORANGE = "#E65100"
GOLD = "#B8860B"
GRAY = "#5A6068"
LIGHT = "#F4F6F8"
WHITE = "#FFFFFF"
PURPLE = "#6A1B9A"

CJK_BOLD = "/work/ymj1123ntu/NTU-Thesis-Writing-Template/fonts/chinese/cwTeX_Hei.ttf"
CJK_REG = "/work/ymj1123ntu/NTU-Thesis-Writing-Template/fonts/chinese/MOE_Song.ttf"
fp_bold = fm.FontProperties(fname=CJK_BOLD)
fp_reg = fm.FontProperties(fname=CJK_REG)

CJK_FONT_NAME = "Microsoft JhengHei"  # rendered by whatever app opens the pptx


def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


plt.rcParams.update({
    "axes.spines.top": False, "axes.spines.right": False,
    "figure.facecolor": "white", "savefig.dpi": 200,
    "savefig.bbox": "tight", "savefig.pad_inches": 0.15,
})


# ═══════════════════════════════════════════════════════════════════════
# FIGURES
# ═══════════════════════════════════════════════════════════════════════

def fig_timeline():
    fig, ax = plt.subplots(figsize=(13, 4.9))
    ax.set_xlim(0, 13); ax.set_ylim(0, 5.6); ax.axis("off")
    ax.add_patch(Rectangle((0.3, 2.55), 12.4, 0.13, facecolor=NAVY))

    events = [
        (0.8, "7/3", "確認投稿方向：\nBIB Problem-solving\nProtocol，初稿", TEAL, True),
        (2.8, "7/4-5", "修編譯/引用錯誤\n統一圖表字體", GRAY, False),
        (4.6, "7/5", "6-mer NB control\n分離 k-mer 長度\n混淆變數", TEAL, True),
        (6.3, "7/6", "依 reviewer 觀點\n大改：scope、\nrelated work", GRAY, False),
        (8.0, "7/8", "誠實修正：\n豐度宣稱撤回\n(Bracken r=0.997)", RED, True),
        (9.9, "7/9-11", "真實 mock\ncommunity 驗證\n上線 (D6331)", ORANGE, False),
        (11.9, "7/13", "今天：敘事對齊\n+ 圖表/作者\n資訊定稿", GOLD, True),
    ]
    for x, date, label, color, up in events:
        ax.scatter(x, 2.615, s=220, color=color, zorder=4, edgecolor="white", linewidth=2)
        if up:
            ax.text(x, 3.05, date, ha="center", fontsize=11, fontweight="bold", color=color)
            ax.text(x, 3.45, label, ha="center", fontsize=9.5, color=NAVY, va="bottom", fontproperties=fp_reg)
            ax.plot([x, x], [2.68, 3.0], color=color, lw=1, alpha=0.5)
        else:
            ax.text(x, 2.10, date, ha="center", fontsize=11, fontweight="bold", color=color)
            ax.text(x, 1.95, label, ha="center", fontsize=9.5, color=NAVY, va="top", fontproperties=fp_reg)
            ax.plot([x, x], [2.50, 2.20], color=color, lw=1, alpha=0.5)

    ax.text(6.5, 5.3, "本週進度時間軸 · 7/3（確認方向）→ 7/13（今天）",
            ha="center", fontsize=14, fontweight="bold", color=NAVY, fontproperties=fp_bold)
    plt.savefig(TMP / "timeline.png"); plt.close()


def fig_abundance_correction():
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.8))

    # Left: before/after Bracken correction
    ax = axes[0]
    methods = ["Kraken2\n(raw)", "Kraken2\n+Bracken", "NT-v2", "MT 13-mer"]
    before = [0.823, None, 0.992, 0.999]
    colors = [RED, GREEN, TEAL, NAVY]
    r_vals = [0.823, 0.997, 0.992, 0.999]
    bars = ax.bar(methods, r_vals, color=colors, width=0.55, zorder=3)
    ax.set_ylim(0, 1.15)
    ax.set_ylabel("Pearson r (genus abundance)", fontsize=11)
    ax.axhline(1.0, color=GRAY, lw=0.8, ls="--", alpha=0.5)
    for b, v in zip(bars, r_vals):
        ax.text(b.get_x() + b.get_width()/2, v + 0.03, f"{v:.3f}",
                ha="center", fontsize=11, fontweight="bold")
    ax.annotate("+0.174\n(Bracken 修正)", xy=(1, 0.997), xytext=(1.55, 0.55),
                fontsize=10, color=RED, fontweight="bold", ha="center",
                fontproperties=fp_bold,
                arrowprops=dict(arrowstyle="->", color=RED, lw=1.5))
    ax.set_xlim(-0.6, 3.6)
    ax.set_title("撤回前 vs 撤回後：豐度估計基準", fontsize=13, fontweight="bold",
                 color=NAVY, fontproperties=fp_bold)
    ax.grid(axis="y", alpha=0.25)

    # Right: narrative box
    ax2 = axes[1]
    ax2.axis("off")
    ax2.set_xlim(0, 10); ax2.set_ylim(0, 10)
    box_before = FancyBboxPatch((0.3, 5.6), 9.4, 3.6, boxstyle="round,pad=0.15",
                                 facecolor="#FFEBEE", edgecolor=RED, linewidth=1.5)
    ax2.add_patch(box_before)
    ax2.text(5, 8.7, "撤回前的說法（草稿）", fontsize=12, fontweight="bold",
              color=RED, ha="center", fontproperties=fp_bold)
    ax2.text(5, 7.6,
             "NT-v2 豐度估計 r=0.992\n勝過 Kraken2 raw r=0.823\n→「Neural 模型勝過 Kraken2」",
             fontsize=10.5, ha="center", va="center", color=NAVY, fontproperties=fp_reg)

    box_after = FancyBboxPatch((0.3, 0.6), 9.4, 4.4, boxstyle="round,pad=0.15",
                                facecolor="#E8F5E9", edgecolor=GREEN, linewidth=1.5)
    ax2.add_patch(box_after)
    ax2.text(5, 4.6, "本週補 Bracken baseline 後", fontsize=12, fontweight="bold",
              color=GREEN, ha="center", fontproperties=fp_bold)
    ax2.text(5, 2.4,
             "Kraken2+Bracken r=0.997\n(raw 0.823 只是 30% 棄權\n造成的可修復假象)\n→ neural 無豐度優勢，\n賣點改放 read-level 準確度\n+ out-of-database 穩健性",
             fontsize=10.5, ha="center", va="center", color=NAVY, fontproperties=fp_reg)

    plt.savefig(TMP / "abundance_correction.png"); plt.close()


def fig_mock_community():
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.9))

    # Left: simulated vs real collapse
    ax = axes[0]
    methods = ["Kraken2+Bracken", "NT-v2", "MT 13-mer"]
    sim_r = [0.997, 0.992, 0.999]
    real_r = [0.580, 0.420, 0.545]
    x = np.arange(len(methods)); w = 0.32
    ax.bar(x - w/2, sim_r, width=w, color=TEAL, label="模擬數據 (in-DB)", zorder=3)
    ax.bar(x + w/2, real_r, width=w, color=ORANGE, label="真實 mock community (D6331)", zorder=3)
    ax.set_xticks(x); ax.set_xticklabels(methods, fontsize=10.5)
    ax.set_ylabel("Pearson r (genus abundance)", fontsize=11)
    ax.set_ylim(0, 1.32)
    ax.legend(fontsize=10, prop=fp_reg, loc="upper center", bbox_to_anchor=(0.5, 1.19), ncol=2, frameon=False)
    ax.set_title("模擬 → 真實：全面崩落，沒有明顯贏家", fontsize=13, fontweight="bold",
                 color=NAVY, fontproperties=fp_bold, pad=45)
    for xi, sv, rv in zip(x, sim_r, real_r):
        ax.text(xi - w/2, sv + 0.03, f"{sv:.3f}", ha="center", fontsize=9)
        ax.text(xi + w/2, rv + 0.03, f"{rv:.3f}", ha="center", fontsize=9)
    ax.grid(axis="y", alpha=0.25)

    # Right: trade-off table-like summary
    ax2 = axes[1]
    ax2.axis("off"); ax2.set_xlim(0, 10); ax2.set_ylim(0, 10)
    ax2.text(5, 9.4, "真實 D6331 上的三方比較", fontsize=13, fontweight="bold",
              color=NAVY, ha="center", fontproperties=fp_bold)
    rows = [
        ("Kraken2+Bracken", "r=0.580（最佳）", "BC=0.437（最差）", "敏感度 54.5%（最低）"),
        ("NT-v2", "r=0.420", "—", "敏感度 72.7%，100% 有分類"),
        ("MT 13-mer", "r=0.545", "BC=0.344（最佳）", "敏感度 72.7%，100% 有分類"),
    ]
    y0 = 8.0
    for name, a, b, c in rows:
        ax2.text(0.2, y0, name, fontsize=11, fontweight="bold", color=NAVY, fontproperties=fp_reg)
        ax2.text(0.2, y0 - 0.55, f"  {a}   {b}   {c}", fontsize=10, color=GRAY, fontproperties=fp_reg)
        y0 -= 1.55
    box = FancyBboxPatch((0.2, 0.3), 9.6, 2.7, boxstyle="round,pad=0.15",
                          facecolor="#FFF8E1", edgecolor=GOLD, linewidth=1.5)
    ax2.add_patch(box)
    ax2.text(5, 1.65,
             "Kraken2 在真實 reads 上棄權率\n30%（模擬）→ 56%（真實）\nclassical pipeline 的模擬優勢\n無法直接轉移到真實世界",
             fontsize=10.5, ha="center", va="center", color=NAVY, fontweight="bold",
             fontproperties=fp_bold)

    plt.savefig(TMP / "mock_community.png"); plt.close()


# ═══════════════════════════════════════════════════════════════════════
# PPTX HELPERS
# ═══════════════════════════════════════════════════════════════════════

def add_slide(prs, layout=6):
    return prs.slides.add_slide(prs.slide_layouts[layout])


def title_bar(slide, title, subtitle=""):
    W = Inches(13.33)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(NAVY)
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = rgb(WHITE)
    r.font.name = CJK_FONT_NAME
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(13); r2.font.color.rgb = rgb("#90CAF9")
        r2.font.name = CJK_FONT_NAME


def add_image(slide, path, left, top, width, height):
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))


def add_text_box(slide, text, left, top, width, height,
                  fontsize=12, color=NAVY, bold=False, align=PP_ALIGN.LEFT, bg=None, border=None):
    if bg is not None or border is not None:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                      Inches(width), Inches(height))
        if bg is not None:
            shp.fill.solid(); shp.fill.fore_color.rgb = rgb(bg)
        else:
            shp.fill.background()
        if border is not None:
            shp.line.color.rgb = rgb(border); shp.line.width = Pt(1.25)
        else:
            shp.line.fill.background()
        tf = shp.text_frame
    else:
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        r = p.add_run(); r.text = line
        r.font.size = Pt(fontsize); r.font.bold = bold
        r.font.color.rgb = rgb(color)
        r.font.name = CJK_FONT_NAME


def bullets(slide, items, left, top, width, height, fontsize=14, color=NAVY, gap=Pt(6)):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        if isinstance(item, tuple):
            text, kwargs = item
        else:
            text, kwargs = item, {}
        r = p.add_run(); r.text = text
        r.font.size = Pt(kwargs.get("fontsize", fontsize))
        r.font.bold = kwargs.get("bold", False)
        r.font.color.rgb = rgb(kwargs.get("color", color))
        r.font.name = CJK_FONT_NAME
        p.level = kwargs.get("level", 0)


# ═══════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════

def build():
    print("Generating figures...")
    fig_timeline()
    fig_abundance_correction()
    fig_mock_community()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    sl = add_slide(prs)
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = rgb(NAVY); bg.line.fill.background()
    add_text_box(sl, "本週進度報告", 1, 1.3, 11.3, 1.1, fontsize=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, "GFM Metagenomic Benchmark → Briefings in Bioinformatics 投稿",
                 1, 2.5, 11.3, 0.6, fontsize=17, color="#90CAF9", align=PP_ALIGN.CENTER)
    add_text_box(sl, "2026-07-03（確認投稿方向） → 2026-07-13（今天）",
                 1, 3.15, 11.3, 0.5, fontsize=14, color="#BBDEFB", align=PP_ALIGN.CENTER)
    add_text_box(sl,
                 "確認方向 → 補嚴謹度 → 誠實修正豐度宣稱 → 真實群落驗證 → 手稿定稿潤色",
                 1, 4.0, 11.3, 0.6, fontsize=13.5, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(sl, "楊明儒 · 台大生醫電資所 · 2026-07-13", 1, 6.6, 11.3, 0.5,
                 fontsize=11, color="#78909C", align=PP_ALIGN.CENTER)

    # ── Slide 2: Timeline ──
    sl = add_slide(prs)
    title_bar(sl, "本週時間軸總覽", "從確認投稿方向到手稿定稿潤色")
    add_image(sl, TMP / "timeline.png", 0.2, 1.1, 12.9, 4.75)
    rows = [
        "7/3  確定重新定位為 BIB Problem-solving Protocol，同時點名兩個最危險的 reviewer 質疑",
        "7/5-6  補 6-mer NB control 隔離混淆變數；依 reviewer 觀點大改 scope / related work",
        "7/8  誠實修正：發現豐度宣稱的 baseline 不夠強，主動撤回並補上 Kraken2+Bracken",
        "7/9-11  真實 ZymoBIOMICS mock community (D6331) 三方比較上線",
        "7/13  敘事全文對齊，Figure 1 / compute 圖 / Box 1 潤色，第一作者資訊填入",
    ]
    y = 6.0
    for row in rows:
        add_text_box(sl, row, 0.35, y, 12.6, 0.32, fontsize=11.5, color=NAVY)
        y += 0.3

    # ── Slide 3: 起點 - 確認投稿方向 ──
    sl = add_slide(prs)
    title_bar(sl, "起點：確認投稿方向", "2026-07-03 — 為什麼是 BIB，為什麼要改寫成 Problem-solving Protocol")
    bullets(sl, [
        ("結論：碩論以「thesis 形式」直接投 BIB 不適合", dict(bold=True, color=RED, fontsize=15)),
        ("BIB 偏 review 期刊、不收 pure original research", dict(level=1, fontsize=13, color=GRAY)),
        ("對策：重新定位為 Problem-solving Protocol（2,000–5,000 字）", dict(bold=True, color=GREEN, fontsize=15)),
        ("敘事從「我提出 NT-v2+LoRA 方法」→「回答方法論問題：GFM 該如何被評估用於 "
         "metagenomic taxonomic classification」", dict(level=1, fontsize=13, color=GRAY)),
        ("選定標題：How should genomic foundation models be evaluated for metagenomic "
         "taxonomic classification?", dict(level=1, fontsize=12.5, color=TEAL)),
        ("投稿前主動點名的兩大風險（本週的工作正是逐一處理它們）", dict(bold=True, color=ORANGE, fontsize=15)),
        ("① Kraken2+Bracken baseline 不完整 —— 沒用最強 baseline 會被 reviewer 抓到",
         dict(level=1, fontsize=13, color=NAVY)),
        ("② 沒測真實 metagenomics —— 只有 ART 模擬數據，對 Kraken2 有結構性優勢",
         dict(level=1, fontsize=13, color=NAVY)),
    ], 0.5, 1.25, 12.3, 5.6)

    # ── Slide 4: 初稿 + 技術地基修復 ──
    sl = add_slide(prs)
    title_bar(sl, "初稿完成 + 技術地基修復", "2026-07-03 → 2026-07-06")
    bullets(sl, [
        ("7/3  Initial BIB manuscript draft（OUP Problem-solving Protocol 版型）", dict(fontsize=14)),
        ("Repo 建立 + GitHub Action 自動編譯（本機無 TeX，push 觸發，釘 TeX Live 2024）",
         dict(level=1, fontsize=12.5, color=GRAY)),
        ("7/4  修編譯錯誤：society logo 空白、tabular* 表格、CI 版本問題", dict(fontsize=14)),
        ("統一全部圖表字體（DejaVu Sans）、逐一核對數據來源、修圖說不一致",
         dict(level=1, fontsize=12.5, color=GRAY)),
        ("7/5  修正引用格式：numbered [1,3–5]（unsrtnat + sort&compress）", dict(fontsize=14)),
        ("7/6  依 reviewer 觀點大改：scope、related work、protocol/decision boxes、reproducibility",
         dict(fontsize=14)),
    ], 0.5, 1.25, 12.3, 5.4)

    # ── Slide 5: 6-mer NB control ──
    sl = add_slide(prs)
    title_bar(sl, "補強因果嚴謹度：6-mer Naive Bayes Control", "2026-07-05 — 隔離 k-mer 長度這個混淆變數")
    bullets(sl, [
        ("動機：先前只比較過「13-mer NB 74.9% vs NT-v2 6-mer 67%」，k-mer 長度本身是個沒控制的混淆變數",
         dict(fontsize=14, color=NAVY)),
        ("補上 6-mer NB control（同方法、公平比較，只換 k 值）", dict(bold=True, fontsize=15, color=TEAL)),
        ("結果：NB 13-mer 74.9%  vs  NB 6-mer 25.9%", dict(bold=True, fontsize=16, color=GREEN)),
        ("→ k-mer 長度才是決定性因子，不是「Naive Bayes 這個方法比較強」",
         dict(level=1, fontsize=13, color=GRAY)),
        ("鞏固論文核心論點：overlapping 13-mer 13-mer 87.4%（MT）vs NT-v2 6-mer 67.1% 的差距，"
         "來自 tokenization（k-mer 長度），不是 backbone pre-training", dict(fontsize=14, color=NAVY)),
    ], 0.5, 1.25, 12.3, 5.4)

    # ── Slide 6: 誠實修正 - 豐度宣稱撤回 ──
    sl = add_slide(prs)
    title_bar(sl, "誠實轉折：豐度宣稱撤回", "2026-07-08 — 補上最強 baseline 後，主動撤回先前結論")
    add_image(sl, TMP / "abundance_correction.png", 0.3, 1.15, 12.7, 4.7)
    add_text_box(sl,
                 "呼應鐵律「不確定就驗，不硬凹」：先自己補上最強 baseline（Kraken2+Bracken），"
                 "再誠實地把「neural 勝過 Kraken2」的宣稱改成「read-level 準確度 + out-of-database 穩健性」",
                 0.3, 5.95, 12.7, 0.9, fontsize=12, color=RED, bold=True)

    # ── Slide 7: 真實 mock community ──
    sl = add_slide(prs)
    title_bar(sl, "真實世界驗證：ZymoBIOMICS Mock Community D6331", "2026-07-09 → 07-11 — 回應「只測模擬數據」的風險")
    add_image(sl, TMP / "mock_community.png", 0.2, 1.1, 12.9, 4.75)
    add_text_box(sl,
                 "關鍵發現：classical pipeline 在模擬（in-DB）數據上的壓倒性優勢是模擬數據的假象，"
                 "無法直接轉移到真實 Illumina reads —— 新增 Sec 6.4",
                 0.2, 5.95, 12.9, 0.9, fontsize=12, color=ORANGE, bold=True)

    # ── Slide 8: 敘事對齊 + 最終潤色 ──
    sl = add_slide(prs)
    title_bar(sl, "敘事對齊 + 最終視覺潤色", "2026-07-13（今天）")
    bullets(sl, [
        ("全文敘事對齊：修掉「mock/Bracken 仍是 future work」等過時說法，統一 87.4% 的四捨五入，"
         "把模擬 vs 真實群落的建議分開講", dict(fontsize=14)),
        ("Figure 1（benchmark design）重畫：修好右側箭頭錯位/重疊，語意修正為 "
         "reads → 逐一改變因子 → 三層級評估 → Kraken2+Bracken", dict(fontsize=14)),
        ("rc_tta 圖：內部版本代號（v3/v9/v11...）換成有意義的設定名稱", dict(fontsize=14)),
        ("Supplementary 新增統一 compute 圖（throughput + peak GPU memory）", dict(fontsize=14)),
        ("Box 1 protocol checklist 8 步驟，每步都補上 rationale/pitfall", dict(fontsize=14)),
        ("填入第一作者（楊明儒）與單位（台大生醫電資所）", dict(fontsize=14)),
    ], 0.5, 1.25, 12.3, 5.4)

    # ── Slide 9: 目前狀態與待辦 ──
    sl = add_slide(prs)
    title_bar(sl, "目前狀態與待辦")
    bullets(sl, [
        ("手稿現況：main 12 頁 + supplementary 4 頁，0 undefined refs，TeX Live 2024 全部乾淨編譯",
         dict(bold=True, color=GREEN, fontsize=15)),
        ("本週已補齊投稿前識別的兩大風險", dict(bold=True, fontsize=15, color=NAVY)),
        ("✅ Kraken2+Bracken baseline 完整", dict(level=1, fontsize=13, color=GREEN)),
        ("✅ 真實 mock community 驗證", dict(level=1, fontsize=13, color=GREEN)),
        ("尚待處理", dict(bold=True, fontsize=15, color=ORANGE)),
        ("P0：out-of-genome generalization 實驗（13-mer 98.7% 被質疑是 closed-set k-mer lookup 的最大風險，"
         "§8 已主動揭露但尚未做實驗）", dict(level=1, fontsize=13, color=RED)),
        ("待補作者資訊：第二作者、ORCID、funding、Zenodo DOI、CRediT initials、acknowledgments",
         dict(level=1, fontsize=13, color=NAVY)),
        ("更完整真實群落驗證：多個 community、matched read length、Centrifuge 等",
         dict(level=1, fontsize=13, color=NAVY)),
    ], 0.5, 1.25, 12.3, 5.6)

    # ── Slide 10: 總結 ──
    sl = add_slide(prs)
    title_bar(sl, "本週總結：四個關鍵訊息")
    bullets(sl, [
        ("① 確認方向：重新定位為 BIB Problem-solving Protocol，並提前鎖定兩個最危險的 reviewer 質疑",
         dict(fontsize=15, bold=True, color=NAVY)),
        ("② 補嚴謹度：6-mer NB control 把「NB 比較強」收斂成「k-mer 長度才是關鍵」",
         dict(fontsize=15, bold=True, color=TEAL)),
        ("③ 誠實自我修正：發現先前豐度宣稱只是 baseline 不夠強，主動撤回並重新定位論文賣點",
         dict(fontsize=15, bold=True, color=RED)),
        ("④ 真實數據試煉：mock community 顯示模擬排名無法照搬到真實世界，是本週最重要的科學發現",
         dict(fontsize=15, bold=True, color=ORANGE)),
    ], 0.6, 1.4, 12.1, 3.6, gap=Pt(18))
    add_text_box(sl, "下一步：out-of-genome generalization 實驗（P0）+ 補齊作者資訊",
                 0.6, 5.6, 12.1, 0.6, fontsize=14, color=GOLD, bold=True)

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()

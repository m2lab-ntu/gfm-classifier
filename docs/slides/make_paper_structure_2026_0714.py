#!/usr/bin/env python3
"""
Paper content & structure walkthrough deck · 2026-07-14
Manuscript: "How should genomic foundation models be evaluated for
metagenomic taxonomic classification? A controlled benchmark of
pre-training, tokenization, and data scaling" (Briefings in Bioinformatics,
Problem-solving Protocol).

Purpose: introduce the FULL planned content and structure of the manuscript
(title/abstract/key points, all 9 sections, supplementary, endmatter status),
not weekly progress. Built after `git pull` on the paper repo (picked up
Cursor's two-replicate mock-community strengthening + Fig.1 TikZ fixes).
"""

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from matplotlib.patches import Rectangle, FancyBboxPatch, FancyArrowPatch
import numpy as np
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DOCS = Path(__file__).parent
TMP = DOCS / "paper_structure_figs"
TMP.mkdir(exist_ok=True)
OUT = DOCS / "paper_structure_2026_0714.pptx"

NAVY = "#0A2540"
TEAL = "#1E6091"
GREEN = "#2E7D32"
RED = "#C62828"
ORANGE = "#E65100"
GOLD = "#B8860B"
GRAY = "#5A6068"
LIGHT = "#F4F6F8"
WHITE = "#FFFFFF"
PURPLE = "#6A1B9A"

CJK_BOLD = "/work/ymj1123ntu/NTU-Thesis-Writing-Template/fonts/chinese/cwTeX_Hei.ttf"
CJK_REG = "/work/ymj1123ntu/NTU-Thesis-Writing-Template/fonts/chinese/MOE_Song.ttf"
fp_bold = fm.FontProperties(fname=CJK_BOLD)
fp_reg = fm.FontProperties(fname=CJK_REG)
CJK_FONT_NAME = "Microsoft JhengHei"


def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


plt.rcParams.update({
    "axes.spines.top": False, "axes.spines.right": False,
    "figure.facecolor": "white", "savefig.dpi": 200,
    "savefig.bbox": "tight", "savefig.pad_inches": 0.15,
})


# ═══════════════════════════════════════════════════════════════════════
# FIGURES
# ═══════════════════════════════════════════════════════════════════════

def fig_roadmap():
    """9-section flow of the manuscript."""
    fig, ax = plt.subplots(figsize=(13, 5.4))
    ax.set_xlim(0, 13); ax.set_ylim(0, 6); ax.axis("off")

    sections = [
        ("§1", "Introduction", TEAL),
        ("§2", "為何 GFM\n逆風", TEAL),
        ("§3", "Benchmark\n設計", NAVY),
        ("§4", "資料規模\n飽和", ORANGE),
        ("§5", "Pretrain vs\nTokenization", RED),
        ("§6", "Read vs Sample\n+ 真實驗證", PURPLE),
        ("§7", "實務\n建議", GREEN),
        ("§8", "限制與\n未來方向", GRAY),
    ]
    n = len(sections)
    xs = np.linspace(0.9, 12.1, n)
    y = 3.3
    w, h = 1.25, 1.5
    for x, (tag, label, color) in zip(xs, sections):
        box = FancyBboxPatch((x - w/2, y - h/2), w, h, boxstyle="round,pad=0.08",
                              facecolor=color, edgecolor="none", alpha=0.92)
        ax.add_patch(box)
        ax.text(x, y + 0.42, tag, ha="center", fontsize=13, fontweight="bold", color=WHITE)
        ax.text(x, y - 0.15, label, ha="center", va="center", fontsize=10.5, color=WHITE,
                fontproperties=fp_bold)
    for i in range(n - 1):
        ax.annotate("", xy=(xs[i+1] - w/2 - 0.03, y), xytext=(xs[i] + w/2 + 0.03, y),
                    arrowprops=dict(arrowstyle="-|>", color=GRAY, lw=1.6))

    ax.text(6.5, 5.5, "全文架構：破題 → 設計 → 三個因果實驗 → 實務轉譯 → 誠實揭露限制",
            ha="center", fontsize=14, fontweight="bold", color=NAVY, fontproperties=fp_bold)
    ax.text(6.5, 1.0,
            "核心論點貫穿全文：Tokenization（overlapping 13-mer）—— 而非 backbone pre-training —— 決定 150bp 短讀分類的效能上限",
            ha="center", fontsize=11.5, color=NAVY, fontproperties=fp_reg)
    plt.savefig(TMP / "roadmap.png"); plt.close()


def fig_design_pipeline():
    """Sec 3: benchmark design schematic (mirrors manuscript Figure 1 semantics)."""
    fig, ax = plt.subplots(figsize=(13, 4.9))
    ax.set_xlim(0, 13); ax.set_ylim(0, 6); ax.axis("off")

    stages = [
        (1.5, "腸道基因體目錄\n120 屬 / 1,535 種", TEAL),
        (4.7, "洩漏控制模擬 reads\n(ART, 150bp)\nread + genome 雙重切分", NAVY),
        (7.9, "三軸逐一改變\n資料規模 / Pre-training /\nTokenization", ORANGE),
        (11.1, "三層級評估\nRead-level / Sample-level /\nDetection", PURPLE),
    ]
    w, h = 2.3, 2.2
    y = 3.6
    for x, label, color in stages:
        box = FancyBboxPatch((x - w/2, y - h/2), w, h, boxstyle="round,pad=0.1",
                              facecolor=color, edgecolor="none", alpha=0.92)
        ax.add_patch(box)
        ax.text(x, y, label, ha="center", va="center", fontsize=10.5, color=WHITE, fontproperties=fp_bold)
    for i in range(len(stages) - 1):
        x0 = stages[i][0] + w/2 + 0.05
        x1 = stages[i+1][0] - w/2 - 0.05
        ax.annotate("", xy=(x1, y), xytext=(x0, y), arrowprops=dict(arrowstyle="-|>", color=GRAY, lw=1.8))

    box2 = FancyBboxPatch((9.7, 0.5), 2.6, 1.2, boxstyle="round,pad=0.1",
                           facecolor=GOLD, edgecolor="none", alpha=0.92)
    ax.add_patch(box2)
    ax.text(11.0, 1.1, "對照 Kraken2\n(+Bracken)", ha="center", va="center", fontsize=10.5,
            color=WHITE, fontproperties=fp_bold)
    ax.annotate("", xy=(11.0, 1.75), xytext=(11.0, 2.45), arrowprops=dict(arrowstyle="-|>", color=GRAY, lw=1.6))

    ax.text(6.5, 5.6, "Figure 1 概念圖：Benchmark 設計與評估協定（§3）", ha="center",
            fontsize=14, fontweight="bold", color=NAVY, fontproperties=fp_bold)
    ax.text(6.5, 0.3,
            "關鍵細節：258.67M reads 來源池只有 42.64M 唯一序列（6.07× 冗餘）；三個測試池：自然分布 / 較難 leftover / Kraken2 coverage-matched",
            ha="center", fontsize=10, color=GRAY, fontproperties=fp_reg)
    plt.savefig(TMP / "design_pipeline.png"); plt.close()


def fig_scaling_tokenizer():
    fig, ax = plt.subplots(figsize=(12.5, 5.3))
    x6 = [0.5, 5, 50, 250]
    y6 = [55.29, 63.05, 67.07, 67.29]
    x13 = [50, 250]
    y13 = [87.42, 98.7]

    ax.plot(x6, y6, "o-", color=TEAL, lw=2.5, ms=9, label="NT-v2 non-overlapping 6-mer（pretrained）")
    ax.plot(x13, y13, "o-", color=NAVY, lw=2.5, ms=9, label="MT overlapping 13-mer（from scratch）")
    ax.set_xscale("log")
    ax.set_xticks([0.5, 5, 50, 250])
    ax.set_xticklabels(["500K", "5M", "50M", "250M"], fontsize=11)
    ax.set_xlabel("訓練 reads 數量（log scale）", fontsize=12, fontproperties=fp_reg)
    ax.set_ylabel("Genus Top-1 準確度 (%)", fontsize=12, fontproperties=fp_reg)
    ax.set_ylim(45, 105)
    ax.grid(axis="y", alpha=0.25)
    ax.legend(fontsize=11, prop=fp_reg, loc="upper left", frameon=False)

    ax.annotate("50M→250M 只 +0.22pp\n（飽和）", xy=(250, 67.29), xytext=(60, 50),
                fontsize=10.5, color=TEAL, fontweight="bold", fontproperties=fp_bold,
                arrowprops=dict(arrowstyle="->", color=TEAL, lw=1.4))
    ax.annotate("50M→250M +11.3pp\n（持續提升）", xy=(250, 98.7), xytext=(90, 92),
                fontsize=10.5, color=NAVY, fontweight="bold", fontproperties=fp_bold,
                arrowprops=dict(arrowstyle="->", color=NAVY, lw=1.4))
    ax.set_title("§4 資料規模效應：飽和與否取決於 tokenizer", fontsize=14, fontweight="bold",
                 color=NAVY, fontproperties=fp_bold, pad=14)
    plt.savefig(TMP / "scaling_tokenizer.png"); plt.close()


def fig_decomposition():
    fig, axes = plt.subplots(1, 2, figsize=(13, 5.0))

    # Left: bridge chart pretrain -> tokenization
    ax = axes[0]
    labels = ["Random-init\n6-mer", "Pretrained\nNT-v2 6-mer", "13-mer\nscratch"]
    vals = [53.88, 67.07, 87.42]
    colors = [GRAY, TEAL, NAVY]
    bars = ax.bar(labels, vals, color=colors, width=0.55, zorder=3)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width()/2, v + 3.5, f"{v:.2f}", ha="center", fontsize=11,
                fontweight="bold", zorder=5,
                bbox=dict(facecolor="white", edgecolor="none", pad=1, alpha=0.85))
    ax.annotate("", xy=(1, 67.07 - 8), xytext=(0, 53.88 + 8),
                arrowprops=dict(arrowstyle="->", color=GREEN, lw=2))
    ax.text(0.5, 57.5, "+13.19pp\nPre-training", ha="center", fontsize=10, color=GREEN, fontweight="bold",
            fontproperties=fp_bold)
    ax.annotate("", xy=(2, 87.42 - 8), xytext=(1, 67.07 + 8),
                arrowprops=dict(arrowstyle="->", color=RED, lw=2))
    ax.text(1.5, 74, "+20.35pp\nTokenization", ha="center", fontsize=10, color=RED, fontweight="bold",
            fontproperties=fp_bold)
    ax.set_ylim(0, 110); ax.set_ylabel("Genus Top-1 (%)", fontsize=11)
    ax.set_title("固定 50M reads：兩個因子的獨立貢獻", fontsize=12.5, fontweight="bold",
                 color=NAVY, fontproperties=fp_bold)
    ax.grid(axis="y", alpha=0.25)

    # Right: is it lookup? ladder
    ax2 = axes[1]
    labels2 = ["Unique-key\nlookup", "Majority\nvote", "Naive\nBayes", "Neural\n13-mer"]
    vals2 = [0.72, 35.3, 74.9, 87.42]
    colors2 = [RED, ORANGE, GOLD, NAVY]
    bars2 = ax2.barh(labels2, vals2, color=colors2, height=0.55, zorder=3)
    for b, v in zip(bars2, vals2):
        ax2.text(v + 1.5, b.get_y() + b.get_height()/2, f"{v:.1f}%", va="center", fontsize=11, fontweight="bold")
    ax2.set_xlim(0, 105); ax2.set_xlabel("Genus Top-1 (%)", fontsize=11)
    ax2.set_title("「13-mer 98.7% 是不是查表？」— 不是", fontsize=12.5, fontweight="bold",
                  color=NAVY, fontproperties=fp_bold)
    ax2.grid(axis="x", alpha=0.25)
    ax2.invert_yaxis()

    plt.savefig(TMP / "decomposition.png"); plt.close()


def fig_kraken_bracken():
    fig, ax = plt.subplots(figsize=(11, 5.6))
    methods = [
        ("Kraken2 (raw)", 0.823, 93.5, RED),
        ("Kraken2+Bracken", 0.997, 93.5, GREEN),
        ("NT-v2 6-mer", 0.992, 17.5, TEAL),
        ("MT 13-mer", 0.999, 66.5, NAVY),
    ]
    for name, r, sens, color in methods:
        ax.scatter(r, sens, s=800, color=color, alpha=0.85, edgecolor="white", linewidth=2, zorder=3)
        ax.annotate(name, (r, sens), xytext=(0, 22), textcoords="offset points",
                    ha="center", fontsize=11, fontweight="bold", color=color, fontproperties=fp_bold)
    ax.set_xlabel("豐度估計 Pearson r（in-database, coverage-matched）", fontsize=12, fontproperties=fp_reg)
    ax.set_ylabel("偵測敏感度 @95% 特異度 (%)", fontsize=12, fontproperties=fp_reg)
    ax.set_xlim(0.78, 1.03); ax.set_ylim(0, 105)
    ax.grid(alpha=0.25)
    ax.annotate("Bracken 修正後\n兩者兼優（右上角）", xy=(0.997, 93.5), xytext=(0.88, 75),
                fontsize=10.5, color=GREEN, fontweight="bold", fontproperties=fp_bold,
                arrowprops=dict(arrowstyle="->", color=GREEN, lw=1.4))
    ax.set_title("§6：in-database 上「豐度 vs 偵測」的權衡（含 Bracken 修正）", fontsize=13.5,
                 fontweight="bold", color=NAVY, fontproperties=fp_bold, pad=14)
    plt.savefig(TMP / "kraken_bracken.png"); plt.close()


def fig_mock_collapse():
    fig, axes = plt.subplots(1, 2, figsize=(13, 5.0))

    ax = axes[0]
    methods = ["Kraken2+Bracken", "MT 13-mer", "NT-v2 6-mer"]
    sim_r = [0.997, 0.999, 0.992]
    real_r = [0.578, 0.538, 0.412]
    x = np.arange(len(methods)); w = 0.32
    ax.bar(x - w/2, sim_r, width=w, color=TEAL, label="模擬（in-DB）", zorder=3)
    ax.bar(x + w/2, real_r, width=w, color=ORANGE, label="真實 D6331（2 重複均值）", zorder=3)
    for xi, sv, rv in zip(x, sim_r, real_r):
        ax.text(xi - w/2, sv + 0.03, f"{sv:.3f}", ha="center", fontsize=9)
        ax.text(xi + w/2, rv + 0.03, f"{rv:.3f}", ha="center", fontsize=9)
    ax.set_xticks(x); ax.set_xticklabels(methods, fontsize=10.5)
    ax.set_ylabel("Pearson r（genus 豐度）", fontsize=11, fontproperties=fp_reg)
    ax.set_ylim(0, 1.3)
    ax.legend(fontsize=10, prop=fp_reg, loc="upper center", bbox_to_anchor=(0.5, 1.14), ncol=2, frameon=False)
    ax.set_title("模擬 → 真實：全面崩落（n=11 屬，CI 重疊、無明顯排名）", fontsize=12,
                 fontweight="bold", color=NAVY, fontproperties=fp_bold, pad=42)
    ax.grid(axis="y", alpha=0.25)

    ax2 = axes[1]
    ax2.axis("off"); ax2.set_xlim(0, 10); ax2.set_ylim(0, 10)
    ax2.text(5, 9.4, "真實 D6331（2 重複）關鍵數字", fontsize=13, fontweight="bold",
             color=NAVY, ha="center", fontproperties=fp_bold)
    rows = [
        ("樣本", "ZymoBIOMICS D6331，SRR33710519/518，各 3M reads，135bp"),
        ("In-set 屬", "11/21（81.5% reads），其餘含 Escherichia(14%) 不在訓練集內"),
        ("Kraken2 分類率", "真實僅 ~44%（模擬時 ~70%）"),
        ("偵測敏感度", "Bracken 55.6%  ·  NT-v2 83.3%  ·  MT13-mer 77.8%"),
        ("Bray-Curtis", "Kraken2+Bracken 最差 0.434（Veillonella 高估）"),
        ("共同失敗模式", "Roseburia 兩神經模型都嚴重低估；Clostridium 吸收多餘質量"),
    ]
    y0 = 8.2
    for name, val in rows:
        ax2.text(0.2, y0, name, fontsize=11, fontweight="bold", color=TEAL, fontproperties=fp_reg)
        ax2.text(0.2, y0 - 0.5, f"  {val}", fontsize=9.8, color=GRAY, fontproperties=fp_reg)
        y0 -= 1.35

    plt.savefig(TMP / "mock_collapse.png"); plt.close()


# ═══════════════════════════════════════════════════════════════════════
# PPTX HELPERS
# ═══════════════════════════════════════════════════════════════════════

def add_slide(prs, layout=6):
    return prs.slides.add_slide(prs.slide_layouts[layout])


def title_bar(slide, title, subtitle=""):
    W = Inches(13.33)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(NAVY)
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = rgb(WHITE)
    r.font.name = CJK_FONT_NAME
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(13); r2.font.color.rgb = rgb("#90CAF9")
        r2.font.name = CJK_FONT_NAME


def add_image(slide, path, left, top, width, height):
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))


def add_text_box(slide, text, left, top, width, height,
                  fontsize=12, color=NAVY, bold=False, align=PP_ALIGN.LEFT, bg=None, border=None):
    if bg is not None or border is not None:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                      Inches(width), Inches(height))
        if bg is not None:
            shp.fill.solid(); shp.fill.fore_color.rgb = rgb(bg)
        else:
            shp.fill.background()
        if border is not None:
            shp.line.color.rgb = rgb(border); shp.line.width = Pt(1.25)
        else:
            shp.line.fill.background()
        tf = shp.text_frame
    else:
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        r = p.add_run(); r.text = line
        r.font.size = Pt(fontsize); r.font.bold = bold
        r.font.color.rgb = rgb(color)
        r.font.name = CJK_FONT_NAME


def bullets(slide, items, left, top, width, height, fontsize=14, color=NAVY, gap=Pt(6)):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        if isinstance(item, tuple):
            text, kwargs = item
        else:
            text, kwargs = item, {}
        r = p.add_run(); r.text = text
        r.font.size = Pt(kwargs.get("fontsize", fontsize))
        r.font.bold = kwargs.get("bold", False)
        r.font.color.rgb = rgb(kwargs.get("color", color))
        r.font.name = CJK_FONT_NAME
        p.level = kwargs.get("level", 0)


def quad_box(slide, left, top, width, height, title, body, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(LIGHT)
    shp.line.color.rgb = rgb(color); shp.line.width = Pt(1.75)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = rgb(color); r.font.name = CJK_FONT_NAME
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT; p2.space_before = Pt(4)
    r2 = p2.add_run(); r2.text = body
    r2.font.size = Pt(11.5); r2.font.color.rgb = rgb(NAVY); r2.font.name = CJK_FONT_NAME


# ═══════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════

def build():
    print("Generating figures...")
    fig_roadmap()
    fig_design_pipeline()
    fig_scaling_tokenizer()
    fig_decomposition()
    fig_kraken_bracken()
    fig_mock_collapse()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    sl = add_slide(prs)
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = rgb(NAVY); bg.line.fill.background()
    add_text_box(sl, "論文完整內容與架構", 1, 1.0, 11.3, 1.0, fontsize=38, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl,
                 "How should genomic foundation models be evaluated for\nmetagenomic taxonomic classification?",
                 1, 2.1, 11.3, 1.0, fontsize=16, color="#90CAF9", align=PP_ALIGN.CENTER)
    add_text_box(sl, "A controlled benchmark of pre-training, tokenization, and data scaling",
                 1, 3.0, 11.3, 0.5, fontsize=13, color="#BBDEFB", align=PP_ALIGN.CENTER)
    add_text_box(sl, "Briefings in Bioinformatics · Problem-solving Protocol",
                 1, 3.55, 11.3, 0.5, fontsize=13, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(sl,
                 "一句話：150bp 短讀分類的效能上限由 tokenization 決定，不是 backbone pre-training",
                 1, 4.3, 11.3, 0.6, fontsize=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, "楊明儒 · 台大生醫電資所 · 2026-07-14", 1, 6.6, 11.3, 0.5,
                 fontsize=11, color="#78909C", align=PP_ALIGN.CENTER)

    # ── Slide 2: 論文定位 ──
    sl = add_slide(prs)
    title_bar(sl, "論文定位", "投稿現況與文章類型")
    bullets(sl, [
        ("期刊：Briefings in Bioinformatics（BIB）", dict(fontsize=15, bold=True)),
        ("文章類型：Problem-solving Protocol（2,000–5,000 字），而非傳統 original research",
         dict(level=1, fontsize=13, color=GRAY)),
        ("敘事定位：不是「我提出一個新方法」，而是「回答一個方法論問題 —— GFM 該如何被評估用於 "
         "metagenomic taxonomic classification」", dict(fontsize=15, bold=True, color=TEAL)),
        ("第一作者：楊明儒（通訊作者）；單位：台大生醫電資所；第二作者/ORCID/CRediT 待補",
         dict(fontsize=13, color=GRAY)),
        ("目前完整度：main 12 頁 + supplementary，科學內容（摘要～限制）已完整且前後一致；"
         "行政資訊（作者、funding、data DOI）仍有待補齊", dict(fontsize=13, color=ORANGE)),
    ], 0.5, 1.3, 12.3, 5.0)

    # ── Slide 3: Abstract 四象限 ──
    sl = add_slide(prs)
    title_bar(sl, "摘要結構", "Background → Methods → Results → Conclusion")
    quad_box(sl, 0.4, 1.2, 6.1, 2.55, "Background",
             "GFM 在多數基因體任務上表現不錯，但對 150bp 短讀、訊號稀疏的 metagenomic "
             "分類任務效果不明；且不清楚是 pre-training、tokenization 還是資料量在決定表現。", TEAL)
    quad_box(sl, 6.7, 1.2, 6.2, 2.55, "Methods",
             "以腸道基因體目錄（120屬/1,535種）模擬 reads，逐一改變資料規模（500K–250M）、"
             "pre-training（498M pretrained vs random-init）、tokenization（6/12/13-mer），"
             "並在 read/sample/detection 三層級評估，對照 Kraken2。", NAVY)
    quad_box(sl, 0.4, 3.9, 6.1, 2.7, "Results",
             "13-mer 87.4% vs 6-mer 67.1%（+20.35pp）；6-mer 於 50M 後幾乎飽和(+0.22pp)，"
             "13-mer 持續提升到 98.7%；pre-training 貢獻 +13.2pp。豐度估計：67%準確度模型 r=0.99；"
             "Kraken2 raw r=0.82，Bracken 修正到 r=0.997。真實 mock community：全部方法崩落到 r≈0.4–0.6。", GREEN)
    quad_box(sl, 6.7, 3.9, 6.2, 2.7, "Conclusion",
             "Tokenization 是決定性且被低估的因子；未來的 microbial foundation model "
             "應該用 overlapping long-k tokenization 預訓練，才能同時拿到兩個獨立增益。", GOLD)

    # ── Slide 4: 5 Key Points ──
    sl = add_slide(prs)
    title_bar(sl, "五個 Key Points", "main.tex boxed text（規定至多 5 點）")
    bullets(sl, [
        ("① 資料規模效益取決於 tokenizer：50M→250M 對 6-mer 幾乎無增益(+0.22pp)，"
         "對 13-mer 卻有 +11pp", dict(fontsize=14.5)),
        ("② Tokenization、不是 pre-training，決定上限：13-mer scratch(87.4%) 打敗 "
         "498M 參數 pretrained 6-mer(67.1%)", dict(fontsize=14.5)),
        ("③ 固定 tokenization 時 pre-training 仍有幫助(+13.2pp) —— 兩個因子彼此獨立、都重要",
         dict(fontsize=14.5)),
        ("④ Read-level 與 sample-level 效用會分歧：67%準確度模型在模擬數據 r=0.99，"
         "但這排名無法轉移到真實 mock community(r≈0.4–0.6，含 Kraken2+Bracken)",
         dict(fontsize=14.5)),
        ("⑤ 提供實務建議，並主張未來 microbial FM 應以 overlapping long-k tokenization 預訓練",
         dict(fontsize=14.5)),
    ], 0.5, 1.3, 12.3, 5.2, gap=Pt(14))

    # ── Slide 5: 全文架構地圖 ──
    sl = add_slide(prs)
    title_bar(sl, "全文架構地圖", "9 個章節如何串起一個完整論證")
    add_image(sl, TMP / "roadmap.png", 0.2, 1.1, 12.9, 5.3)

    # ── Slide 6: Sec 1-2 ──
    sl = add_slide(prs)
    title_bar(sl, "§1–2　破題與挑戰", "為什麼 metagenomic read classification 對 GFM 是逆風局")
    bullets(sl, [
        ("§1 Introduction：alignment-free 工具（Kraken2/Centrifuge）在 in-database 又快又準，"
         "但遇到新物種就失效；GFM（NT、DNABERT、HyenaDNA、Evo...）在其他基因體任務上很強 —— "
         "能不能用在這裡？", dict(fontsize=13.5)),
        ("先前的神經 read classifier（DeepMicrobes、BERTax、MetaTransformer）證明可行，"
         "但沒有拆解「為什麼」；本文用多個 transformer GFM 在 150bp reads 上系統性拆解 "
         "資料/pre-training/tokenization 三個因子", dict(fontsize=13.5, level=1, color=GRAY)),
        ("§2 為何 GFM 逆風：", dict(fontsize=14, bold=True, color=RED)),
        ("① 短 reads 讓長距離建模優勢（HyenaDNA/Evo 的強項）沒有用武之地",
         dict(level=1, fontsize=13)),
        ("② Tokenization 決定訊號存留：NT-v2 用 non-overlapping 6-mer（~25 tokens/read，"
         "motif 常被切斷）vs Kraken2 用 overlapping 長 k-mer（k≈35）——這是全文核心矛盾",
         dict(level=1, fontsize=13)),
        ("③ 120–1,535 類、長尾分布的標籤空間：選什麼 metric（per-read vs macro）會改變"
         "「誰是贏家」", dict(level=1, fontsize=13)),
    ], 0.5, 1.25, 12.3, 5.5)

    # ── Slide 7: Sec 3 design ──
    sl = add_slide(prs)
    title_bar(sl, "§3　Benchmark 設計與評估協定", "Figure 1：資料 → 三軸 ablation → 三層級評估 → 對照 Kraken2")
    add_image(sl, TMP / "design_pipeline.png", 0.2, 1.1, 12.9, 4.6)
    add_text_box(sl,
                 "三個測試池：自然分布（99,742 reads，主要）／較難 leftover（保守下界）／"
                 "Kraken2 coverage-matched（85,819 reads，公平比較用）",
                 0.2, 5.85, 12.9, 0.9, fontsize=12, color=NAVY)

    # ── Slide 8: Sec 4 scaling ──
    sl = add_slide(prs)
    title_bar(sl, "§4　資料規模與飽和效應", "同樣的資料量，6-mer 飽和、13-mer 繼續漲")
    add_image(sl, TMP / "scaling_tokenizer.png", 0.9, 1.1, 11.5, 5.1)
    add_text_box(sl,
                 "額外證據：6-mer 模型連自己的訓練集都 fit 不到 66% 以上（train≈val）——"
                 "代表性瓶頸，不是優化問題；genus-balance 會讓 per-read 準確度掉 23pp、豐度 r 掉到 0.862",
                 0.9, 6.25, 11.5, 0.9, fontsize=11.5, color=GRAY)

    # ── Slide 9: Sec 5 decomposition ──
    sl = add_slide(prs)
    title_bar(sl, "§5　Pre-training vs Tokenization 分解", "兩個因子互相獨立，且 tokenization 效應更大")
    add_image(sl, TMP / "decomposition.png", 0.2, 1.1, 12.9, 4.75)
    add_text_box(sl,
                 "「是不是查表」的檢驗鏈：exact-match lookup 僅 0.72%（13-mer 在 4^13 空間裡不是屬層級唯一鍵）"
                 "→ 神經網路加值 +13pp 於 naive Bayes 之上 → 結論：學到的是組成表徵，不是記憶",
                 0.2, 5.95, 12.9, 0.9, fontsize=11.5, color=GRAY)

    # ── Slide 10: Sec 6a Kraken/Bracken ──
    sl = add_slide(prs)
    title_bar(sl, "§6a　Read-level vs Sample-level：與 Kraken2 的權衡", "In-database 數據上：豐度 vs 偵測，兩件事")
    add_image(sl, TMP / "kraken_bracken.png", 1.3, 1.1, 10.7, 5.2)

    # ── Slide 11: Sec 6b mock community ──
    sl = add_slide(prs)
    title_bar(sl, "§6b　真實試煉：ZymoBIOMICS Mock Community (D6331)", "全文最關鍵的誠實揭露：模擬排名無法轉移到真實世界")
    add_image(sl, TMP / "mock_collapse.png", 0.2, 1.1, 12.9, 4.9)
    add_text_box(sl,
                 "§6 收尾訊息：neural 模型在 in-database 豐度估計上並不比配置完整的 Kraken2+Bracken 更強；"
                 "真正的價值在 read-level 準確度與 out-of-database 穩健性",
                 0.2, 6.1, 12.9, 0.8, fontsize=12, color=PURPLE, bold=True)

    # ── Slide 12: Sec 7 recommendations ──
    sl = add_slide(prs)
    title_bar(sl, "§7　實務建議", "Box 1 協定 + Box 2 決策指南 + 建模者建議")
    bullets(sl, [
        ("Box 1：8 步驟可重用評估協定（每步都附上 rationale/pitfall）", dict(fontsize=14, bold=True)),
        ("先固定 read 長度 → read+genome 雙重切分 → 算唯一 reads 不只算 reads → 一次只改一個軸 → "
         "三層級分開評估 → 同時用 raw 與修正後的 classical baseline → 統一 RC-TTA 協定 → 用真實數據驗證並公開一切",
         dict(level=1, fontsize=12, color=GRAY)),
        ("Box 2 決策指南（依 use case）：低豐度偵測用 Kraken2；closed-set read 分類選長 k 模型；"
         "只有 6-mer GFM 別指望資料量救得了；真實群落豐度目前沒有穩定贏家",
         dict(fontsize=13.5)),
        ("給建模者的建議：tokenization 是槓桿最大的設計選擇 —— 未來 microbial FM 應以 "
         "overlapping 12–13-mer 預訓練；但點名一個真實工程挑戰：naive 長k詞彙表意味著"
         "數十億參數的 embedding table，需要 factorization/hashing 等技術",
         dict(fontsize=13.5, color=ORANGE)),
    ], 0.5, 1.3, 12.3, 5.3)

    # ── Slide 13: Sec 8 limitations ──
    sl = add_slide(prs)
    title_bar(sl, "§8　限制與未來方向", "主動揭露，而非等 reviewer 抓到")
    bullets(sl, [
        ("P0：Closed-set 評估 —— 98.7%/r≈1 都是同一目錄的 train/test；雖已證明是「學到組成」而非查表，"
         "但 out-of-genome generalization 仍是最重要的待辦，發表前最好先做",
         dict(fontsize=14, bold=True, color=RED)),
        ("模擬 reads 的侷限：ART 模擬可能對 alignment/k-mer 方法有利；真實 mock community 已顯示"
         "closed-set-to-real 的落差，但目前只有 1 個標準品/11 個屬，需要更廣的真實驗證",
         dict(fontsize=13.5)),
        ("Baseline 廣度：只比較 Kraken2（raw + Bracken），建議未來加入 Centrifuge 等",
         dict(fontsize=13.5)),
        ("單一 read 長度：模擬用 150bp、mock 用 135bp，其他長度（75/250bp）未測試；"
         "RC-TTA 只用在 NT-v2，MetaTransformer 少了約 0.7pp（已揭露、不影響定性結論）",
         dict(fontsize=13.5)),
    ], 0.5, 1.3, 12.3, 5.3)

    # ── Slide 14: Supplementary + Endmatter ──
    sl = add_slide(prs)
    title_bar(sl, "Supplementary 內容 + 投稿前待補清單")
    bullets(sl, [
        ("Supplementary（S1–S8）：資料切分 · 模型訓練設定 · 其他 GFM backbone 比較 · "
         "完整 read-level 結果 · train-fit 上限 · 查表分析 · 計算成本 · 真實 mock 細節",
         dict(fontsize=13.5, bold=True)),
        ("✅ 已填：Conflicts of interest、Code availability（GitHub repo）、AI 使用揭露",
         dict(fontsize=13, color=GREEN)),
        ("⚠️ 待補：Funding 金額/grant number、Data DOI（repository 名稱/URL）、"
         "第二作者 + ORCID + CRediT initials（目前是 X.X./Y.Y. 佔位）、Acknowledgments 名單",
         dict(fontsize=13, color=ORANGE)),
        ("⚠️ 待補：Supplementary S2 的 optimizer/scheduler 名稱需再次確認",
         dict(fontsize=13, color=ORANGE)),
    ], 0.5, 1.3, 12.3, 4.8)
    add_text_box(sl, "科學內容已完整一致；行政資訊是目前唯一的「非科學」缺口",
                 0.5, 6.15, 12.3, 0.6, fontsize=13, color=GRAY, bold=True)

    # ── Slide 15: Summary ──
    sl = add_slide(prs)
    title_bar(sl, "總結：一張圖看完整論文")
    bullets(sl, [
        ("核心貢獻：把「data / pre-training / tokenization」三個常被混在一起的因子，"
         "用嚴謹的 one-factor-at-a-time 設計拆開，找出 tokenization 才是決定性變數",
         dict(fontsize=15, bold=True, color=NAVY)),
        ("最誠實的部分：主動加入 Bracken baseline 撤回過度宣稱的豐度優勢，"
         "又主動加入真實 mock community 揭露模擬-真實的落差", dict(fontsize=15, bold=True, color=RED)),
        ("最大待辦：out-of-genome generalization 實驗（P0）+ 作者/funding/DOI 行政資訊",
         dict(fontsize=15, bold=True, color=ORANGE)),
    ], 0.6, 1.5, 12.1, 3.4, gap=Pt(20))
    add_text_box(sl, "下一步：先做 P0 實驗，再補齊行政資訊，準備投稿",
                 0.6, 5.6, 12.1, 0.6, fontsize=14, color=GOLD, bold=True)

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()

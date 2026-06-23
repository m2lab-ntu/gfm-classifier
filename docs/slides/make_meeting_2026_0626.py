#!/usr/bin/env python3
"""
Weekly meeting deck · 2026-06-26
Covers progress 2026-06-19 → 2026-06-26.

Key story (two headlines):
  (A) SELF-CORRECTION of last week's claim. Last week I said the 66.6% v9
      baseline was leakage-inflated. With a clean, train-disjoint common
      validation set (99,742 reads, zero seq_id overlap with EITHER 50M or
      250M training), v9 still scores 66.4% (fwd) / 67.1% (RC-TTA). So it is
      NOT leakage. The 66-vs-42 gap I saw was a TEST-DISTRIBUTION artifact
      (natural per-read vs per-genus-balanced 1000/genus), not data leakage.
  (B) THE REAL FINDING: data SATURATES. With the bug fixed and data correct,
      scaling balanced data 50M → 250M (5x) moves warm-start accuracy only
      67.07% → 67.29% (+0.22 pp). v14 from-scratch 250M is even lower (64.8%).
      This refutes the log-linear "need more data" extrapolation. The ceiling
      is TOKENIZATION (MT overlapping 13-mer 87.42% vs NT-v2 6-mer 67%, +20 pp),
      because the reference is one genome per species — beyond ~50M reads,
      extra reads only redundantly recover the same 1,535 sequences.
  Thesis updated accordingly (figure + abstract + ch01 + ch05). DNABERT-2 50M
  supplementary run stopped (converged 59.3%, stuck re-running epoch 18).
"""

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle
import numpy as np
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DOCS = Path(__file__).parent
TMP  = DOCS / "meeting_0626_figs"
TMP.mkdir(exist_ok=True)
OUT  = DOCS / "meeting_2026_0626.pptx"

NAVY   = "#0A2540"
TEAL   = "#1E6091"
GREEN  = "#2E7D32"
RED    = "#C62828"
ORANGE = "#E65100"
GOLD   = "#B8860B"
GRAY   = "#5A6068"
LIGHT  = "#F4F6F8"
WHITE  = "#FFFFFF"
PURPLE = "#6A1B9A"

def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

plt.rcParams.update({"font.family":"DejaVu Sans","font.size":11,
                     "axes.spines.top":False,"axes.spines.right":False,
                     "figure.facecolor":"white","savefig.dpi":200,
                     "savefig.bbox":"tight","savefig.pad_inches":0.15})


# ═══════════════════════════════════════════════════════════════════════════════
# FIGURES
# ═══════════════════════════════════════════════════════════════════════════════

def fig_timeline():
    fig, ax = plt.subplots(figsize=(13, 4.8))
    ax.set_xlim(0, 13); ax.set_ylim(0, 5.5); ax.axis("off")
    ax.add_patch(Rectangle((0.3, 2.5), 12.4, 0.13, facecolor=NAVY))

    events = [
        (0.9,  "6/19", "Last report:\nv14/v15 running\n+ leakage suspicion", TEAL,   True),
        (3.0,  "6/21", "v14/v15 converged\n64.8% / 67.3%",                   GREEN,  False),
        (5.2,  "6/22", "Clean common val\n(99,742 reads,\nzero overlap)",    TEAL,   True),
        (7.2,  "6/23", "RC-TTA fixed\n(seq-level RC)\n+ DNABERT-2 stopped",  ORANGE, False),
        (9.4,  "6/23", "CORRECTION:\n66% is NOT leakage\n(distribution)",    RED,    True),
        (11.6, "6/26", "TODAY\nthesis updated\n(saturation)",               GOLD,   False),
    ]
    for x, date, label, color, up in events:
        ax.scatter(x, 2.565, s=200, color=color, zorder=4, edgecolor="white", linewidth=2)
        if up:
            ax.text(x, 3.05, date,  ha="center", fontsize=10, fontweight="bold", color=color)
            ax.text(x, 3.45, label, ha="center", fontsize=9,  color=NAVY, va="bottom")
            ax.plot([x, x], [2.62, 3.0], color=color, lw=1, alpha=0.5)
        else:
            ax.text(x, 2.10, date,  ha="center", fontsize=10, fontweight="bold", color=color)
            ax.text(x, 1.95, label, ha="center", fontsize=9,  color=NAVY, va="top")
            ax.plot([x, x], [2.50, 2.15], color=color, lw=1, alpha=0.5)

    ax.text(6.5, 5.2, "Progress Timeline · 6/19 → 6/26 (today)",
            ha="center", fontsize=13, fontweight="bold", color=NAVY)
    plt.savefig(TMP/"timeline.png", bbox_inches="tight"); plt.close()


def fig_clean_results():
    """Clean train-disjoint test (RC-TTA) bars: v9 / v15 / v14, vs old collapse."""
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.9))
    fig.patch.set_facecolor("white")

    # Left: clean RC-TTA bars
    ax = axes[0]
    names  = ["v9\n50M", "v15\n250M warm", "v14\n250M scratch"]
    rc     = [67.1, 67.3, 64.8]
    fwd    = [66.4, 66.6, 64.1]
    cols   = [GREEN, GREEN, TEAL]
    x = np.arange(len(names))
    ax.bar(x-0.18, fwd, 0.36, color=[c+"66" for c in cols], label="forward")
    bars = ax.bar(x+0.18, rc, 0.36, color=cols, label="+ RC-TTA")
    for xi, v in zip(x+0.18, rc):
        ax.text(xi, v+0.4, f"{v:.1f}", ha="center", fontsize=10.5, fontweight="bold", color=NAVY)
    ax.axhline(44.5, color=RED, ls=":", lw=1.6, alpha=0.8)
    ax.text(2.4, 45.6, "last week's wrong-data\ncollapse ~44.5%", color=RED,
            fontsize=8.5, ha="right", va="bottom")
    ax.set_xticks(x); ax.set_xticklabels(names, fontsize=10)
    ax.set_ylim(40, 72); ax.set_ylabel("Genus accuracy (%)")
    ax.set_title("Clean train-disjoint test — bug fix confirmed",
                 fontsize=12, fontweight="bold", color=NAVY)
    ax.legend(loc="lower left", fontsize=9, framealpha=0.9)

    # Right: notes
    ax2 = axes[1]; ax2.axis("off"); ax2.set_xlim(0,1); ax2.set_ylim(0,1)
    ax2.text(0.5, 0.97, "Same Ruler for All Three", ha="center",
             fontsize=13, fontweight="bold", color=NAVY, va="top")
    rows = [
        (GREEN, "Bug fix fully confirmed",
                "Correct CrucialX9 1,535-sp data + single-line FASTA.\nEven from-scratch v14 (64.8%) clears the old 44.5%\nceiling by 20 pp; warm-start v15 sits at 67.3%."),
        (TEAL,  "Clean common validation set",
                "99,742 reads, seq_id zero-overlap with BOTH the 50M\nand the 250M training sets — a fair, leak-free ruler."),
        (GOLD,  "RC-TTA reported for all final numbers",
                "Sequence-level reverse complement, +~0.7 pp each\n(see RC-TTA slide). All thesis numbers use RC-TTA."),
    ]
    y = 0.83
    for color, title, body in rows:
        ax2.add_patch(FancyBboxPatch((0.03, y-0.135), 0.94, 0.16,
                                     boxstyle="round,pad=0.01",
                                     facecolor=color+"18", edgecolor=color, linewidth=1.4))
        ax2.text(0.07, y+0.008, title, fontsize=10.5, color=color, fontweight="bold", va="top")
        ax2.text(0.07, y-0.045, body, fontsize=9.2, color=NAVY, va="top")
        y -= 0.215

    plt.tight_layout()
    plt.savefig(TMP/"clean_results.png"); plt.close()


def fig_leakage_correction():
    """Self-correction: 66% is NOT leakage; it's a distribution effect."""
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.9))
    fig.patch.set_facecolor("white")

    # Left: two evaluation distributions, both clean
    ax = axes[0]
    labels = ["Natural\nper-read", "Per-genus\nbalanced\n(1000/genus)"]
    vals   = [66.4, 42.0]
    cols   = [GREEN, ORANGE]
    bars = ax.bar(labels, vals, color=cols, width=0.55)
    for b, v in zip(bars, vals):
        ax.text(b.get_x()+b.get_width()/2, v+1.0, f"{v:.0f}%", ha="center",
                fontsize=12, fontweight="bold", color=NAVY)
    ax.set_ylim(0, 78); ax.set_ylabel("v9 accuracy on CLEAN test (%)")
    ax.set_title("Same model, two clean metrics", fontsize=12, fontweight="bold", color=NAVY)
    ax.text(0.5, 70, "both leak-free —\ndifference = distribution,\nnot leakage",
            ha="center", fontsize=9.5, color=GRAY)

    # Right: the correction text
    ax2 = axes[1]; ax2.axis("off"); ax2.set_xlim(0,1); ax2.set_ylim(0,1)
    ax2.text(0.5, 0.97, "Correcting Last Week", ha="center",
             fontsize=13, fontweight="bold", color=RED, va="top")
    rows = [
        (RED,   "Last week's claim (wrong)",
                "\"66.6% is leakage-inflated\" — based only on 300\nsampled seq_ids found inside the 50M train set."),
        (GREEN, "This week's test (correct)",
                "On a zero-overlap clean test, v9 still scores 66.4%\n(fwd) / 67.1% (RC-TTA). Removing overlap barely moves\nit → it is NOT leakage."),
        (TEAL,  "The real explanation",
                "66 vs 42 = two test distributions. Per-read favors\ncommon genera (~66%); per-genus-balanced weights rare\nhard genera equally (~42, macro-like). Both are clean."),
    ]
    y = 0.84
    for color, title, body in rows:
        ax2.add_patch(FancyBboxPatch((0.03, y-0.145), 0.94, 0.17,
                                     boxstyle="round,pad=0.01",
                                     facecolor=color+"18", edgecolor=color, linewidth=1.4))
        ax2.text(0.07, y+0.008, title, fontsize=10.3, color=color, fontweight="bold", va="top")
        ax2.text(0.07, y-0.045, body, fontsize=9.0, color=NAVY, va="top")
        y -= 0.225

    plt.tight_layout()
    plt.savefig(TMP/"leakage_correction.png"); plt.close()


def fig_saturation():
    """The headline: data scaling saturates beyond 50M."""
    fig, ax = plt.subplots(figsize=(11.5, 5.2))
    reads = np.array([500_000, 5_000_000, 50_000_000])
    acc   = np.array([55.29, 63.05, 67.07])

    # log-linear fit on the 3 original points (the projection to be refuted)
    coeffs = np.polyfit(np.log10(reads), acc, 1)
    x_fit  = np.logspace(np.log10(4e5), np.log10(6e8), 300)
    y_fit  = np.polyval(coeffs, np.log10(x_fit))
    ax.plot(x_fit, y_fit, color=TEAL, lw=1.6, ls="--", alpha=0.65,
            label="Log-linear fit (3 points)")

    ax.scatter(reads, acc, color=GREEN, s=95, zorder=5, label="This work")
    for x, y, lbl in zip(reads, acc, ["v4 500K", "v8 5M", "v9 50M"]):
        ax.annotate(f"{lbl}\n({y:.1f}%)", xy=(x, y), xytext=(x, y+2.8),
                    fontsize=9, ha="center",
                    arrowprops=dict(arrowstyle="-", color="gray", lw=0.8))

    # 250M empirical point — same series, falls below the projection
    ax.scatter([250_000_000], [67.29], color=GREEN, s=95, zorder=5)
    ax.annotate("v15 250M\n(67.3%)", xy=(2.5e8, 67.29), xytext=(2.5e8, 60.0),
                fontsize=9, ha="center",
                arrowprops=dict(arrowstyle="-", color="gray", lw=0.8))
    proj = np.polyval(coeffs, np.log10(2.5e8))
    ax.scatter([250_000_000], [proj], facecolors="none", edgecolors=RED,
               s=130, lw=1.8, zorder=4)
    ax.annotate(f"projected ~{proj:.0f}%", xy=(2.5e8, proj), xytext=(6.0e8, proj+1.5),
                fontsize=8.5, ha="center", color=RED,
                arrowprops=dict(arrowstyle="->", color=RED, lw=1))
    ax.annotate("", xy=(2.5e8, 67.6), xytext=(2.5e8, proj-0.3),
                arrowprops=dict(arrowstyle="<->", color=RED, lw=1.4, alpha=0.7))
    ax.text(2.9e8, (proj+67.3)/2, "gap = saturation", fontsize=8.5, color=RED, va="center")

    ax.set_xscale("log")
    ax.set_xlabel("Training reads (log scale)")
    ax.set_ylabel("Genus RC-TTA accuracy (%)")
    ax.set_title("Data Saturates Beyond 50M: 5× more data → only +0.22 pp",
                 fontsize=13, fontweight="bold", color=NAVY)
    ax.set_ylim(50, 75); ax.set_xlim(3e5, 1e9)
    ax.legend(loc="lower right", framealpha=0.9, fontsize=9)
    ax.annotate("", xy=(5e7, 67.07), xytext=(5e6, 63.05),
                arrowprops=dict(arrowstyle="<->", color="gray", lw=1.1))
    ax.text(1.5e7, 65.8, "+4.02 pp", fontsize=8.5, ha="center", color="gray")
    ax.annotate("", xy=(2.5e8, 67.29), xytext=(5e7, 67.07),
                arrowprops=dict(arrowstyle="<->", color=GREEN, lw=1.1))
    ax.text(1.1e8, 68.1, "+0.22 pp (5×)", fontsize=9, ha="center", color=GREEN, fontweight="bold")
    plt.savefig(TMP/"saturation.png", bbox_inches="tight"); plt.close()


def fig_tokenization():
    """Why it saturates + the real ceiling is tokenization."""
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.9))
    fig.patch.set_facecolor("white")

    # Left: tokenization bar
    ax = axes[0]
    labels = ["NT-v2\nnon-overlap\n6-mer", "MetaTransformer\noverlap\n13-mer"]
    vals   = [67.07, 87.42]
    cols   = [TEAL, PURPLE]
    bars = ax.bar(labels, vals, color=cols, width=0.55)
    for b, v in zip(bars, vals):
        ax.text(b.get_x()+b.get_width()/2, v+1.0, f"{v:.1f}%", ha="center",
                fontsize=12, fontweight="bold", color=NAVY)
    ax.annotate("", xy=(1, 87.42), xytext=(1, 67.07),
                arrowprops=dict(arrowstyle="<->", color=RED, lw=2))
    ax.text(1.18, 77, "+20.4 pp", fontsize=12, color=RED, fontweight="bold", va="center")
    ax.set_ylim(0, 100); ax.set_ylabel("Genus Top-1 (%)")
    ax.set_title("Same 50M, from scratch — tokenization decides",
                 fontsize=12, fontweight="bold", color=NAVY)

    # Right: explanation
    ax2 = axes[1]; ax2.axis("off"); ax2.set_xlim(0,1); ax2.set_ylim(0,1)
    ax2.text(0.5, 0.97, "Why 50M Is Enough — and What Isn't", ha="center",
             fontsize=12.5, fontweight="bold", color=NAVY, va="top")
    rows = [
        (GREEN, "Why data saturates",
                "Reference = 1 genome per species. Beyond ~50M reads,\nextra reads only redundantly re-cover the same 1,535\nsequences — no new discriminative signal."),
        (PURPLE,"The real ceiling is tokenization",
                "MT overlapping 13-mer from scratch = 87.4%, +20 pp over\nNT-v2's 67%. NT-v2's non-overlapping 6-mer loses\ninformation at tokenize time — unrecoverable by more data."),
        (TEAL,  "Both scaling paths exhausted",
                "More data → tops out at 50M. To break through, change\ntokenization, not volume or backbone capacity."),
    ]
    y = 0.84
    for color, title, body in rows:
        ax2.add_patch(FancyBboxPatch((0.03, y-0.145), 0.94, 0.17,
                                     boxstyle="round,pad=0.01",
                                     facecolor=color+"18", edgecolor=color, linewidth=1.4))
        ax2.text(0.07, y+0.008, title, fontsize=10.3, color=color, fontweight="bold", va="top")
        ax2.text(0.07, y-0.045, body, fontsize=9.0, color=NAVY, va="top")
        y -= 0.225

    plt.tight_layout()
    plt.savefig(TMP/"tokenization.png"); plt.close()


def fig_rctta():
    """RC-TTA done right: +~0.7pp across models."""
    fig, ax = plt.subplots(figsize=(11.5, 4.6))
    names = ["v9 50M", "v15 250M warm", "v14 250M scratch"]
    fwd   = [66.4, 66.6, 64.1]
    rc    = [67.1, 67.3, 64.8]
    x = np.arange(len(names))
    ax.bar(x-0.18, fwd, 0.36, color=GRAY+"99", label="forward")
    ax.bar(x+0.18, rc, 0.36, color=GREEN, label="+ RC-TTA")
    for xi, a, b in zip(x, fwd, rc):
        ax.annotate(f"+{b-a:.1f}", xy=(xi+0.18, b+0.25), ha="center",
                    fontsize=10, color=GREEN, fontweight="bold")
    ax.set_xticks(x); ax.set_xticklabels(names, fontsize=10)
    ax.set_ylim(60, 70); ax.set_ylabel("Genus accuracy (%)")
    ax.set_title("RC-TTA done right: sequence-level reverse complement, +~0.7 pp each",
                 fontsize=12, fontweight="bold", color=NAVY)
    ax.legend(loc="upper right", fontsize=9.5, framealpha=0.9)
    ax.text(0.0, 60.4, "Was token-flip (wrong); now true RC of the DNA sequence, re-tokenized, softmax averaged.",
            fontsize=8.8, color=GRAY)
    plt.savefig(TMP/"rctta.png", bbox_inches="tight"); plt.close()


def fig_thesis_updated():
    fig, ax = plt.subplots(figsize=(13, 5.0))
    ax.axis("off"); ax.set_xlim(0,1); ax.set_ylim(0,1)
    ax.text(0.5, 0.97, "Thesis Updated — 250M Saturation Propagated", ha="center",
            fontsize=14, fontweight="bold", color=NAVY, va="top")

    items = [
        (TEAL,  "Figure (data scaling)",
                "Added 4th green point (250M, 67.3%) falling below the fit's projected ~71% — the curve flattens off its own projection. 250M held out of the fit."),
        (GREEN, "Abstract (zh + en)",
                "After the 50M result: a further 5× to 250M yields only +0.22 pp → saturates; ceiling set by tokenization, not data volume."),
        (ORANGE,"Ch.1 contributions + Ch.5 conclusion",
                "Same +0.22 pp / saturation note; residual gap to MetaTransformer attributed to tokenization, not data."),
        (GRAY,  "Housekeeping",
                "Unified 250M label color; removed a stray '\\' in two plot titles. All committed + pushed; CI builds the PDF."),
    ]
    y = 0.83
    for color, title, body in items:
        ax.add_patch(FancyBboxPatch((0.03, y-0.135), 0.94, 0.155,
                                    boxstyle="round,pad=0.01",
                                    facecolor=color+"15", edgecolor=color, linewidth=1.3))
        ax.text(0.06, y+0.005, title, fontsize=11, color=color, fontweight="bold", va="top")
        ax.text(0.06, y-0.05, body, fontsize=9.3, color=NAVY, va="top")
        y -= 0.195

    plt.tight_layout()
    plt.savefig(TMP/"thesis_updated.png"); plt.close()


def fig_next_steps():
    fig, ax = plt.subplots(figsize=(13, 5.0))
    ax.axis("off"); ax.set_xlim(0,1); ax.set_ylim(0,1)
    ax.text(0.5, 0.97, "Next Steps · Direction Set by the Saturation Result", ha="center",
            fontsize=14, fontweight="bold", color=NAVY, va="top")

    steps = [
        ("Primary: tokenization, not scaling", PURPLE, [
            ("Verify MT 13-mer", "Re-test the 87% on the clean train-disjoint set; add species-level numbers"),
            ("Overlapping long-k", "Best place to invest — beats stacking more data, which now provably saturates"),
        ]),
        ("Thesis story can close", GREEN, [
            ("Three independent results", "Data saturates @50M ✓ · pre-training +13 pp (fixed tok.) ✓ · tokenization +20 pp ✓"),
            ("Honest baseline", "66.4/67.1% on clean test; leakage claim retracted; distribution effect documented"),
        ]),
        ("Loose ends", GRAY, [
            ("DNABERT-2 50M", "Stopped at 59.3% (resume bug looped epoch 18); 59% < NT-v2 6-mer supports tokenization story"),
            ("Optional", "Fix resume off-by-one only if epochs 19–30 are needed — convergence says low value"),
        ]),
    ]
    y = 0.86
    for section, color, its in steps:
        ax.text(0.02, y, section, fontsize=11.5, color=color, fontweight="bold")
        y -= 0.05
        for label, detail in its:
            ax.add_patch(FancyBboxPatch((0.02, y-0.058), 0.96, 0.064,
                                        boxstyle="round,pad=0.01",
                                        facecolor=color+"15", edgecolor=color+"55", linewidth=1))
            ax.text(0.05, y-0.012, label, fontsize=9.8, color=color, fontweight="bold")
            ax.text(0.30, y-0.012, detail, fontsize=9.0, color=NAVY)
            y -= 0.078
        y -= 0.022

    plt.tight_layout()
    plt.savefig(TMP/"next_steps.png"); plt.close()


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def add_slide(prs, layout=6):
    return prs.slides.add_slide(prs.slide_layouts[layout])

def title_bar(slide, title, subtitle=""):
    W = Inches(13.33)
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(NAVY)
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = rgb(WHITE)
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(14); r2.font.color.rgb = rgb("#90CAF9")

def add_image(slide, path, left, top, width, height):
    slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                             Inches(width), Inches(height))

def add_text_box(slide, text, left, top, width, height,
                 fontsize=12, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fontsize); r.font.bold = bold; r.font.color.rgb = rgb(color)


def build():
    print("Generating figures...")
    fig_timeline()
    fig_clean_results()
    fig_leakage_correction()
    fig_saturation()
    fig_tokenization()
    fig_rctta()
    fig_thesis_updated()
    fig_next_steps()

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    sl = add_slide(prs)
    bg = sl.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = rgb(NAVY); bg.line.fill.background()
    add_text_box(sl, "Weekly Progress Update", 1, 1.4, 11, 1,
                 fontsize=38, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, "Bug fixed → data saturates at 50M · the ceiling is tokenization",
                 1, 2.6, 11, 0.6, fontsize=16, color="#90CAF9", align=PP_ALIGN.CENTER)
    add_text_box(sl, "June 26, 2026", 1, 3.3, 11, 0.5,
                 fontsize=14, color="#BBDEFB", align=PP_ALIGN.CENTER)
    add_text_box(sl,
                 "50M → 250M = +0.22 pp  ·  66% is NOT leakage (distribution)  ·  MT 13-mer 87% > NT-v2 67%",
                 1, 4.1, 11, 0.6, fontsize=13, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(sl, "YMJ · M2 Lab · NTU CSIE", 1, 6.6, 11, 0.5,
                 fontsize=11, color="#78909C", align=PP_ALIGN.CENTER)

    # ── Slide 2: Timeline ──
    sl = add_slide(prs)
    title_bar(sl, "Progress Timeline · 6/19 → 6/26")
    add_image(sl, TMP/"timeline.png", 0.2, 1.1, 12.9, 4.6)
    rows = [
        ("6/21",  "v14/v15 converged on correct data — 64.8% / 67.3% (RC-TTA)"),
        ("6/22",  "Built clean common val: 99,742 reads, zero seq_id overlap with 50M AND 250M training"),
        ("6/23",  "Fixed RC-TTA (true sequence-level RC); stopped stuck DNABERT-2 run"),
        ("6/23–26","Retracted the leakage claim (it's distribution); updated thesis with the saturation result"),
    ]
    y = 5.9
    for date, text in rows:
        add_text_box(sl, f"  {date}   {text}", 0.3, y, 12.7, 0.35, fontsize=11, color=NAVY)
        y += 0.37

    # ── Slide 3: Clean results ──
    sl = add_slide(prs)
    title_bar(sl, "Fix Confirmed — Clean Numbers",
              "Train-disjoint test, RC-TTA: v9 67.1% · v15 67.3% · v14 64.8% (vs last week's 44.5% collapse)")
    add_image(sl, TMP/"clean_results.png", 0.2, 1.2, 12.9, 5.2)

    # ── Slide 4: Leakage correction ──
    sl = add_slide(prs)
    title_bar(sl, "Correction: 66% is NOT Leakage",
              "Retracting last week's claim — the 66-vs-42 gap is a test-distribution artifact, not data leakage")
    add_image(sl, TMP/"leakage_correction.png", 0.2, 1.2, 12.9, 5.2)

    # ── Slide 5: Saturation (headline) ──
    sl = add_slide(prs)
    title_bar(sl, "The Real Finding: Data Saturates",
              "5× more balanced data (50M → 250M) moves accuracy only +0.22 pp — refutes the log-linear projection")
    add_image(sl, TMP/"saturation.png", 1.0, 1.3, 11.3, 5.1)

    # ── Slide 6: Tokenization ceiling ──
    sl = add_slide(prs)
    title_bar(sl, "Why It Saturates → Tokenization is the Ceiling",
              "1 genome/species → 50M is enough · MT overlap 13-mer 87% vs NT-v2 6-mer 67% (+20 pp)")
    add_image(sl, TMP/"tokenization.png", 0.2, 1.2, 12.9, 5.2)

    # ── Slide 7: RC-TTA ──
    sl = add_slide(prs)
    title_bar(sl, "RC-TTA Done Right", "Sequence-level reverse complement (was token-flip) — consistent +~0.7 pp, zero training cost")
    add_image(sl, TMP/"rctta.png", 1.0, 1.4, 11.3, 4.6)

    # ── Slide 8: Thesis updated ──
    sl = add_slide(prs)
    title_bar(sl, "Thesis Updated")
    add_image(sl, TMP/"thesis_updated.png", 0.2, 1.2, 12.9, 5.2)

    # ── Slide 9: Next steps ──
    sl = add_slide(prs)
    title_bar(sl, "Next Steps")
    add_image(sl, TMP/"next_steps.png", 0.2, 1.2, 12.9, 5.2)

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()

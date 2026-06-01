#!/usr/bin/env python3
"""Summary of Contributions slides — standalone 5-slide deck for weekly meeting."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

TEMPLATE = Path("/work/ymj1123ntu/token_level_gfm_classifier/docs/0512.pptx")
OUT      = Path("/work/ymj1123ntu/token_level_gfm_classifier/docs/contribution_summary.pptx")

BLUE   = RGBColor(0x00, 0x70, 0xC0)
BLACK  = RGBColor(0x00, 0x00, 0x00)
GRAY   = RGBColor(0x55, 0x55, 0x55)
GREEN  = RGBColor(0x00, 0x70, 0x00)
RED    = RGBColor(0xC0, 0x00, 0x00)
ORANGE = RGBColor(0xD0, 0x60, 0x00)
DKBLUE = RGBColor(0x00, 0x40, 0x80)

prs = Presentation(TEMPLATE)
xml_slides = prs.slides._sldIdLst
for sld in list(xml_slides):
    xml_slides.remove(sld)

LAY_TITLE   = prs.slide_layouts[0]
LAY_CONTENT = prs.slide_layouts[1]
LAY_SECTION = prs.slide_layouts[19]

slide_w = prs.slide_width
slide_h = prs.slide_height


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(LAY_TITLE)
    slide.placeholders[0].text = title
    for para in slide.placeholders[0].text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.bold = True
    sub_ph = slide.placeholders[1]
    sub_ph.text = subtitle
    for para in sub_ph.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(18)


def add_content(title, bullets):
    slide = prs.slides.add_slide(LAY_CONTENT)
    title_ph = slide.placeholders[0]
    title_ph.text = title
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(24)
    body_ph = slide.placeholders[1]
    tf = body_ph.text_frame
    tf.clear()
    first = True
    for (level, text, bold, color) in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.size = Pt(15) if level == 0 else Pt(13)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: Title
# ─────────────────────────────────────────────────────────────────────────────
add_title_slide(
    "Summary of Contributions",
    "Genomic Foundation Models for Metagenomic Taxonomic Classification\n楊明儒  ·  資工所碩二  ·  2026"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: Research Setup
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "研究設定 — 任務、框架、規模",
    [
        (0, "任務", True, BLUE),
        (1, "Metagenomic DNA 短讀段 (150 bp) 分類到 genus（120 類）/ species（1,535 類）", False, BLACK),
        (1, "資料集：2,505 HGR-UMGS 基因組 × ART 模擬 Illumina 讀段，共 258M reads", False, BLACK),
        (0, "Framework：NT-v2 + LoRA（token-level pipeline）", True, BLUE),
        (1, "NT-v2 backbone（498M params）+ AttentionPool head，LoRA on Q/K/V（1.1% trainable）", False, BLACK),
        (1, "6-mer tokenization（繼承自 NT-v2 pre-training）", False, BLACK),
        (0, "比較基準：MetaTransformer（from-scratch, 5M params）", True, BLUE),
        (1, "可變 k-mer tokenization（6/12/13-mer × overlapping/non-overlapping）", False, BLACK),
        (1, "同一 50M 訓練資料 → 同等訓練條件下的 controlled ablation", False, BLACK),
        (0, "核心問題", True, DKBLUE),
        (1, "Pre-training、data volume、tokenization 三者對分類準確率的相對貢獻是什麼？", True, DKBLUE),
        (1, "Read-level accuracy 如何轉換成 sample-level 的實用性？", True, DKBLUE),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3: Finding 1 — Tokenization > Pre-training
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "Finding 1 — Tokenization 是跨架構的主導因素",
    [
        (0, "核心對比（genus Top-1 accuracy，50M reads）", True, BLUE),
        (1, "NT-v2 + LoRA  (pre-trained, 6-mer non-overlapping)        67.07%  RC TTA", False, BLACK),
        (1, "MT 13-mer      (scratch,      13-mer overlapping, stride=1)  87.42%         ← +20.35 pp", True, GREEN),
        (1, "⟹  Tokenization 設計超越 pre-training 帶來的優勢", True, GREEN),
        (0, "Tokenization 效應細分", True, BLUE),
        (1, "Non-overlapping vs overlapping（固定 k=13）：+59.1 pp（28.3% → 87.4%）", True, RED),
        (1, "k 增長（fixed stride=1）：6-mer 48.9% < 12-mer 78.8% < 13-mer 87.4%", True, ORANGE),
        (1, "兩個效應疊加：overlapping 是關鍵，長 k 是放大器", False, BLACK),
        (0, "Pre-training 的貢獻（固定 6-mer tokenization）", True, BLUE),
        (1, "NT-v2 pre-trained vs random-init（同架構、同 6-mer、50M reads）：+13.19 pp", False, BLACK),
        (1, "Pre-trained GFM 比較：NT-v2 > DNABERT (+1.27 pp) > DNABERT-2 (+4.17 pp)  @5M", False, BLACK),
        (1, "⟹  Pre-training 有效，但受限於繼承的 6-mer tokenization constraint", True, ORANGE),
        (0, "結論：NT-v2 的主要瓶頸不是 backbone 容量或 LoRA，而是 6-mer tokenization", True, RED),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4: Finding 2 — Within GFM: data >> tricks
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "Finding 2 — Within GFM：Data Volume >> 所有技巧",
    [
        (0, "Data scaling（genus RC TTA accuracy）", True, BLUE),
        (1, "500K → 5M：+7.76 pp  |  5M → 50M：+4.02 pp  ← 每 10× 資料遞減", True, GREEN),
        (1, "任何單一技巧的最大效益：±0.5 pp", True, RED),
        (0, "各項 ablation 結果（genus，5M reads）", True, BLUE),
        (1, "Head type（attention pool vs mean pool）：+0.3 pp", False, GRAY),
        (1, "Logit Adjustment（針對 class imbalance）：±0.2 pp", False, GRAY),
        (1, "RC Consistency loss（訓練時正向/互補一致）：+0.1 pp", False, GRAY),
        (1, "Class balance（balanced subsampling）：macro F1 +12.27 pp，F1=0 class 9→2", True, ORANGE),
        (0, "RC TTA（test-time，免訓練）", True, BLUE),
        (1, "Pre-trained backbone：+0.78 ~ +1.54 pp  |  Random-init：+0.08 pp", False, BLACK),
        (1, "效果來自修正 pre-training 造成的 strand bias，而非降低 stochastic noise", False, BLACK),
        (0, "結論：50M reads 之後繼續 scale data 仍有空間（預測 258M → 82–90%）", True, DKBLUE),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5: Finding 3 — Sample-level utility
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "Finding 3 — Tokenization 決定 Sample-Level 實用性",
    [
        (0, "Genus level（120 類，5M pool，50K reads/sample）", True, BLUE),
        (1, "三個模型 Pearson r 全部 ≥ 0.984，genus 聚合具容錯性", False, BLACK),
        (1, "NT-Genus（67%）r = 0.993，BC = 0.094  ←  接近 MT 13-mer（87%）r = 0.999", True, GREEN),
        (0, "Species level（1,535 類，100K test set，1K reads/sample）", True, BLUE),
        (1, "MT 13-mer flat  49.7% → Pearson r = 0.466，ROC AUC = 0.966  ✓ 實用", True, GREEN),
        (1, "MT 13-mer hier. 50.9% → Pearson r = 0.478，ROC AUC = 0.967  ✓ 最佳", True, GREEN),
        (1, "NT-Species      17.8% → Pearson r = 0.135  ← noise floor 限制（1,535 類，SNR≈1.2）", True, RED),
        (1, "MT 6-mer        6–9%  → Pearson r = 0.034–0.065  ← 無法使用", True, RED),
        (0, "Noise floor 機制（species level）", True, BLUE),
        (1, "1,535 species，每個 absent species 平均收到 ~0.54 spurious reads / 1K sample", False, BLACK),
        (1, "每個 present species 平均 signal ~0.65 reads → SNR ≈ 1.2 → r 被壓到 0.14", False, BLACK),
        (1, "需要 ~40–50% read accuracy 才能讓 signal 超過 noise floor", True, ORANGE),
        (0, "結論：6-mer tokenization（NT-v2 上限 ~18%）無法突破 species 實用門檻", True, RED),
        (1, "Overlapping 13-mer（MT 13-mer 49.7%）剛好跨過，pre-training 無法補償 tokenization 缺口", True, RED),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6: Finding 4 — Hierarchical classification
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "Finding 4 — Hierarchical Classification：驗證 + 瓶頸量化",
    [
        (0, "Logit masking（MT 13-mer）：小但一致的提升", True, BLUE),
        (1, "Flat → Hier.：Pearson r +0.012，BC −0.009，ROC AUC +0.0015，read acc +1.12 pp", True, GREEN),
        (1, "四個指標全部改善——genus-guided masking 有效排除跨 genus 的混淆", True, GREEN),
        (1, "NT-Species 同樣應用 logit masking（NT-Genus 66.1% router）：大多數指標反而下降", True, RED),
        (1, "Read acc 17.8%→15.8%（−2.0 pp）；Pearson r 0.135→0.106（−0.029）；BC 0.589→0.608（worse）", True, RED),
        (1, "⟹  Masking 效果取決於 router 品質：87.4% 有效；66.1% 有害；48.9% 有害", True, ORANGE),
        (0, "Per-genus 50M LoRA classifiers（100K test set，oracle routing）", True, BLUE),
        (1, "Oracle routing：per-genus 27.8% > flat sp_v4 27.4%  ← 驗證 per-genus specialisation 有效", True, GREEN),
        (1, "Predicted routing（66.1% genus model）：14.7% < flat 15.9%  ← routing error 蓋過 specialisation 優勢", True, RED),
        (1, "Per-genus classifier 不認識自己 genus 以外的 read → out-of-distribution problem", False, BLACK),
        (0, "Subgenus routing 嘗試（NT-v2 6-mer embedding + K-means）", True, BLUE),
        (1, "Silhouette score ≈ 0.03（Collinsella 為例）——NT-v2 6-mer 無法在 embedding 空間分離同 genus species", True, RED),
        (1, "Hard/soft subgenus routing 均降低 oracle Top-1：27.8% → 25.6–26.1%", True, RED),
        (1, "⟹  Subgenus routing 需要更好的 representation，不是更多的 routing tier", False, BLACK),
        (0, "核心瓶頸：genus router 準確率（66–67%）決定了 hierarchical pipeline 的上限", True, ORANGE),
    ]
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7: Take-home + Future
# ─────────────────────────────────────────────────────────────────────────────
add_content(
    "Take-home Messages & Future Direction",
    [
        (0, "三個跨架構可推廣的結論", True, BLUE),
        (1, "① Pre-training 有效，但受限於繼承的 tokenization 方案（+13 pp，但 tokenization 差距 +20 pp）", True, BLACK),
        (1, "② Overlapping long-k tokenization 是 short-read 分類的關鍵設計選擇，非架構深度", True, BLACK),
        (1, "③ Tokenization 決定的不只是 read accuracy，而是 sample-level 方法有無實用性的門檻", True, BLACK),
        (1, "④ Hierarchical routing 有效但受限於 router 準確率；per-genus specialisation 在 oracle 下成立", True, BLACK),
        (0, "未完成的驗證", True, ORANGE),
        (1, "Taiwana2：MT 6-mer per-genus classifiers（81 個）→ 確認 per-genus decomposition 上界", False, BLACK),
        (1, "NT-Species hierarchical（樣本級）：已完成——Pearson r 0.135→0.106（more evidence masking requires >80% router）", True, ORANGE),
        (0, "最有潛力的下一步", True, GREEN),
        (1, "Pre-train GFM with overlapping 13-mer tokenization on GTDB（>400K microbial genomes）", True, GREEN),
        (1, "同時獲得 13-mer tokenization 的區分力 + pre-training 的表示能力 → 預期突破兩個瓶頸", True, GREEN),
        (0, "本研究的定位", True, DKBLUE),
        (1, "提供 controlled ablation evidence：告訴社群「哪個設計因素最重要」", True, DKBLUE),
        (1, "為 GFM + overlapping tokenization 的下一代模型提供明確的設計方向", True, DKBLUE),
    ]
)

prs.save(OUT)
print(f"Saved: {OUT}")

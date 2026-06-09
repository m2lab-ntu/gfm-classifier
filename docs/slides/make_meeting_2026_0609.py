#!/usr/bin/env python3
"""
Weekly meeting deck · 2026-06-09
Covers progress 2026-06-02 → 2026-06-09.
Key story: Nano4 onboarded + DNABERT-2/MT benchmark scripts ready → ready to fire.
"""

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle, FancyArrowPatch
import numpy as np
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DOCS = Path("/work/ymj1123ntu/gfm-classifier/docs/slides")
TMP  = DOCS / "meeting_0609_figs"
TMP.mkdir(exist_ok=True)
OUT  = DOCS / "meeting_2026_0609.pptx"

NAVY   = "#0A2540"
TEAL   = "#1E6091"
GREEN  = "#2E7D32"
RED    = "#C62828"
ORANGE = "#E65100"
GOLD   = "#B8860B"
GRAY   = "#5A6068"
LIGHT  = "#F4F6F8"
WHITE  = "#FFFFFF"
PURPLE = "#6A1B9A"

def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

plt.rcParams.update({"font.family":"DejaVu Sans","font.size":11,
                     "axes.spines.top":False,"axes.spines.right":False,
                     "figure.facecolor":"white","savefig.dpi":200,
                     "savefig.bbox":"tight","savefig.pad_inches":0.15})


# ═══════════════════════════════════════════════════════════════════════════════
# FIGURES
# ═══════════════════════════════════════════════════════════════════════════════

def fig_timeline():
    # Alternating up/down labels to prevent horizontal crowding
    fig, ax = plt.subplots(figsize=(13, 5.0))
    ax.set_xlim(0, 13); ax.set_ylim(0, 5.5); ax.axis("off")
    ax.add_patch(Rectangle((0.3, 2.5), 12.4, 0.13, facecolor=NAVY))

    # (x, date, label, color, up=True means label above the line)
    events = [
        (0.8,  "6/2",  "Advisor meeting\n(deck 0602)", TEAL,   True),
        (2.4,  "6/2",  "Nano4 env\nsetup started",     TEAL,   False),
        (4.0,  "6/3",  "Sanity check\nPASSED ✓",       GREEN,  True),
        (5.6,  "6/3-5","peft / vocab /\ndevice fixes",  ORANGE, False),
        (7.2,  "6/5",  "DNABERT-2\nscript ready",      GREEN,  True),
        (8.8,  "6/6",  "MT benchmark\nscript ready",   GREEN,  False),
        (10.4, "6/7",  "data_loader +\nresource_monitor", TEAL, True),
        (12.0, "6/9",  "TODAY\nFire P0", GOLD, False),
    ]
    for x, date, label, color, up in events:
        ax.scatter(x, 2.565, s=200, color=color, zorder=4, edgecolor="white", linewidth=2)
        if up:
            ax.text(x, 3.05, date,  ha="center", fontsize=10, fontweight="bold", color=color)
            ax.text(x, 3.45, label, ha="center", fontsize=9,  color=NAVY, va="bottom")
            ax.plot([x, x], [2.62, 3.0], color=color, lw=1, alpha=0.5)
        else:
            ax.text(x, 2.10, date,  ha="center", fontsize=10, fontweight="bold", color=color)
            ax.text(x, 1.95, label, ha="center", fontsize=9,  color=NAVY, va="top")
            ax.plot([x, x], [2.50, 2.15], color=color, lw=1, alpha=0.5)

    ax.text(6.5, 5.2, "Progress Timeline · 6/2 advisor meeting → 6/9 (today)",
            ha="center", fontsize=13, fontweight="bold", color=NAVY)
    plt.savefig(TMP/"timeline.png", bbox_inches="tight"); plt.close()


def fig_nano4_status():
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.5))
    fig.patch.set_facecolor("white")

    # Left: checklist
    ax = axes[0]; ax.axis("off")
    ax.set_xlim(0,1); ax.set_ylim(0,1)
    ax.text(0.5, 0.97, "Nano4 Migration Checklist", ha="center",
            fontsize=13, fontweight="bold", color=NAVY, va="top")

    items = [
        ("✓", "SSH + conda gfm env (Python 3.11)", GREEN),
        ("✓", "git clone gfm-classifier", GREEN),
        ("✓", "Training data (reads_50M, labels_50M)", GREEN),
        ("✓", "Test data (reads_100K, in_db_mask)", GREEN),
        ("✓", "NT-v2 sp_v4 checkpoint (1.9 GB)", GREEN),
        ("✓", "Sanity check PASSED (NT-Sp = 17.83%)", GREEN),
        ("⟳", "DNABERT-2 last.pt from Nano5 (473 MB)", ORANGE),
        ("⟳", "MT models from Taiwana-2 (benchmark)", ORANGE),
        ("⟳", "HMP mock community FASTQ", ORANGE),
    ]
    y = 0.86
    for mark, text, color in items:
        ax.text(0.05, y, mark, fontsize=13, color=color, fontweight="bold")
        ax.text(0.16, y, text, fontsize=10.5, color=NAVY, va="center")
        y -= 0.09

    # Right: hardware spec
    ax2 = axes[1]; ax2.axis("off")
    ax2.set_xlim(0,1); ax2.set_ylim(0,1)
    ax2.text(0.5, 0.97, "Nano4 Hardware (Free Trial)", ha="center",
             fontsize=13, fontweight="bold", color=NAVY, va="top")

    specs = [
        ("GPU",       "NVIDIA H200 · 143,771 MiB HBM3e"),
        ("CUDA",      "12.6 / 13.0"),
        ("Partition", "dev: 1 hr · 1 GPU"),
        ("",          "normal: 12 hr · min 64 GPU"),
        ("PyTorch",   "2.5.1+cu121"),
        ("transformers","4.46.3 + peft 0.19.1"),
        ("Period",    "Free until ~end of June"),
    ]
    y = 0.86
    for k, v in specs:
        ax2.text(0.05, y, k, fontsize=10.5, color=TEAL, fontweight="bold", va="center")
        ax2.text(0.30, y, v, fontsize=10.5, color=NAVY, va="center")
        y -= 0.09

    # Sanity check result box
    ax2.add_patch(FancyBboxPatch((0.02,0.06), 0.95, 0.22,
                                  boxstyle="round,pad=0.01",
                                  facecolor="#E8F5E9", edgecolor=GREEN, linewidth=2))
    ax2.text(0.5, 0.22, "Sanity Check · Job 71907 · 2026-06-03", ha="center",
             fontsize=10, fontweight="bold", color=GREEN)
    ax2.text(0.5, 0.15, "NT-Species sp_v4 → 17.83%  vs  TWCC 17.83%  =  MATCH ✓",
             ha="center", fontsize=10, color=NAVY)
    ax2.text(0.5, 0.09, "Elapsed: 1 min 32 s  |  batch_size=1024  |  100K reads",
             ha="center", fontsize=9, color=GRAY)

    plt.tight_layout()
    plt.savefig(TMP/"nano4_status.png"); plt.close()


def fig_dnabert2_plan():
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.5))
    fig.patch.set_facecolor("white")

    # Left: training progress bar
    ax = axes[0]; ax.axis("off")
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.text(0.5, 0.97, "DNABERT-2 Genus 50M — Resume Status", ha="center",
            fontsize=13, fontweight="bold", color=NAVY, va="top")

    # Epoch bar
    done = 17/30
    ax.add_patch(FancyBboxPatch((0.05,0.72), 0.90*done, 0.10,
                                 boxstyle="round,pad=0.005", facecolor=GREEN, edgecolor="none"))
    ax.add_patch(FancyBboxPatch((0.05+0.90*done, 0.72), 0.90*(1-done), 0.10,
                                 boxstyle="round,pad=0.005", facecolor="#E0E0E0", edgecolor="none"))
    ax.text(0.5, 0.77, "Epoch 17 / 30  (57% done)", ha="center",
            fontsize=11, color=WHITE, fontweight="bold", va="center")

    rows = [
        ("Last checkpoint", "Epoch 17  |  val_acc = 59.22%"),
        ("Best so far",     "Epoch 17  |  val_acc = 59.22%"),
        ("Original env",    "TWCC Nano5 (H100 NVL 96 GB)"),
        ("New env",         "Nano4 H200 143 GB  ←  faster"),
        ("Epochs left",     "13 epochs  ≈  5-7 dev job resubmits"),
        ("SLURM script",    "slurm/run_dnabert2_genus_50M.nano4.sh"),
        ("Blocker",         "rsync last.pt (473 MB) from Nano5"),
    ]
    y = 0.61
    for k, v in rows:
        ax.text(0.07, y, k+":", fontsize=10, color=TEAL, fontweight="bold")
        ax.text(0.43, y, v,    fontsize=10, color=NAVY)
        y -= 0.08

    # Right: expected final result
    ax2 = axes[1]; ax2.axis("off")
    ax2.set_xlim(0,1); ax2.set_ylim(0,1)
    ax2.text(0.5, 0.97, "Expected Final Result (Genus Comparison)", ha="center",
             fontsize=13, fontweight="bold", color=NAVY, va="top")

    models = ["MT 13-mer\nflat", "NT-v2\n(LoRA)", "DNABERT-2\n(LoRA, 5M)",
              "DNABERT-2\n(LoRA, 50M)\n[pending]", "Kraken2\n(in-DB)"]
    vals   = [94.25, 64.45, 58.88, 62.0, 77.68]  # 62 is estimate
    colors = [GREEN, TEAL, TEAL, ORANGE, GRAY]

    bars = ax2.barh(models, vals, color=colors, height=0.55, left=35)
    for bar, v, color in zip(bars, vals, colors):
        ax2.text(v+0.5, bar.get_y()+bar.get_height()/2,
                 f"{v:.1f}%{'*' if v==62.0 else ''}", va="center",
                 fontsize=9.5, color=color, fontweight="bold")

    ax2.set_xlim(35, 102)
    ax2.set_xlabel("Genus Read Accuracy (%)", fontsize=10)
    ax2.tick_params(labelsize=9)
    ax2.axvline(64.45, color=TEAL, linestyle="--", alpha=0.4, linewidth=1)
    ax2.text(64.45+0.3, -0.7, "NT-v2\nbaseline", fontsize=8, color=TEAL, alpha=0.7)
    ax2.text(99, 3.1, "* estimate", fontsize=8, color=ORANGE)

    plt.savefig(TMP/"dnabert2_plan.png", bbox_inches="tight"); plt.close()


def fig_mt_benchmark():
    fig, axes = plt.subplots(1, 2, figsize=(13, 5))
    fig.patch.set_facecolor("white")

    # Left: model list
    ax = axes[0]; ax.axis("off")
    ax.set_xlim(0,1); ax.set_ylim(0,1)
    ax.text(0.5, 0.97, "MT Speed Benchmark — 6 Models on H200", ha="center",
            fontsize=13, fontweight="bold", color=NAVY, va="top")

    models = [
        ("1", "MT 13-mer Genus",   "stride-1, 94.25%", GREEN),
        ("2", "MT 6-mer Genus",    "stride-1, 48.76%", TEAL),
        ("3", "MT 13-mer Species", "stride-1, 53.70%", GREEN),
        ("4", "MT 6-mer Species",  "stride-1, 9.00%",  TEAL),
        ("5", "NT-v2 + LoRA",      "sp_v4,  17.83%",   PURPLE),
        ("6", "DNABERT-2 + LoRA",  "genus,  ≥59%",     ORANGE),
    ]
    y = 0.84
    for num, name, detail, color in models:
        ax.add_patch(FancyBboxPatch((0.03, y-0.04), 0.94, 0.085,
                                    boxstyle="round,pad=0.01",
                                    facecolor=color+"22", edgecolor=color, linewidth=1.2))
        ax.text(0.08, y+0.002, f"#{num}", fontsize=11, color=color, fontweight="bold")
        ax.text(0.18, y+0.002, name, fontsize=11, color=NAVY, fontweight="bold")
        ax.text(0.18, y-0.028, detail, fontsize=9.5, color=GRAY)
        y -= 0.115

    ax.text(0.5, 0.04, "Metrics: throughput (reads/sec) · latency (ms/read) · peak GPU (MiB)",
            ha="center", fontsize=9.5, color=GRAY)

    # Right: why this matters
    ax2 = axes[1]; ax2.axis("off")
    ax2.set_xlim(0,1); ax2.set_ylim(0,1)
    ax2.text(0.5, 0.97, "Why Benchmark Matters", ha="center",
             fontsize=13, fontweight="bold", color=NAVY, va="top")

    points = [
        (TEAL,  "Thesis §4.x",
                "MT is 498× smaller (5M vs 498M) — speed\nadvantage is a core selling point"),
        (GREEN, "Journal claim",
                "Need hard numbers: reads/sec, GPU memory\nto support practical deployment argument"),
        (ORANGE,"Architecture comparison",
                "6-mer vs 13-mer: does larger k-mer hurt\nspeed? (tokenise overhead vs model size)"),
        (PURPLE,"H200 baseline",
                "Establishes reference numbers on H200\nfor reproducibility"),
    ]
    y = 0.82
    for color, title, body in points:
        ax2.add_patch(FancyBboxPatch((0.03, y-0.085), 0.94, 0.105,
                                     boxstyle="round,pad=0.01",
                                     facecolor=color+"18", edgecolor=color, linewidth=1.2))
        ax2.text(0.07, y+0.005, title, fontsize=10.5, color=color, fontweight="bold")
        ax2.text(0.07, y-0.03, body,  fontsize=9.5, color=NAVY)
        y -= 0.135

    ax2.add_patch(FancyBboxPatch((0.03,0.03), 0.94, 0.065,
                                  boxstyle="round,pad=0.01",
                                  facecolor="#FFF8E1", edgecolor=GOLD, linewidth=1.5))
    ax2.text(0.5, 0.075, "Blocker: MT models need rsync Taiwana-2 → Nano4",
             ha="center", fontsize=10.5, color=GOLD, fontweight="bold")
    ax2.text(0.5, 0.048, "ETA once rsync done: ~1 dev job (1 hr), all 6 models in one shot",
             ha="center", fontsize=9.5, color=NAVY)

    plt.tight_layout()
    plt.savefig(TMP/"mt_benchmark.png"); plt.close()


def fig_engineering():
    fig, ax = plt.subplots(figsize=(13, 6.5))
    ax.axis("off"); ax.set_xlim(0,1); ax.set_ylim(0,1)
    ax.text(0.5, 0.975, "Engineering Fixes This Week (6 commits)", ha="center",
            fontsize=14, fontweight="bold", color=NAVY, va="top")

    fixes = [
        (GREEN,  "peft key remap",        "480b7dc",
                 "peft ≥0.6 renamed query.weight → query.base_layer.weight. "
                 "Auto-remapped in train.py resume so old TWCC checkpoints load on Nano4."),
        (GREEN,  "Vocab size mismatch",   "5087e3f",
                 "MT 13-mer genus checkpoint had embedding dim 4097 vs expected 4096. "
                 "extract_mt_predictions now truncates/pads embedding row to match."),
        (TEAL,   "data_loader.py added",  "d6fd99b",
                 "data_loader.py was missing from repo (only on TWCC local). "
                 "Added to scripts/ so all envs can use the unified FASTA loader."),
        (TEAL,   "PYTHONPATH / sys.path", "39f9291",
                 "extract_mt_predictions was shadowing MT's own utils package via sys.path. "
                 "Fixed import order; MetaTransformer src now loads correctly."),
        (ORANGE, "Device handler init",   "61239b6",
                 "MT device_handler required explicit cuda init before model.to(device). "
                 "Added init call; also added mid-epoch time_limit_sec checkpoint save."),
        (PURPLE, "resource_monitor.py",   "39f9291",
                 "New script: logs GPU memory + CPU usage every N sec during inference. "
                 "Used by benchmark for peak GPU MiB measurement."),
    ]

    # 2-column layout with generous spacing
    ITEM_H   = 0.125   # normalized height per item box
    ROW_GAP  = 0.025   # gap between rows
    ROW_STEP = ITEM_H + ROW_GAP
    TOP_Y    = 0.905   # y of first row top

    cols = [(0.01, 0.49), (0.51, 0.99)]
    for i, (color, title, commit, body) in enumerate(fixes):
        col_x, col_end = cols[i % 2]
        row = i // 2
        y_top = TOP_Y - row * ROW_STEP
        w = col_end - col_x

        ax.add_patch(FancyBboxPatch((col_x, y_top - ITEM_H), w, ITEM_H,
                                     boxstyle="round,pad=0.008",
                                     facecolor=color+"18", edgecolor=color, linewidth=1.5))
        # Title + commit hash on same line
        ax.text(col_x+0.015, y_top - 0.015, title, fontsize=10.5, color=color,
                fontweight="bold", va="top")
        ax.text(col_end - 0.01, y_top - 0.015, commit, fontsize=8, color=GRAY,
                fontweight="normal", va="top", ha="right", family="monospace")
        # Body (single line, no \n)
        ax.text(col_x+0.015, y_top - 0.055, body, fontsize=8.8, color=NAVY,
                va="top", wrap=True)

    plt.savefig(TMP/"engineering.png", bbox_inches="tight"); plt.close()


def fig_decisions():
    fig, ax = plt.subplots(figsize=(13, 6.5))
    ax.axis("off"); ax.set_xlim(0,1); ax.set_ylim(0,1)
    ax.text(0.5, 0.975, "Decision Board · 6/9", ha="center",
            fontsize=14, fontweight="bold", color=NAVY, va="top")

    # Each item: (title, body_line1, body_line2_or_None, color)
    resolved = [
        ("✓  TWCC budget crisis",
         "Nano4 H200 free trial covers all pending compute — TWCC top-up NOT needed.",
         None, GREEN),
        ("✓  Nano4 onboarding",
         "Sanity check PASSED (Job 71907). Pipeline identical to TWCC baseline (17.83%).",
         None, GREEN),
    ]
    pending = [
        ("Q1  DNABERT-1 50M",
         "11.7 hr/epoch → ~350 dev jobs. Recommend: cancel.",
         "5M checkpoint (61.78% RC TTA) is sufficient for thesis comparison.", ORANGE),
        ("Q2  MT 13-mer hier retraining",
         "Checkpoint corrupted on Taiwana-2. ~1-2 days V100 to retrain.",
         "Adds +2-3 pp to Table 4.29. Decide before 6/15 defense.", ORANGE),
        ("Q3  HMP real-dataset inference",
         "Download SRR072232 mock community (~1-2 GB), run NT-v2 sp_v4.",
         "Pearson r vs known abundance → real-world validation for journal.", ORANGE),
    ]

    ITEM_H  = 0.115
    ITEM_GAP = 0.015
    STEP    = ITEM_H + ITEM_GAP
    HDR_H   = 0.045

    def draw_section(title, color, items, y_start):
        ax.add_patch(Rectangle((0, y_start - HDR_H), 1, HDR_H,
                                facecolor=color+"33", edgecolor="none"))
        ax.text(0.02, y_start - HDR_H/2, title, fontsize=11, color=color,
                fontweight="bold", va="center")
        y = y_start - HDR_H - ITEM_GAP
        for t, b1, b2, c in items:
            ax.add_patch(FancyBboxPatch((0.01, y - ITEM_H), 0.98, ITEM_H,
                                         boxstyle="round,pad=0.008",
                                         facecolor=c+"12", edgecolor=c+"88", linewidth=1.2))
            ax.text(0.03, y - 0.018, t,  fontsize=10.5, color=NAVY, fontweight="bold", va="top")
            ax.text(0.03, y - 0.052, b1, fontsize=9.5,  color=GRAY, va="top")
            if b2:
                ax.text(0.03, y - 0.081, b2, fontsize=9.5, color=GRAY, va="top")
            y -= STEP
        return y - ITEM_GAP

    y = draw_section("RESOLVED ✓", GREEN, resolved, 0.935)
    draw_section("PENDING — need advisor input today", ORANGE, pending, y)

    plt.savefig(TMP/"decisions.png", bbox_inches="tight"); plt.close()


def fig_next_steps():
    fig, ax = plt.subplots(figsize=(13, 5))
    ax.axis("off"); ax.set_xlim(0,1); ax.set_ylim(0,1)
    ax.text(0.5, 0.97, "Next Steps · 6/9 → 6/15 (defense)", ha="center",
            fontsize=14, fontweight="bold", color=NAVY, va="top")

    steps = [
        ("This week",  RED, [
            ("P0-A  DNABERT-2 50M",   "rsync last.pt Nano5→Nano4  →  sbatch dev  →  ~7 resubmits to epoch 30"),
            ("P0-B  MT benchmark",     "rsync MT models Taiwana-2→Nano4  →  sbatch run_mt_benchmark_nano4.sh"),
            ("P0-C  HMP inference",    "download SRR072232  →  run NT-v2 sp_v4  →  Pearson r vs known abundance"),
        ]),
        ("On completion", GREEN, [
            ("RC TTA eval",            "Run eval_rc_tta for DNABERT-2 50M  →  update Table 4.x"),
            ("Benchmark numbers",      "Fill in throughput/latency/peak-GPU in thesis §4.speed"),
            ("Thesis final push",      "Update figures + tables  →  compile  →  submit for binding"),
        ]),
        ("Timeline", TEAL, [
            ("6/15",  "Thesis defense"),
            ("6/30",  "Final thesis submission"),
            ("7-8/2026", "Rewrite Ch4 → Bioinformatics journal manuscript"),
        ]),
    ]

    y = 0.86
    for section, color, items in steps:
        ax.text(0.02, y, section, fontsize=11.5, color=color, fontweight="bold")
        y -= 0.05
        for label, detail in items:
            ax.add_patch(FancyBboxPatch((0.02, y-0.065), 0.96, 0.07,
                                         boxstyle="round,pad=0.01",
                                         facecolor=color+"15", edgecolor=color+"55", linewidth=1))
            ax.text(0.05, y-0.015, label, fontsize=10, color=color, fontweight="bold")
            ax.text(0.30, y-0.015, detail, fontsize=9.5, color=NAVY)
            y -= 0.085
        y -= 0.02

    plt.tight_layout()
    plt.savefig(TMP/"next_steps.png"); plt.close()


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def add_slide(prs, layout=6):
    return prs.slides.add_slide(prs.slide_layouts[layout])

def title_bar(slide, title, subtitle=""):
    W, H = Inches(13.33), Inches(7.5)
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(NAVY)
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = rgb(WHITE)
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(14); r2.font.color.rgb = rgb("#90CAF9")

def add_image(slide, path, left, top, width, height):
    slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                              Inches(width), Inches(height))

def add_text_box(slide, text, left, top, width, height,
                 fontsize=12, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fontsize); r.font.bold = bold; r.font.color.rgb = rgb(color)


def build():
    # Generate figures
    print("Generating figures...")
    fig_timeline()
    fig_nano4_status()
    fig_dnabert2_plan()
    fig_mt_benchmark()
    fig_engineering()
    fig_decisions()
    fig_next_steps()

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    sl = add_slide(prs)
    bg = sl.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = rgb(NAVY); bg.line.fill.background()

    add_text_box(sl, "Weekly Progress Update", 1, 1.6, 11, 1,
                 fontsize=38, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(sl, "Token-Level Metagenomic Classification · NT-v2 + LoRA", 1, 2.8, 11, 0.6,
                 fontsize=16, color="#90CAF9", align=PP_ALIGN.CENTER)
    add_text_box(sl, "June 9, 2026", 1, 3.5, 11, 0.5,
                 fontsize=14, color="#BBDEFB", align=PP_ALIGN.CENTER)

    bullets = "Nano4 H200 onboarded  ·  DNABERT-2 50M resume ready  ·  MT benchmark ready"
    add_text_box(sl, bullets, 1, 4.3, 11, 0.6,
                 fontsize=13, color=GOLD, align=PP_ALIGN.CENTER)

    add_text_box(sl, "YMJ · M2 Lab · NTU CSIE", 1, 6.6, 11, 0.5,
                 fontsize=11, color="#78909C", align=PP_ALIGN.CENTER)

    # ── Slide 2: Timeline ────────────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "Progress Timeline · 6/2 → 6/9")
    add_image(sl, TMP/"timeline.png", 0.2, 1.1, 12.9, 3.8)

    rows = [
        ("6/2", "Advisor meeting (deck 0602) — 3 open questions"),
        ("6/3", "Nano4 sanity check PASSED — NT-Sp 17.83% matches TWCC exactly"),
        ("6/3–7", "6 engineering commits: peft compat, vocab fix, data_loader, benchmark script"),
        ("6/9", "TODAY — scripts ready, waiting to fire P0 jobs"),
    ]
    y = 5.1
    for date, text in rows:
        add_text_box(sl, f"  {date}  {text}", 0.3, y, 12.7, 0.35, fontsize=11, color=NAVY)
        y += 0.37

    # ── Slide 3: Nano4 Status ────────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "Nano4 H200 — Migration Complete", "Free 1-month trial · dev 1hr · normal 12hr (min 64 GPU)")
    add_image(sl, TMP/"nano4_status.png", 0.2, 1.1, 12.9, 5.8)

    # ── Slide 4: DNABERT-2 Plan ──────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "P0-A: DNABERT-2 50M Resume", "Epoch 17/30 · val_acc 59.22% · ~7 dev-job resubmits to complete")
    add_image(sl, TMP/"dnabert2_plan.png", 0.2, 1.1, 12.9, 5.8)

    # ── Slide 5: MT Benchmark ────────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "P0-B: MT Speed / Memory Benchmark", "6 models · H200 · 100K reads · throughput + peak GPU")
    add_image(sl, TMP/"mt_benchmark.png", 0.2, 1.1, 12.9, 5.8)

    # ── Slide 6: Engineering ─────────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "Engineering Highlights (6 commits this week)")
    add_image(sl, TMP/"engineering.png", 0.2, 1.1, 12.9, 5.8)

    # ── Slide 7: Decisions ───────────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "Decision Board · 6/9")
    add_image(sl, TMP/"decisions.png", 0.2, 1.1, 12.9, 5.8)

    # ── Slide 8: Next Steps ──────────────────────────────────────────────────
    sl = add_slide(prs)
    title_bar(sl, "Next Steps · 6/9 → 6/15 (defense)")
    add_image(sl, TMP/"next_steps.png", 0.2, 1.1, 12.9, 5.8)

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()

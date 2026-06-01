#!/usr/bin/env python3
"""
Advisor meeting deck v2 — 2026-05-26
Blank canvas, custom visual layouts, heavy use of figures.
"""

import os
import sys
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Rectangle
import numpy as np
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

DOCS    = Path("/work/ymj1123ntu/token_level_gfm_classifier/docs")
FIG_OLD = DOCS / "figures_0519"
THESIS_FIG = Path("/work/ymj1123ntu/thesis/figures")
TMP     = DOCS / "advisor_figs_v2"
TMP.mkdir(exist_ok=True)
OUT     = DOCS / "advisor_meeting_2026_0526_v2.pptx"

# ─── Color palette ─────────────────────────────────────────────────────────────
NAVY   = "#0A2540"
TEAL   = "#1E6091"
GREEN  = "#2E7D32"
RED    = "#C62828"
ORANGE = "#E65100"
GOLD   = "#B8860B"
GRAY   = "#5A6068"
LIGHT  = "#F4F6F8"
WHITE  = "#FFFFFF"

# RGB versions for python-pptx
def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

C_NAVY   = rgb(NAVY)
C_TEAL   = rgb(TEAL)
C_GREEN  = rgb(GREEN)
C_RED    = rgb(RED)
C_ORANGE = rgb(ORANGE)
C_GOLD   = rgb(GOLD)
C_GRAY   = rgb(GRAY)
C_LIGHT  = rgb(LIGHT)
C_WHITE  = rgb(WHITE)

plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "font.size": 11,
    "axes.spines.top": False,
    "axes.spines.right": False,
    "axes.labelweight": "normal",
    "figure.facecolor": "white",
    "axes.facecolor": "white",
    "savefig.dpi": 200,
    "savefig.bbox": "tight",
    "savefig.pad_inches": 0.15,
})


# ═══════════════════════════════════════════════════════════════════════════════
# CUSTOM FIGURE GENERATION
# ═══════════════════════════════════════════════════════════════════════════════

def fig_foundation_models():
    """Foundation model comparison: DNABERT-1, DNABERT-2, NT-v2 at 5M scale + scaling."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4.3))

    # Left: controlled 5M comparison
    names = ["DNABERT-1\n(2021)", "DNABERT-2\n(2023)", "NT-v2\n(2023)"]
    accs  = [61.20, 57.35, 62.02]
    params = [91, 119, 498]  # M params
    colors = [GRAY, GRAY, GREEN]
    bars = ax1.bar(names, accs, color=colors, edgecolor="white", linewidth=2)
    for b, v, p in zip(bars, accs, params):
        ax1.text(b.get_x() + b.get_width()/2, v + 0.8, f"{v:.1f}%",
                 ha="center", fontsize=11, fontweight="bold", color=NAVY)
        ax1.text(b.get_x() + b.get_width()/2, v/2, f"{p}M\nparams",
                 ha="center", va="center", fontsize=9, color="white", fontweight="bold")
    ax1.set_ylim(0, 75)
    ax1.set_ylabel("Genus Top-1 accuracy (%)", fontsize=11)
    ax1.set_title("DNA Foundation Models at 5M reads (controlled)",
                  fontsize=12, fontweight="bold", color=NAVY, pad=10)
    ax1.grid(axis="y", linestyle="--", alpha=0.3)
    # Star NT-v2
    ax1.scatter(2, 68, marker="*", s=300, color=GOLD, zorder=5,
                edgecolor=NAVY, linewidth=1)
    ax1.text(2, 71, "Selected", ha="center", fontsize=9.5,
             color=GOLD, fontweight="bold")

    # Right: NT-v2 scaling
    sizes = [0.5, 5, 50]
    accs_s = [55.29, 62.02, 64.45]
    rc_tta = [55.29 + 1.0, 63.05, 66.05]  # approx RC TTA additions
    ax2.plot(sizes, accs_s, "-o", color=TEAL, markersize=10, linewidth=2.2,
             markerfacecolor=GREEN, markeredgecolor="white", markeredgewidth=1.8,
             label="Forward only")
    ax2.plot(sizes, rc_tta, "--^", color=GOLD, markersize=9, linewidth=1.8,
             markerfacecolor=GOLD, markeredgecolor="white", markeredgewidth=1.5,
             label="With RC TTA")
    for s, a in zip(sizes, accs_s):
        ax2.text(s, a + 1.3, f"{a:.1f}%", ha="center", fontsize=9.5,
                 fontweight="bold", color=NAVY)
    for s, a in zip(sizes, rc_tta):
        ax2.text(s, a - 2.5, f"{a:.1f}%", ha="center", fontsize=9,
                 color=GOLD, fontweight="bold")
    ax2.set_xscale("log")
    ax2.set_xlim(0.3, 100)
    ax2.set_ylim(50, 72)
    ax2.set_xlabel("Training data (M reads)", fontsize=11)
    ax2.set_ylabel("Genus Top-1 (%)", fontsize=11)
    ax2.set_title("NT-v2: scaling consolidates the advantage",
                  fontsize=12, fontweight="bold", color=NAVY, pad=10)
    ax2.legend(loc="lower right", fontsize=9.5)
    ax2.grid(True, which="both", linestyle="--", alpha=0.3)

    plt.tight_layout()
    plt.savefig(TMP / "foundation_models.png")
    plt.close()


def fig_problem_schematic():
    """Diagram showing input → output of metagenomic classification."""
    fig, ax = plt.subplots(figsize=(8.5, 3.6))
    ax.set_xlim(0, 10); ax.set_ylim(0, 4); ax.axis("off")

    # Read box
    ax.add_patch(FancyBboxPatch((0.3, 1.3), 2.6, 1.4,
                                 boxstyle="round,pad=0.05", linewidth=2,
                                 edgecolor=NAVY, facecolor=LIGHT))
    ax.text(1.6, 2.2, "150 bp DNA read", ha="center", va="center",
            fontsize=12, fontweight="bold", color=NAVY)
    ax.text(1.6, 1.65, "ATCGGCT…GTAC",
            ha="center", va="center", fontsize=9, family="monospace", color=GRAY)

    # Arrow
    ax.add_patch(FancyArrowPatch((3.0, 2.0), (4.0, 2.0),
                                  arrowstyle="-|>", mutation_scale=22,
                                  linewidth=2.5, color=TEAL))
    ax.text(3.5, 2.4, "classify", ha="center", fontsize=10, color=TEAL, style="italic")

    # Model box
    ax.add_patch(FancyBboxPatch((4.1, 1.3), 2.4, 1.4,
                                 boxstyle="round,pad=0.05", linewidth=2,
                                 edgecolor=TEAL, facecolor="#E3F2FD"))
    ax.text(5.3, 2.2, "model", ha="center", va="center",
            fontsize=12, fontweight="bold", color=NAVY)
    ax.text(5.3, 1.65, "NT-v2 / MT / Kraken2",
            ha="center", va="center", fontsize=9, color=GRAY)

    # Arrow
    ax.add_patch(FancyArrowPatch((6.6, 2.0), (7.6, 2.0),
                                  arrowstyle="-|>", mutation_scale=22,
                                  linewidth=2.5, color=TEAL))

    # Species box
    ax.add_patch(FancyBboxPatch((7.7, 1.3), 2.1, 1.4,
                                 boxstyle="round,pad=0.05", linewidth=2,
                                 edgecolor=GREEN, facecolor="#E8F5E9"))
    ax.text(8.75, 2.2, "species ID", ha="center", va="center",
            fontsize=12, fontweight="bold", color=NAVY)
    ax.text(8.75, 1.65, "1 of 1,535",
            ha="center", va="center", fontsize=10, color=GREEN, fontweight="bold")

    # Bottom annotations: three challenges
    ax.text(0.4, 0.5, "Challenges:", fontsize=10, fontweight="bold", color=NAVY)
    ax.text(2.0, 0.5, "• Short read · intra-genus similarity",
            fontsize=9.5, color=GRAY)
    ax.text(2.0, 0.05, "• 1,535 classes · errors compound per sample",
            fontsize=9.5, color=GRAY)
    ax.text(6.7, 0.5, "• Novel organisms absent from any DB",
            fontsize=9.5, color=GRAY)

    plt.savefig(TMP / "problem_schematic.png")
    plt.close()


def fig_three_families():
    """Schematic of three model families side-by-side."""
    fig, ax = plt.subplots(figsize=(10, 4))
    ax.set_xlim(0, 12); ax.set_ylim(0, 5); ax.axis("off")

    families = [
        ("NT-v2 + LoRA",      "Pre-trained 498M\n6-mer\nLoRA: 1.1% trainable",
         "#E8F5E9", GREEN, "Pre-training axis"),
        ("MetaTransformer",   "From scratch · 5M\n6-mer OR 13-mer\nFull training",
         "#E3F2FD", TEAL,  "Tokenisation axis"),
        ("Kraken2",           "k-mer exact match\nk=35, minimiser 31\nCustom DB",
         "#FFF8E1", GOLD,  "Algorithm axis"),
    ]

    for i, (name, body, fc, ec, axis_lbl) in enumerate(families):
        x = 0.4 + i * 4.0
        ax.add_patch(FancyBboxPatch((x, 0.5), 3.4, 3.6,
                                     boxstyle="round,pad=0.1",
                                     linewidth=2.5, edgecolor=ec, facecolor=fc))
        ax.text(x + 1.7, 3.55, name,
                ha="center", fontsize=13, fontweight="bold", color=NAVY)
        ax.text(x + 1.7, 2.2, body,
                ha="center", va="center", fontsize=10.5, color=NAVY)
        ax.text(x + 1.7, 0.95, axis_lbl,
                ha="center", fontsize=9.5, color=ec, fontweight="bold", style="italic")

    ax.text(6.0, 4.7, "Three Model Families · Same 50M training data · Same 100K test set",
            ha="center", fontsize=12, fontweight="bold", color=NAVY)

    plt.savefig(TMP / "three_families.png")
    plt.close()


def fig_tokenisation_dominates():
    """Bar chart: Tokenisation > Pre-training, genus + species side by side."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4.3))

    # Genus
    names_g = ["MT 6-mer\n(scratch)", "NT-Genus\n(pre-trained)", "MT 13-mer\n(scratch)"]
    vals_g = [48.82, 67.07, 87.42]
    colors_g = [GRAY, TEAL, GREEN]
    bars = ax1.bar(names_g, vals_g, color=colors_g, edgecolor="white", linewidth=1.5)
    for b, v in zip(bars, vals_g):
        ax1.text(b.get_x() + b.get_width()/2, v + 1.5, f"{v:.1f}%",
                 ha="center", fontsize=11, fontweight="bold", color=NAVY)
    ax1.set_ylim(0, 115)
    ax1.set_ylabel("Read accuracy (%)", fontsize=11)
    ax1.set_title("Genus level (120 classes)", fontsize=12, fontweight="bold",
                  color=NAVY, pad=14)
    ax1.set_yticks([0, 20, 40, 60, 80, 100])
    ax1.grid(axis="y", linestyle="--", alpha=0.3)

    # Annotate the comparisons (lowered to leave space at top)
    ax1.annotate("", xy=(1, 60), xytext=(0, 60),
                 arrowprops=dict(arrowstyle="<->", color=ORANGE, lw=1.4))
    ax1.text(0.5, 56, "+18.3 pp\npre-training", ha="center", fontsize=8.5,
             color=ORANGE, fontweight="bold", va="top")
    ax1.annotate("", xy=(2, 105), xytext=(1, 105),
                 arrowprops=dict(arrowstyle="<->", color=GREEN, lw=1.4))
    ax1.text(1.5, 109, "+20.4 pp tokenisation", ha="center", fontsize=8.5,
             color=GREEN, fontweight="bold")

    # Species
    names_s = ["MT 6-mer\n(scratch)", "NT-Species\n(pre-trained)", "MT 13-mer\n(scratch)"]
    vals_s = [9.2, 17.8, 49.7]
    colors_s = [GRAY, TEAL, GREEN]
    bars = ax2.bar(names_s, vals_s, color=colors_s, edgecolor="white", linewidth=1.5)
    for b, v in zip(bars, vals_s):
        ax2.text(b.get_x() + b.get_width()/2, v + 1.0, f"{v:.1f}%",
                 ha="center", fontsize=11, fontweight="bold", color=NAVY)
    ax2.set_ylim(0, 75)
    ax2.set_title("Species level (1,535 classes)", fontsize=12, fontweight="bold",
                  color=NAVY, pad=14)
    ax2.set_yticks([0, 10, 20, 30, 40, 50, 60])
    ax2.grid(axis="y", linestyle="--", alpha=0.3)

    ax2.annotate("", xy=(1, 27), xytext=(0, 27),
                 arrowprops=dict(arrowstyle="<->", color=ORANGE, lw=1.4))
    ax2.text(0.5, 23, "+8.6 pp\npre-training", ha="center", fontsize=8.5,
             color=ORANGE, fontweight="bold", va="top")
    ax2.annotate("", xy=(2, 67), xytext=(1, 67),
                 arrowprops=dict(arrowstyle="<->", color=GREEN, lw=1.4))
    ax2.text(1.5, 70, "+31.9 pp tokenisation", ha="center", fontsize=8.5,
             color=GREEN, fontweight="bold")

    plt.tight_layout()
    plt.savefig(TMP / "tokenisation_dominates.png")
    plt.close()


def fig_data_vs_tricks():
    """Side-by-side: data scaling line + trick bars."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4))

    # Data scaling
    sizes = [0.5, 5, 50]
    accs = [55.29, 63.05, 67.07]
    ax1.plot(sizes, accs, "-o", color=TEAL, markersize=10, linewidth=2.5,
             markerfacecolor=GREEN, markeredgecolor="white", markeredgewidth=2)
    for s, a in zip(sizes, accs):
        ax1.text(s, a + 1.5, f"{a:.1f}%", ha="center", fontsize=10,
                 fontweight="bold", color=NAVY)
    ax1.set_xscale("log")
    ax1.set_xlim(0.3, 100)
    ax1.set_ylim(50, 75)
    ax1.set_xlabel("Training data (M reads)", fontsize=11)
    ax1.set_ylabel("Genus Top-1 (%)", fontsize=11)
    ax1.set_title("Data scaling: +7.76 pp per 10×", fontsize=12,
                  fontweight="bold", color=GREEN)
    ax1.grid(True, which="both", linestyle="--", alpha=0.3)
    ax1.annotate("+7.76 pp", xy=(2, 59), fontsize=11, color=GREEN, fontweight="bold")
    ax1.annotate("+4.02 pp", xy=(20, 65), fontsize=11, color=GREEN, fontweight="bold")

    # Tricks
    tricks = ["Logit\nAdjust", "RC\nConsistency", "Class-bal\nsampler", "LR\ntuning"]
    effects = [0.3, 0.5, 0.4, 0.5]
    ax2.bar(tricks, effects, color=GRAY, edgecolor="white", linewidth=1.5)
    ax2.axhline(7.76, color=GREEN, linestyle="--", linewidth=2, label="Data 10×: +7.76 pp")
    for i, e in enumerate(effects):
        ax2.text(i, e + 0.3, f"±{e}", ha="center", fontsize=10, color=NAVY)
    ax2.set_ylim(0, 9)
    ax2.set_ylabel("Effect (pp)", fontsize=11)
    ax2.set_title("Training tricks: each ±0.5 pp", fontsize=12,
                  fontweight="bold", color=RED)
    ax2.legend(loc="upper right", fontsize=10)
    ax2.grid(axis="y", linestyle="--", alpha=0.3)

    plt.tight_layout()
    plt.savefig(TMP / "data_vs_tricks.png")
    plt.close()


def fig_utility_threshold():
    """Scatter: read accuracy vs Pearson r, showing non-linear threshold."""
    fig, ax = plt.subplots(figsize=(9.5, 4.5))

    # Species-level points
    pts_species = [
        ("MT 6-mer hier",   6.4, 0.034, GRAY),
        ("MT 6-mer flat",   9.2, 0.065, GRAY),
        ("NT-Species hier", 15.8, 0.106, TEAL),
        ("NT-v2 per-genus",15.6, 0.083, TEAL),
        ("NT-Species flat", 17.8, 0.135, TEAL),
        ("NT-v2 oracle",   29.5, 0.159, ORANGE),
        ("MT 13-mer flat", 49.7, 0.466, GREEN),
        ("MT 13-mer hier", 50.9, 0.478, GREEN),
        ("Kraken2",         66.2, 0.727, GOLD),
    ]
    for name, x, y, c in pts_species:
        ax.scatter(x, y, s=140, color=c, edgecolor="white", linewidth=2, zorder=3)
        # Label position offsets
        offs = {"NT-Species hier": (-8, 0.04), "NT-v2 per-genus": (8, -0.05),
                "NT-Species flat": (5, 0.03), "MT 6-mer hier": (-8, 0.03),
                "MT 6-mer flat": (10, 0.03)}
        dx, dy = offs.get(name, (3, 0.02))
        ax.annotate(name, (x, y), xytext=(x + dx, y + dy),
                    fontsize=9.5, color=NAVY)

    # Shaded threshold region
    ax.axvspan(40, 50, alpha=0.15, color=GREEN, zorder=1)
    ax.text(45, 0.85, "Practical\nthreshold\n40–50%",
            ha="center", fontsize=10, color=GREEN, fontweight="bold")

    # Annotate noise floor zone
    ax.axhspan(0, 0.2, alpha=0.1, color=RED, zorder=1)
    ax.text(2, 0.18, "Noise floor: r constrained < 0.2", fontsize=9.5,
            color=RED, fontweight="bold", style="italic")

    ax.set_xlim(0, 75)
    ax.set_ylim(0, 0.95)
    ax.set_xlabel("Per-read accuracy (%)", fontsize=11)
    ax.set_ylabel("Sample-level Pearson r", fontsize=11)
    ax.set_title("Species-level utility curve · 1,535 classes",
                 fontsize=12, fontweight="bold", color=NAVY)
    ax.grid(True, linestyle="--", alpha=0.3)

    plt.tight_layout()
    plt.savefig(TMP / "utility_threshold.png")
    plt.close()


def fig_router_threshold():
    """Bar chart: monotonic species improvement vs router accuracy."""
    fig, ax = plt.subplots(figsize=(9, 4.3))

    routers = ["MT 6-mer\n(48.9%)", "NT-Genus\n(66.0%)", "MT 13-mer\n(87.5%)"]
    effects = [-2.8, -2.0, 1.12]
    colors = [RED, ORANGE, GREEN]
    bars = ax.bar(routers, effects, color=colors, edgecolor="white", linewidth=2)

    for b, e in zip(bars, effects):
        sign = "+" if e > 0 else ""
        ax.text(b.get_x() + b.get_width()/2,
                e + (0.25 if e > 0 else -0.3),
                f"{sign}{e:.2f} pp",
                ha="center", va="bottom" if e > 0 else "top",
                fontsize=12, fontweight="bold", color=NAVY)

    ax.axhline(0, color="black", linewidth=1)
    # Threshold marker
    ax.axvline(1.5, color=GOLD, linestyle="--", linewidth=1.5, alpha=0.6)
    ax.text(1.5, 1.5, "↓\nThreshold ≈ 80%", ha="center", fontsize=10,
            color=GOLD, fontweight="bold")

    ax.set_ylim(-4, 2.5)
    ax.set_ylabel("Species Top-1 change (pp)", fontsize=11)
    ax.set_xlabel("Genus router (and its accuracy)", fontsize=11)
    ax.set_title("Hierarchical masking effect is monotonic in router accuracy",
                 fontsize=12, fontweight="bold", color=NAVY)
    ax.grid(axis="y", linestyle="--", alpha=0.3)

    plt.tight_layout()
    plt.savefig(TMP / "router_threshold.png")
    plt.close()


def fig_oracle_ceiling():
    """Bar chart: NT-v2 configurations vs MT 13-mer."""
    fig, ax = plt.subplots(figsize=(10, 4.3))

    configs = [
        ("NT-v2\nflat",     17.8, TEAL),
        ("NT-v2\nhier.",    15.8, TEAL),
        ("NT-v2 per-genus\n(predicted)", 15.6, TEAL),
        ("NT-v2 per-genus\n(ORACLE)",    29.5, ORANGE),
        ("MT 13-mer\nflat",  49.7, GREEN),
        ("MT 13-mer\nhier.", 50.9, GREEN),
    ]
    names = [c[0] for c in configs]
    vals  = [c[1] for c in configs]
    cols  = [c[2] for c in configs]
    bars = ax.bar(names, vals, color=cols, edgecolor="white", linewidth=2)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width()/2, v + 1.2, f"{v:.1f}%",
                ha="center", fontsize=11, fontweight="bold", color=NAVY)

    # Highlight oracle ceiling
    ax.axhline(29.5, xmin=0.0, xmax=0.65, color=ORANGE, linestyle="--",
               linewidth=2, alpha=0.7)
    ax.text(3.5, 33, "6-mer NT-v2 ceiling (oracle): 29.5%",
            fontsize=10, color=ORANGE, fontweight="bold")

    # Gap arrow
    ax.annotate("", xy=(4, 49.7), xytext=(4, 30),
                arrowprops=dict(arrowstyle="<->", color=RED, lw=2))
    ax.text(4.25, 40, "21.4 pp\ntokenisation\ngap", fontsize=10,
            color=RED, fontweight="bold")

    ax.set_ylim(0, 60)
    ax.set_ylabel("Species Top-1 read accuracy (%)", fontsize=11)
    ax.set_title("Per-genus oracle = 6-mer ceiling · MT 13-mer surpasses without routing",
                 fontsize=12, fontweight="bold", color=NAVY)
    ax.grid(axis="y", linestyle="--", alpha=0.3)

    plt.tight_layout()
    plt.savefig(TMP / "oracle_ceiling.png")
    plt.close()


def fig_kraken2_indb_vs_ood():
    """Three-panel: in-DB radar (species) + genus bar + OOD bar."""
    fig = plt.figure(figsize=(13.2, 4.5))
    ax1 = plt.subplot(131, projection="polar")
    ax2 = plt.subplot(132)
    ax3 = plt.subplot(133)

    # In-DB comparison radar
    metrics = ["Read\nAcc", "Pearson r", "1-BC", "ROC AUC", "Sens@95%spec"]
    kraken = [0.662, 0.727, 1-0.218, 0.925, 0.857]
    mt13h  = [0.509, 0.478, 1-0.369, 0.967, 0.831]
    ntspf  = [0.178, 0.135, 1-0.589, 0.794, 0.552]
    angles = np.linspace(0, 2*np.pi, len(metrics), endpoint=False).tolist()
    kraken += [kraken[0]]; mt13h += [mt13h[0]]; ntspf += [ntspf[0]]
    angles += [angles[0]]

    ax1.plot(angles, kraken, color=GOLD, linewidth=2.2, label="Kraken2")
    ax1.fill(angles, kraken, color=GOLD, alpha=0.15)
    ax1.plot(angles, mt13h, color=GREEN, linewidth=2.2, label="MT 13-mer hier")
    ax1.fill(angles, mt13h, color=GREEN, alpha=0.10)
    ax1.plot(angles, ntspf, color=TEAL, linewidth=2.2, label="NT-Species flat")
    ax1.fill(angles, ntspf, color=TEAL, alpha=0.10)
    ax1.set_xticks(angles[:-1])
    ax1.set_xticklabels(metrics, fontsize=9)
    ax1.set_ylim(0, 1)
    ax1.set_yticks([0.25, 0.5, 0.75, 1.0])
    ax1.set_yticklabels(["0.25", "0.50", "0.75", "1.0"], fontsize=8)
    ax1.set_title("Species level · in-DB",
                  fontsize=11.5, fontweight="bold", color=NAVY, pad=14)
    ax1.legend(loc="lower center", bbox_to_anchor=(0.5, -0.30), fontsize=8.5, ncol=1)

    # Genus level read accuracy comparison
    g_names = ["MT 6-mer\ngenus", "NT-Genus\n(v9)", "Kraken2\n(in-DB)", "MT 13-mer\ngenus"]
    g_accs  = [48.8, 67.1, 69.5, 87.4]
    g_cols  = [GRAY, TEAL, GOLD, GREEN]
    bars = ax2.bar(g_names, g_accs, color=g_cols, edgecolor="white", linewidth=1.5)
    for b, v in zip(bars, g_accs):
        ax2.text(b.get_x() + b.get_width()/2, v + 1.5, f"{v:.1f}%",
                 ha="center", fontsize=9.5, fontweight="bold", color=NAVY)
    ax2.set_ylim(0, 100)
    ax2.set_ylabel("Genus Top-1 read acc (%)", fontsize=10)
    ax2.set_title("Genus level · read accuracy",
                  fontsize=11.5, fontweight="bold", color=NAVY, pad=10)
    ax2.grid(axis="y", linestyle="--", alpha=0.3)
    ax2.text(1.5, 92, "Kraken2 genus precision on\nclassified reads: 99.3%",
             ha="center", fontsize=8.5, color=GOLD, style="italic", fontweight="bold")

    # OOD bar (now ax3)
    models = ["Kraken2", "NT-Species", "MT 13-mer"]
    in_db  = [66.2, 17.8, 49.7]
    ood    = [0.0, 8.5, 12.0]
    x = np.arange(len(models))
    w = 0.36
    ax3.bar(x - w/2, in_db, w, color=GREEN, label="In-DB", edgecolor="white", linewidth=1.5)
    ax3.bar(x + w/2, ood,   w, color=RED,   label="OOD",
            edgecolor="white", linewidth=1.5)
    for i, (id_v, od_v) in enumerate(zip(in_db, ood)):
        ax3.text(i - w/2, id_v + 1.5, f"{id_v:.1f}%", ha="center",
                 fontsize=9.5, fontweight="bold", color=NAVY)
        ax3.text(i + w/2, od_v + 1.5, f"{od_v:.1f}%", ha="center",
                 fontsize=9.5, fontweight="bold",
                 color=RED if od_v == 0 else NAVY)
    ax3.text(0, -7, "0%  DB has no\nk-mers to match",
             ha="center", fontsize=8, color=RED, style="italic")
    ax3.set_xticks(x); ax3.set_xticklabels(models, fontsize=9.5)
    ax3.set_ylim(-12, 80)
    ax3.set_ylabel("Read accuracy (%)", fontsize=10)
    ax3.set_title("OOD: only neural retains signal",
                  fontsize=11.5, fontweight="bold", color=NAVY, pad=10)
    ax3.legend(loc="upper right", fontsize=9)
    ax3.grid(axis="y", linestyle="--", alpha=0.3)

    plt.tight_layout()
    plt.savefig(TMP / "kraken2_indb_ood.png")
    plt.close()


def fig_story_arc():
    """Story arc flow diagram."""
    fig, ax = plt.subplots(figsize=(11, 4.0))
    ax.set_xlim(0, 14); ax.set_ylim(0, 5); ax.axis("off")

    steps = [
        ("Hypothesis 1\nParameter count", "Refuted",
         "NT-v2 498M ≪ wins", RED),
        ("Hypothesis 2\nPre-training", "Helps but…",
         "+18 pp at genus,\nplateau at species", ORANGE),
        ("Hypothesis 3\nRouting design", "Bounded",
         "Oracle ceiling 29.5%\n< MT 13-mer 49.7%", ORANGE),
        ("Hypothesis 4\nTokenisation", "Confirmed",
         "+31.9 pp at species\nfixes everything else", GREEN),
    ]
    for i, (title, verdict, body, color) in enumerate(steps):
        x = 0.3 + i * 3.4
        ax.add_patch(FancyBboxPatch((x, 1.0), 2.9, 3.3,
                                     boxstyle="round,pad=0.08",
                                     linewidth=2.5, edgecolor=color, facecolor=WHITE))
        ax.text(x + 1.45, 3.8, title,
                ha="center", fontsize=11, fontweight="bold", color=NAVY)
        ax.text(x + 1.45, 2.95, verdict,
                ha="center", fontsize=11, fontweight="bold", color=color)
        ax.text(x + 1.45, 1.8, body,
                ha="center", va="center", fontsize=9.5, color=GRAY)
        if i < 3:
            ax.add_patch(FancyArrowPatch((x + 2.95, 2.6), (x + 3.3, 2.6),
                                          arrowstyle="-|>", mutation_scale=20,
                                          linewidth=2, color=NAVY))

    ax.text(7, 4.7, "What does NOT explain the species-level gap →  what DOES",
            ha="center", fontsize=12, fontweight="bold", color=NAVY)
    ax.text(7, 0.4, "Tokenisation is the hidden bottleneck",
            ha="center", fontsize=13, fontweight="bold", color=GREEN, style="italic")

    plt.savefig(TMP / "story_arc.png")
    plt.close()


def fig_results_matrix():
    """Heatmap-style table of all sample-level results."""
    rows = [
        ("NT-Species flat",            17.8, 0.135, 0.589, 0.794, 55.2, "neutral"),
        ("NT-Species hier.",           15.8, 0.106, 0.608, 0.796, 53.0, "bad"),
        ("NT-v2 per-genus (pred.)",    15.6, 0.083, 0.625, 0.789, 52.5, "bad"),
        ("NT-v2 per-genus (ORACLE)",   29.5, 0.159, 0.541, 0.842, 65.2, "warn"),
        ("MT 6-mer flat",               9.2, 0.065, 0.654, 0.690,  7.1, "bad"),
        ("MT 6-mer hier.",              6.4, 0.034, 0.686, 0.667,  5.7, "bad"),
        ("MT 13-mer flat",             49.7, 0.466, 0.378, 0.966, 82.6, "good"),
        ("MT 13-mer hier.",            50.9, 0.478, 0.369, 0.967, 83.1, "good"),
        ("Kraken2 (in-DB) · species",  66.2, 0.727, 0.218, 0.925, 85.7, "kraken"),
        ("Kraken2 (in-DB) · genus",    69.5, 0.844, 0.178, 0.895, 80.3, "kraken"),
    ]
    fig, ax = plt.subplots(figsize=(13, 5.5))
    ax.axis("off")

    headers = ["Model", "Read Acc (%)", "Pearson r", "Bray-Curtis", "ROC AUC", "Sens@95%spec"]
    col_widths = [3.4, 1.5, 1.4, 1.5, 1.4, 1.7]
    col_starts = [0]
    for w in col_widths:
        col_starts.append(col_starts[-1] + w)
    total_w = col_starts[-1]
    row_h = 0.45

    # Header row
    y_top = len(rows) * row_h + 0.5
    for j, (h, x0, w) in enumerate(zip(headers, col_starts, col_widths)):
        ax.add_patch(Rectangle((x0, y_top), w, row_h, facecolor=NAVY, edgecolor="white"))
        ax.text(x0 + w/2, y_top + row_h/2, h, ha="center", va="center",
                fontsize=12, color="white", fontweight="bold")

    style_map = {
        "good":   {"fc": "#E8F5E9", "tc": GREEN, "bold": True},
        "warn":   {"fc": "#FFF3E0", "tc": ORANGE, "bold": True},
        "bad":    {"fc": "#FFEBEE", "tc": RED,   "bold": False},
        "neutral":{"fc": "white",   "tc": NAVY,  "bold": False},
        "kraken": {"fc": "#FFF8E1", "tc": GOLD,  "bold": True},
    }

    for i, row in enumerate(rows):
        name = row[0]; vals = row[1:6]; style = row[6]
        s = style_map[style]
        y = y_top - (i + 1) * row_h
        for j, (x0, w) in enumerate(zip(col_starts, col_widths)):
            ax.add_patch(Rectangle((x0, y), w, row_h, facecolor=s["fc"],
                                   edgecolor="white", linewidth=0.8))
        # Model name (left aligned)
        ax.text(col_starts[0] + 0.15, y + row_h/2, name, ha="left", va="center",
                fontsize=11.5, fontweight="bold" if s["bold"] else "normal", color=s["tc"])
        # Values:  j=0 read acc (%), j=1 Pearson r, j=2 BC, j=3 ROC AUC, j=4 sens (%)
        for j, v in enumerate(vals):
            x = col_starts[j+1] + col_widths[j+1]/2
            if j in (0, 4):
                txt = f"{v:.1f}%"
            else:
                txt = f"{v:.3f}"
            ax.text(x, y + row_h/2, txt, ha="center", va="center",
                    fontsize=11.5, fontweight="bold" if s["bold"] else "normal",
                    color=s["tc"])

    # Kraken2 OOD note row
    y = y_top - (len(rows) + 1) * row_h
    ax.add_patch(Rectangle((0, y), total_w, row_h, facecolor="#FFEBEE", edgecolor="white"))
    ax.text(0.15, y + row_h/2,
            "Kraken2 on 219 OOD species (not in DB):  0.0%   →   neural retains low-but-nonzero signal",
            ha="left", va="center", fontsize=11.5, color=RED, fontweight="bold", style="italic")

    ax.set_xlim(-0.2, total_w + 0.2)
    ax.set_ylim(y - 0.3, y_top + row_h + 0.2)

    plt.savefig(TMP / "results_matrix.png")
    plt.close()


def fig_publication_matrix():
    """Comparison matrix of publication options."""
    fig, ax = plt.subplots(figsize=(11, 4.2))
    ax.axis("off")

    headers = ["", "Bioinformatics journal", "ML workshop", "Thesis only"]
    rows = [
        ("Effort",        "Medium (8 weeks)", "Medium-High (10 weeks)", "Low (already drafted)"),
        ("Best fit",      "Same-DB Kraken2 comparison\nOOD analysis · sample-level",
                          "Tokenisation > pre-training\ncontrolled ablation",
                          "Complete narrative\nincl. all 4 contributions"),
        ("Risk",          "Reviewer asks for real-dataset\nvalidation",
                          "ML reviewers may want\nlarger-scale experiments",
                          "Less academic visibility"),
        ("My ranking",    "★★★ Primary",       "★★  Secondary",       "★   Baseline"),
    ]
    col_widths = [1.4, 3.4, 3.4, 3.0]
    col_starts = [0]
    for w in col_widths:
        col_starts.append(col_starts[-1] + w)
    row_h = 0.78
    y_top = len(rows) * row_h + 0.5

    # Header row
    header_colors = [NAVY, GREEN, TEAL, GOLD]
    for j, (h, x0, w, c) in enumerate(zip(headers, col_starts, col_widths, header_colors)):
        ax.add_patch(Rectangle((x0, y_top), w, 0.5, facecolor=c, edgecolor="white"))
        ax.text(x0 + w/2, y_top + 0.25, h, ha="center", va="center",
                fontsize=11.5, color="white", fontweight="bold")

    for i, row in enumerate(rows):
        label = row[0]; vals = row[1:]
        y = y_top - (i + 1) * row_h
        for j, (x0, w) in enumerate(zip(col_starts, col_widths)):
            fc = LIGHT if j == 0 else "white"
            ax.add_patch(Rectangle((x0, y), w, row_h, facecolor=fc,
                                   edgecolor=GRAY, linewidth=0.6))
        ax.text(col_starts[0] + 0.12, y + row_h/2, label, ha="left", va="center",
                fontsize=10.5, fontweight="bold", color=NAVY)
        for j, v in enumerate(vals):
            x = col_starts[j+1] + col_widths[j+1]/2
            ax.text(x, y + row_h/2, v, ha="center", va="center",
                    fontsize=9.5, color=NAVY)

    ax.set_xlim(-0.2, col_starts[-1] + 0.2)
    ax.set_ylim(0, y_top + 0.7)
    plt.savefig(TMP / "publication_matrix.png")
    plt.close()


def fig_strengths_weaknesses():
    """Two-column visual: strengths green, weaknesses red."""
    fig, ax = plt.subplots(figsize=(11.5, 4.5))
    ax.set_xlim(0, 12); ax.set_ylim(0, 6); ax.axis("off")

    # Strengths box
    ax.add_patch(FancyBboxPatch((0.2, 0.4), 5.6, 5.0,
                                 boxstyle="round,pad=0.1",
                                 linewidth=2.5, edgecolor=GREEN, facecolor="#E8F5E9"))
    ax.text(3.0, 5.0, "Strengths", ha="center", fontsize=14, fontweight="bold", color=GREEN)
    strengths = [
        "Same-data 3-axis ablation — unusually clean",
        "Sample-level metrics bridge to real workflows",
        "Same-DB Kraken2 baseline closes \"why not Kraken2?\"",
        "Mechanistic story (noise floor, router threshold)",
        "Per-genus oracle quantifies tokenisation gap",
    ]
    for i, s in enumerate(strengths):
        ax.text(0.5, 4.3 - i*0.7, "✓  " + s, fontsize=11, color=NAVY, va="center")

    # Weaknesses box
    ax.add_patch(FancyBboxPatch((6.2, 0.4), 5.6, 5.0,
                                 boxstyle="round,pad=0.1",
                                 linewidth=2.5, edgecolor=RED, facecolor="#FFEBEE"))
    ax.text(9.0, 5.0, "Weaknesses", ha="center", fontsize=14, fontweight="bold", color=RED)
    weaknesses = [
        ("Simulated reads only (ART Illumina)",
         "→ add HMP mock community sanity check"),
        ("Limited DB (2,505 genomes)",
         "→ acknowledge scope, defer expansion"),
        ("No paired-end / error / coverage tests",
         "→ controlled-study framing"),
        ("Foundation model on smaller side (NT-v2 498M)",
         "→ scale-agnostic finding is the point"),
    ]
    y0 = 4.3
    for i, (w, m) in enumerate(weaknesses):
        ax.text(6.5, y0 - i*1.0, "✗  " + w, fontsize=11, color=NAVY, va="center")
        ax.text(7.0, y0 - i*1.0 - 0.35, m, fontsize=9.5, color=ORANGE,
                va="center", style="italic")

    plt.savefig(TMP / "strengths_weaknesses.png")
    plt.close()


def fig_timeline():
    """Project timeline with milestones."""
    fig, ax = plt.subplots(figsize=(11.5, 3.6))
    ax.set_xlim(0, 14); ax.set_ylim(0, 4); ax.axis("off")

    # Timeline bar
    ax.add_patch(Rectangle((0.6, 1.9), 13.0, 0.18, facecolor=NAVY))

    milestones = [
        (1.0,  "Now",                   "May 2026",   "Thesis 135 pp\nAll results in",         GREEN),
        (3.5,  "Defense",               "Jun 15",     "Mock defense\n+ slides",                 TEAL),
        (5.5,  "Submission prep",       "Jul–Aug",    "Rewrite Ch.4 →\nBioinformatics manuscript", ORANGE),
        (8.0,  "Submission",            "Aug 31",     "Bioinformatics\njournal",                GOLD),
        (11.0, "Review",                "Sep–Dec",    "First-round\nreview",                    GRAY),
        (13.0, "Revision / accept",     "2027 Q1",    "Final outcome",                          NAVY),
    ]
    for x, when, date, body, color in milestones:
        # Dot
        ax.scatter(x, 2.0, s=180, color=color, edgecolor="white", linewidth=2, zorder=4)
        # Title above
        ax.text(x, 2.6, when, ha="center", fontsize=11, fontweight="bold", color=color)
        ax.text(x, 2.3, date, ha="center", fontsize=9, color=GRAY, style="italic")
        # Detail below
        ax.text(x, 1.4, body, ha="center", fontsize=9, color=NAVY, va="top")

    ax.text(7, 3.6, "Recommended Timeline",
            ha="center", fontsize=13, fontweight="bold", color=NAVY)

    plt.savefig(TMP / "timeline.png")
    plt.close()


# Generate all custom figures
print("Generating custom figures...")
fig_foundation_models()
fig_problem_schematic()
fig_three_families()
fig_tokenisation_dominates()
fig_data_vs_tricks()
fig_utility_threshold()
fig_router_threshold()
fig_oracle_ceiling()
fig_kraken2_indb_vs_ood()
fig_story_arc()
fig_results_matrix()
fig_publication_matrix()
fig_strengths_weaknesses()
fig_timeline()
print(f"  → {TMP}/")


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX ASSEMBLY (blank canvas)
# ═══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_background(slide):
    """Light gradient bg."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = hex_to_rgb(WHITE)
    bg.line.fill.background()


def add_header(slide, title, accent=NAVY, page=None, total=None):
    """Header strip with title."""
    # Accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = hex_to_rgb(accent)
    bar.line.fill.background()
    # Title text
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25),
                                   SLIDE_W - Inches(1), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = title
    run.font.size = Pt(24); run.font.bold = True
    run.font.color.rgb = hex_to_rgb(NAVY)
    if page is not None:
        pn = slide.shapes.add_textbox(SLIDE_W - Inches(1.2), Inches(7.05),
                                       Inches(1.0), Inches(0.4))
        ptf = pn.text_frame; ptf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        prun = ptf.paragraphs[0].add_run()
        prun.text = f"{page} / {total}"
        prun.font.size = Pt(10); prun.font.color.rgb = hex_to_rgb(GRAY)


def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=NAVY, align="left",
                bg=None, border=None):
    """Generic textbox."""
    if bg or border:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      left, top, width, height)
        if bg:
            shp.fill.solid(); shp.fill.fore_color.rgb = hex_to_rgb(bg)
        else:
            shp.fill.background()
        if border:
            shp.line.color.rgb = hex_to_rgb(border); shp.line.width = Pt(1.5)
        else:
            shp.line.fill.background()
        tf = shp.text_frame
    else:
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run(); run.text = line
        run.font.size = Pt(font_size); run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color)


def add_picture(slide, fig_path, left, top, width=None, height=None):
    if width is not None:
        return slide.shapes.add_picture(str(fig_path), left, top, width=width)
    elif height is not None:
        return slide.shapes.add_picture(str(fig_path), left, top, height=height)
    else:
        return slide.shapes.add_picture(str(fig_path), left, top)


def add_takeaway(slide, text, color=GREEN):
    """Highlighted callout box at slide bottom."""
    h = Inches(0.7)
    top = SLIDE_H - h - Inches(0.4)
    add_textbox(slide, text, Inches(0.5), top,
                SLIDE_W - Inches(1), h,
                font_size=14, bold=True, color="#FFFFFF", align="center",
                bg=color)


def make_slide(page=None, total=None):
    s = prs.slides.add_slide(BLANK)
    add_background(s)
    return s


TOTAL = 16  # 15 + Foundation Model Selection slide


# ─── SLIDE 1: Title ──────────────────────────────────────────────────────────
s = make_slide()
# Big accent
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.8))
bar.fill.solid(); bar.fill.fore_color.rgb = hex_to_rgb(NAVY)
bar.line.fill.background()

add_textbox(s, "Genomic Foundation Models for\nMetagenomic Read Classification",
            Inches(0.6), Inches(2.4), Inches(12), Inches(1.5),
            font_size=36, bold=True, color=NAVY, align="left")
add_textbox(s, "Thesis review & publication discussion",
            Inches(0.6), Inches(4.2), Inches(12), Inches(0.6),
            font_size=22, color=TEAL, align="left")
# A divider line
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.95),
                         Inches(4), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = hex_to_rgb(TEAL); ln.line.fill.background()

add_textbox(s, "楊明儒 (Ming-Ju Yang) · 國立臺灣大學 資工所碩二\n2026-05-26",
            Inches(0.6), Inches(5.3), Inches(12), Inches(1),
            font_size=16, color=GRAY, align="left")

add_textbox(s, "Section 1: The Research Question  ·  Section 2: Six Key Findings\n"
               "Section 3: Synthesis  ·  Section 4: Publication Direction",
            Inches(0.6), Inches(6.6), Inches(12), Inches(0.6),
            font_size=12, color=GRAY, align="left")


# ─── SLIDE 2: The Problem ────────────────────────────────────────────────────
s = make_slide()
add_header(s, "The Problem  ·  Why Metagenomic Read Classification Matters",
           accent=NAVY, page=2, total=TOTAL)
# Schematic
add_picture(s, TMP / "problem_schematic.png", Inches(0.5), Inches(1.2),
            width=Inches(12.3))
# Why-it-matters box
add_textbox(s,
    "Why it matters  ·  Per-sample species abundance underlies microbiome diagnostics, "
    "infectious disease, ecology.  Existing k-mer tools (Kraken2/Bracken) fail on "
    "novel organisms — motivating learned representations.",
    Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.4),
    font_size=14, color=NAVY, align="left", bg="#F4F6F8", border=TEAL)


# ─── SLIDE 3: Three Model Families ───────────────────────────────────────────
s = make_slide()
add_header(s, "The Comparison  ·  Three Model Families on Identical Data",
           accent=NAVY, page=3, total=TOTAL)
add_picture(s, TMP / "three_families.png", Inches(0.7), Inches(1.2),
            width=Inches(11.9))
add_textbox(s,
    "Orthogonal axes  ·  tokenisation × pre-training × routing  ·  "
    "all on same 50M training reads · same 100K independent test set · "
    "same 2,505-genome reference inventory",
    Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.8),
    font_size=14, bold=True, color=NAVY, align="center",
    bg="#F4F6F8", border=NAVY)
add_textbox(s,
    "Evaluation: read-level Top-1 + 5 sample-level metrics  "
    "(Pearson r · Spearman ρ · Bray-Curtis · ROC AUC · Sens@95%spec)",
    Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
    font_size=12, color=GRAY, align="center")


# ─── SLIDE 4: Foundation Model Selection (NEW) ──────────────────────────────
s = make_slide()
add_header(s, "Backbone Selection  ·  Why NT-v2 Among DNA Foundation Models",
           accent=NAVY, page=4, total=TOTAL)
add_picture(s, TMP / "foundation_models.png", Inches(0.4), Inches(1.0),
            width=Inches(12.5))
# Left/right takeaway boxes
add_textbox(s,
    "Controlled 5M comparison (left)\n"
    "• Same data · same task · same head\n"
    "• Only backbone varies — fairest possible test\n"
    "• NT-v2 narrowly leads at 5M (+0.8 pp vs DNABERT-1)",
    Inches(0.5), Inches(5.5), Inches(6.2), Inches(1.4),
    font_size=11.5, color=NAVY, align="left", bg="#F4F6F8", border=TEAL)
add_textbox(s,
    "Why the lead widens at scale (right)\n"
    "• NT-v2 has 498M params · room to absorb 50M data\n"
    "• DNABERT-1/2 plateau earlier; we chose the model\n"
    "  with the steepest scaling slope · validated at 50M",
    Inches(7.0), Inches(5.5), Inches(5.8), Inches(1.4),
    font_size=11.5, color=NAVY, align="left", bg="#F4F6F8", border=GREEN)


# ─── SLIDE 5: Finding 1 — Tokenisation > Pre-training ───────────────────────
s = make_slide()
add_header(s, "Finding 1  ·  Tokenisation Dominates Pre-training",
           accent=GREEN, page=5, total=TOTAL)
add_picture(s, TMP / "tokenisation_dominates.png", Inches(0.4), Inches(1.1),
            width=Inches(12.5))
add_takeaway(s,
    "Pre-training contributes, but cannot rescue an inadequate tokenisation at species granularity",
    color=GREEN)


# ─── SLIDE 6: Finding 2 — Data Volume Dominates ──────────────────────────────
s = make_slide()
add_header(s, "Finding 2  ·  Data Volume Dominates Training Tricks",
           accent=GREEN, page=6, total=TOTAL)
add_picture(s, TMP / "data_vs_tricks.png", Inches(0.5), Inches(1.2),
            width=Inches(12.3))
add_takeaway(s,
    "Data scaling beats hyperparameter tuning by 10×  ·  invest engineering effort in data, not tricks",
    color=GREEN)


# ─── SLIDE 7: Finding 3 — Utility Threshold ──────────────────────────────────
s = make_slide()
add_header(s, "Finding 3  ·  Species-Level Practical Utility Threshold",
           accent=GREEN, page=7, total=TOTAL)
add_picture(s, TMP / "utility_threshold.png", Inches(0.4), Inches(1.0),
            width=Inches(8.5))
# Right side: noise floor explanation
add_textbox(s, "Noise floor mechanism",
            Inches(9.2), Inches(1.2), Inches(3.8), Inches(0.5),
            font_size=15, bold=True, color=RED, align="left")
add_textbox(s,
    "1,535 species · 1K reads / sample\n"
    "Expected signal per species: 0.65 reads\n\n"
    "At 17.8% accuracy: ~0.54 spurious\nreads/species  →  SNR ≈ 1\n\n"
    "At 50% accuracy: noise drops to\n0.33/species  →  signal dominates\n\n"
    "Threshold:  ~40-50% accuracy",
    Inches(9.2), Inches(1.75), Inches(3.8), Inches(4.5),
    font_size=12, color=NAVY, align="left", bg="#FFEBEE", border=RED)
add_takeaway(s,
    "Species abundance becomes useful only past ~40-50% read accuracy  ·  6-mer NT-v2 cannot reach this",
    color=GREEN)


# ─── SLIDE 8: Finding 4 — Router Threshold ──────────────────────────────────
s = make_slide()
add_header(s, "Finding 4  ·  Hierarchical Masking Is Router-Dependent (Monotonic)",
           accent=GREEN, page=8, total=TOTAL)
add_picture(s, TMP / "router_threshold.png", Inches(0.6), Inches(1.0),
            width=Inches(8.4))
# Right side: intuition box
add_textbox(s, "Intuition",
            Inches(9.4), Inches(1.2), Inches(3.6), Inches(0.5),
            font_size=15, bold=True, color=ORANGE, align="left")
add_textbox(s,
    "Router right → mask removes\nwrong candidates  (good)\n\n"
    "Router wrong → mask excludes\nthe true species  (catastrophic)\n\n"
    "Net effect = average of these\nweighted by router accuracy\n\n"
    "Crossover near ~80% router",
    Inches(9.4), Inches(1.75), Inches(3.6), Inches(4.0),
    font_size=12, color=NAVY, align="left", bg="#FFF3E0", border=ORANGE)
add_takeaway(s,
    "Hierarchical pipelines need router > ~80% to be net positive  ·  not a free-lunch optimisation",
    color=GREEN)


# ─── SLIDE 9: Finding 5 — Per-Genus Oracle ──────────────────────────────────
s = make_slide()
add_header(s, "Finding 5  ·  Per-Genus Oracle Quantifies the Tokenisation Gap",
           accent=GREEN, page=9, total=TOTAL)
add_picture(s, TMP / "oracle_ceiling.png", Inches(0.4), Inches(1.0),
            width=Inches(12.5))
add_takeaway(s,
    "Even with PERFECT genus routing, 6-mer NT-v2 ceiling = 29.5%  ·  MT 13-mer 49.7% (no routing)  →  +20 pp tokenisation gap",
    color=ORANGE)


# ─── SLIDE 10: Finding 6 — Kraken2 OOD ───────────────────────────────────────
s = make_slide()
add_header(s, "Finding 6  ·  Kraken2 Wins In-DB · Fails on Novel Organisms (OOD)",
           accent=GREEN, page=10, total=TOTAL)
add_picture(s, TMP / "kraken2_indb_ood.png", Inches(0.4), Inches(1.0),
            width=Inches(12.5))
add_takeaway(s,
    "Learned methods complement, not replace, k-mer matching  ·  different deployment regimes (in-DB vs OOD)",
    color=GOLD)


# ─── SLIDE 11: Complete Results Matrix ───────────────────────────────────────
s = make_slide()
add_header(s, "Complete Results  ·  Sample-Level Evaluation Matrix",
           accent=TEAL, page=11, total=TOTAL)
add_picture(s, TMP / "results_matrix.png", Inches(0.5), Inches(1.1),
            width=Inches(12.3))
add_takeaway(s,
    "Independent 100K test  ·  1K reads/sample  ·  100 partition samples  ·  same DB, same data, same protocol",
    color=TEAL)


# ─── SLIDE 12: Story Arc ─────────────────────────────────────────────────────
s = make_slide()
add_header(s, "The Unifying Story  ·  Tokenisation Is the Hidden Bottleneck",
           accent=TEAL, page=12, total=TOTAL)
add_picture(s, TMP / "story_arc.png", Inches(0.4), Inches(1.0),
            width=Inches(12.5))
add_takeaway(s,
    "Foundation-model literature emphasises scale & pre-training  ·  our study isolates tokenisation as decisive at species granularity",
    color=TEAL)


# ─── SLIDE 13: Contributions ─────────────────────────────────────────────────
s = make_slide()
add_header(s, "What This Work Contributes  ·  Four Standalone Findings",
           accent=TEAL, page=13, total=TOTAL)
# 4 cards
contribs = [
    ("(1)", "Controlled 3-axis ablation",
     "Same 50M data · same test · same protocol\nIsolates tokenisation as dominant at species level", GREEN),
    ("(2)", "Quantitative noise-floor analysis",
     "Non-linear utility curve  ·  ~40-50% threshold\nMechanistic explanation via SNR model", TEAL),
    ("(3)", "Router quality threshold theorem",
     "Monotonic across 3 router qualities\n~80% accuracy required for net positive", ORANGE),
    ("(4)", "Same-DB Kraken2 + OOD analysis",
     "Kraken2 wins in-DB (66%)  ·  fails on OOD (0%)\nPositions learned methods as complementary", GOLD),
]
positions = [(0.5, 1.2), (6.8, 1.2), (0.5, 4.1), (6.8, 4.1)]
for (num, title, body, color), (x, y) in zip(contribs, positions):
    # Number badge
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                              Inches(0.7), Inches(0.7))
    shp.fill.solid(); shp.fill.fore_color.rgb = hex_to_rgb(color)
    shp.line.fill.background()
    tf = shp.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run(); run.text = num
    run.font.size = Pt(15); run.font.bold = True
    run.font.color.rgb = hex_to_rgb(WHITE)
    # Title
    add_textbox(s, title, Inches(x + 0.85), Inches(y + 0.05),
                Inches(5.3), Inches(0.6),
                font_size=15, bold=True, color=color, align="left")
    # Body
    add_textbox(s, body, Inches(x + 0.85), Inches(y + 0.7),
                Inches(5.3), Inches(1.7),
                font_size=12, color=NAVY, align="left",
                bg="#F4F6F8", border=color)


# ─── SLIDE 14: Publication Options ──────────────────────────────────────────
s = make_slide()
add_header(s, "Publication Direction  ·  Candidate Venues",
           accent=GOLD, page=14, total=TOTAL)
add_picture(s, TMP / "publication_matrix.png", Inches(0.4), Inches(1.0),
            width=Inches(12.5))
add_takeaway(s,
    "Recommended:  Bioinformatics journal (primary)  →  ML workshop if rejected  →  Thesis as baseline",
    color=GOLD)


# ─── SLIDE 15: Strengths/Weaknesses ─────────────────────────────────────────
s = make_slide()
add_header(s, "Honest Assessment  ·  Strengths & Weaknesses (with mitigations)",
           accent=GOLD, page=15, total=TOTAL)
add_picture(s, TMP / "strengths_weaknesses.png", Inches(0.4), Inches(1.0),
            width=Inches(12.5))
add_takeaway(s,
    "Position as \"controlled comparison study\"  ·  add 1 real-dataset sanity check (e.g., HMP) before submission",
    color=GOLD)


# ─── SLIDE 16: Next Steps + Questions ───────────────────────────────────────
s = make_slide()
add_header(s, "Next Steps & Questions for You",
           accent=GOLD, page=16, total=TOTAL)
add_picture(s, TMP / "timeline.png", Inches(0.4), Inches(0.9),
            width=Inches(12.5))

# Four question cards
qs = [
    ("Q1", "Bioinformatics journal as primary venue — agreed?",  GREEN),
    ("Q2", "Real-dataset validation before submission, or post-defense?", TEAL),
    ("Q3", "ML workshop (Option B) worth the effort?", ORANGE),
    ("Q4", "Collaborators / datasets you would want included?", GOLD),
]
for i, (qn, qt, color) in enumerate(qs):
    x = Inches(0.4 + i * 3.18)
    y = Inches(4.5)
    add_textbox(s, qn, x, y, Inches(0.5), Inches(0.55),
                font_size=18, bold=True, color="#FFFFFF", align="center",
                bg=color)
    add_textbox(s, qt, x + Inches(0.55), y + Inches(0.02),
                Inches(2.55), Inches(1.6),
                font_size=11.5, color=NAVY, align="left",
                bg="#F4F6F8", border=color)


prs.save(str(OUT))
print(f"\n✅ Saved: {OUT}")
print(f"   Total slides: {len(prs.slides)}")

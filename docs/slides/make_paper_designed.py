#!/usr/bin/env python3
"""
Paper content & structure walkthrough — redesigned visual system, ZH + EN.

Manuscript: "How should genomic foundation models be evaluated for
metagenomic taxonomic classification? A controlled benchmark of
pre-training, tokenization, and data scaling" (Briefings in Bioinformatics,
Problem-solving Protocol).

Design system: validated categorical palette (dataviz skill reference,
references/palette.md) — entities keep a FIXED color across every chart
(Kraken2 raw = red/critical, Kraken2+Bracken = aqua/good, NT-v2 = blue,
MT 13-mer = violet, MT 6-mer = orange). Dark-ink cover/divider slides,
light-surface content slides, card-based bullets with accent markers,
eyebrow + title + rule header pattern, consistent footer.

Produces two decks: paper_designed_zh.pptx, paper_designed_en.pptx.
"""

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from matplotlib.patches import FancyBboxPatch
import numpy as np
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DOCS = Path(__file__).parent

# ═══════════════════════════════════════════════════════════════════════
# DESIGN TOKENS (validated palette — dataviz skill references/palette.md)
# ═══════════════════════════════════════════════════════════════════════

INK        = "#0B0B0B"
INK_2      = "#52514E"   # secondary text
MUTED      = "#898781"   # captions / axis labels
SURFACE    = "#FCFCFB"   # card / chart surface
PAGE       = "#F7F7F5"   # slide background
HAIRLINE   = "#E1E0D9"   # borders / gridlines
DARK_BG    = "#0D1B2A"   # cover / divider background
DARK_INK2  = "#AEB8C4"   # muted text on dark bg

BLUE    = "#2A78D6"   # NT-v2
AQUA    = "#1BAF7A"   # Kraken2+Bracken ("good" / corrected)
YELLOW  = "#EDA100"
GREEN   = "#008300"
VIOLET  = "#4A3AA7"   # MT 13-mer
RED     = "#E34948"   # Kraken2 raw ("critical" / deflated)
MAGENTA = "#E87BA4"
ORANGE  = "#EB6834"   # MT 6-mer / secondary highlight

GOOD, WARNING, SERIOUS, CRITICAL = "#0CA30C", "#FAB219", "#EC835A", "#D03B3B"

# fixed entity -> color (never re-cycled across the deck)
C_KRAKEN_RAW   = RED
C_KRAKEN_BRACK = AQUA
C_NTV2         = BLUE
C_MT13         = VIOLET
C_MT6          = ORANGE
C_BASELINE     = MUTED

# section identity colors (roadmap / dividers)
SEC_COLOR = {"12": BLUE, "3": VIOLET, "4": ORANGE, "5": RED, "6": AQUA, "7": GREEN, "8": MUTED}

CJK_BOLD_PATH = "/work/ymj1123ntu/NTU-Thesis-Writing-Template/fonts/chinese/cwTeX_Hei.ttf"
CJK_REG_PATH  = "/work/ymj1123ntu/NTU-Thesis-Writing-Template/fonts/chinese/cwTeX_Yuan.ttf"
fp_bold_zh = fm.FontProperties(fname=CJK_BOLD_PATH)
fp_reg_zh  = fm.FontProperties(fname=CJK_REG_PATH)

FONT_ZH = "Microsoft JhengHei"
FONT_EN = "Calibri"

plt.rcParams.update({
    "axes.spines.top": False, "axes.spines.right": False,
    "axes.spines.left": False,
    "axes.edgecolor": HAIRLINE, "axes.linewidth": 1,
    "figure.facecolor": SURFACE, "axes.facecolor": SURFACE,
    "savefig.facecolor": SURFACE,
    "savefig.dpi": 200, "savefig.bbox": "tight", "savefig.pad_inches": 0.2,
    "text.color": INK, "axes.labelcolor": INK_2, "xtick.color": MUTED, "ytick.color": MUTED,
    "font.family": "DejaVu Sans",
})


def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def fprop(lang, bold=False):
    if lang == "zh":
        return fp_bold_zh if bold else fp_reg_zh
    return None  # default DejaVu Sans reads fine for English


def fontname(lang):
    return FONT_ZH if lang == "zh" else FONT_EN


# ═══════════════════════════════════════════════════════════════════════
# TEXT CONTENT (ZH / EN)
# ═══════════════════════════════════════════════════════════════════════

T = {
"zh": dict(
    cover_eyebrow="BRIEFINGS IN BIOINFORMATICS · PROBLEM-SOLVING PROTOCOL",
    cover_title="論文完整內容與架構",
    cover_sub="How should genomic foundation models be evaluated for\nmetagenomic taxonomic classification?",
    cover_sub2="A controlled benchmark of pre-training, tokenization, and data scaling",
    cover_thesis="150bp 短讀分類的效能上限由 tokenization 決定，不是 backbone pre-training",
    cover_footer="楊明儒 · 台大生醫電資所 · 2026-07-14",

    div_abstract_num="摘要", div_abstract_title="論文摘要", div_abstract_sub="Background → Methods → Results → Conclusion",
    div_walk_num="全文", div_walk_title="逐章導覽", div_walk_sub="§1 破題 → §8 誠實揭露限制",
    div_close_num="收尾", div_close_title="建議、限制與總結", div_close_sub="§7–§8 + Supplementary + 總結",

    s_position_eyebrow="論文定位", s_position_title="投稿現況與文章類型",
    s_position_bullets=[
        ("期刊：Briefings in Bioinformatics（BIB）", BLUE, 0),
        ("文章類型：Problem-solving Protocol（2,000–5,000 字），而非傳統 original research", MUTED, 1),
        ("敘事定位：不是「我提出一個新方法」，而是「回答一個方法論問題 — GFM 該如何被評估用於 metagenomic taxonomic classification」", VIOLET, 0),
        ("第一作者：楊明儒（通訊作者）；單位：台大生醫電資所；第二作者/ORCID/CRediT 待補", MUTED, 1),
        ("目前完整度：科學內容（摘要～限制）已完整一致；行政資訊（作者、funding、data DOI）仍待補齊", ORANGE, 0),
    ],

    s_bg_eyebrow="摘要 · BACKGROUND", s_bg_title="研究背景",
    s_bg_bullets=[
        "Genomic foundation models 在許多基因體任務上表現亮眼",
        "但對 150bp 短讀、訊號稀疏的 metagenomic taxonomic classification 效果不明",
        "不清楚是 pre-training、tokenization 還是資料規模在決定表現",
        "這是本文要系統性拆解的核心問題",
    ],

    s_methods_eyebrow="摘要 · METHODS", s_methods_title="研究方法",
    s_methods_bullets=[
        "以人類腸道基因體目錄模擬 reads，涵蓋 120 屬、1,535 種",
        "逐一改變四個因子：資料規模（500K–250M）、pre-training（498M pretrained vs random-init）、tokenization（6/12/13-mer）、評估層級",
        "評估三層級：read-level 準確度、sample-level 豐度、binary detection",
        "對照 alignment-free 分類器 Kraken2（含 Bracken 修正）",
    ],

    s_results_eyebrow="摘要 · RESULTS", s_results_title="核心結果",
    s_results_bullets=[
        "Overlapping 13-mer 決定效能上限：從零訓練達 87.4%，遠勝 pretrained 6-mer 的 67.1%",
        "6-mer 於 50M reads 後幾乎飽和（+0.22pp），13-mer 持續提升到 98.7%",
        "固定 tokenization 時，pre-training 貢獻 +13.2pp",
        "豐度估計：67% 準確度模型達 r=0.99；Kraken2 raw r=0.82，Bracken 修正後 r=0.997",
        "真實 mock community：所有方法崩落到 r≈0.4–0.6，沒有明顯贏家",
    ],

    s_keypoints_eyebrow="KEY POINTS", s_keypoints_title="五個關鍵發現",
    s_keypoints=[
        ("① 資料規模效益取決於 tokenizer", "50M→250M 對 6-mer 幾乎無增益(+0.22pp)，對 13-mer 卻有 +11pp", BLUE),
        ("② Tokenization、不是 pre-training，決定上限", "13-mer 從零訓練(87.4%)打敗 498M 參數 pretrained 6-mer(67.1%)", VIOLET),
        ("③ Pre-training 仍是獨立且重要的因子", "固定 tokenization 時仍貢獻 +13.2pp — 兩個因子彼此獨立、都重要", ORANGE),
        ("④ Read-level 與 sample-level 效用會分歧", "67% 準確度模型模擬 r=0.99，但排名無法轉移到真實 mock community", RED),
        ("⑤ 給出實務建議與未來方向", "主張未來 microbial foundation model 應以 overlapping long-k tokenization 預訓練", AQUA),
    ],

    s_roadmap_eyebrow="全文架構", s_roadmap_title="9 個章節如何串起一個論證",
    roadmap_labels=["§1\nIntroduction", "§2\n為何 GFM\n逆風", "§3\nBenchmark\n設計", "§4\n資料規模\n飽和",
                    "§5\nPretrain vs\nTokenization", "§6\nRead vs Sample\n+ 真實驗證", "§7\n實務\n建議", "§8\n限制與\n未來方向"],
    roadmap_caption="核心論點貫穿全文：Tokenization（overlapping 13-mer）— 而非 backbone pre-training — 決定 150bp 短讀分類的效能上限",

    s_challenge_eyebrow="§1–2 · 破題與挑戰", s_challenge_title="為什麼 GFM 對短讀分類是逆風局",
    s_challenge_bullets=[
        ("Alignment-free 工具（Kraken2/Centrifuge）在 in-database 又快又準，但遇到新物種就失效；GFM 在其他基因體任務上很強 — 能不能用在這裡？", MUTED, 0),
        ("① 短 reads 讓長距離建模優勢（GFM 的強項）沒有用武之地", BLUE, 1),
        ("② Tokenization 決定訊號存留：NT-v2 用 non-overlapping 6-mer（motif 常被切斷）vs Kraken2 用 overlapping 長 k-mer", BLUE, 1),
        ("③ 120–1,535 類、長尾分布的標籤空間 — 選什麼 metric 會改變「誰是贏家」", BLUE, 1),
    ],

    s_design_eyebrow="§3 · Benchmark 設計", s_design_title="設計與評估協定",
    design_labels=["腸道基因體目錄\n120 屬 / 1,535 種", "洩漏控制模擬 reads\n(ART, 150bp)\nread + genome 雙重切分",
                   "三軸逐一改變\n資料規模 / Pre-training /\nTokenization", "三層級評估\nRead-level / Sample-level /\nDetection"],
    design_kraken="對照 Kraken2\n(+Bracken)",
    design_caption="關鍵細節：258.67M reads 來源池只有 42.64M 唯一序列（6.07× 冗餘）；三個測試池：自然分布 / 較難 leftover / Kraken2 coverage-matched",

    s_scaling_eyebrow="§4 · 資料規模", s_scaling_title="資料規模與飽和效應",
    scaling_legend6="NT-v2 non-overlapping 6-mer（pretrained）", scaling_legend13="MT overlapping 13-mer（from scratch）",
    scaling_xlabel="訓練 reads 數量（log scale）", scaling_ylabel="Genus Top-1 準確度 (%)",
    scaling_ann6="50M→250M 只 +0.22pp\n（飽和）", scaling_ann13="50M→250M +11.3pp\n（持續提升）",
    scaling_caption="額外證據：6-mer 模型連自己的訓練集都 fit 不到 66% 以上（train≈val）— 代表性瓶頸，不是優化問題",

    s_decomp_eyebrow="§5 · 因果拆解", s_decomp_title="Pre-training vs Tokenization",
    decomp_labels=["Random-init\n6-mer", "Pretrained\nNT-v2 6-mer", "13-mer\nscratch"],
    decomp_title_l="固定 50M reads：兩個因子的獨立貢獻",
    decomp_ann_pre="+13.19pp\nPre-training", decomp_ann_tok="+20.35pp\nTokenization",
    decomp_labels2=["Unique-key\nlookup", "Majority\nvote", "Naive\nBayes", "Neural\n13-mer"],
    decomp_title_r="「13-mer 98.7% 是不是查表？」— 不是",
    decomp_caption="檢驗鏈：exact-match lookup 僅 0.72% → 神經網路加值 +13pp 於 naive Bayes 之上 → 結論：學到的是組成表徵，不是記憶",

    s_kraken_eyebrow="§6a · Sample-level 效用", s_kraken_title="與 Kraken2 的權衡",
    kraken_xlabel="豐度估計 Pearson r（in-database, coverage-matched）", kraken_ylabel="偵測敏感度 @95% 特異度 (%)",
    kraken_title="in-database 上「豐度 vs 偵測」的權衡（含 Bracken 修正）",
    kraken_ann="Bracken 修正後\n兩者兼優（右上角）",

    s_mock_eyebrow="§6b · 真實試煉", s_mock_title="ZymoBIOMICS Mock Community (D6331)",
    mock_legend_sim="模擬（in-DB）", mock_legend_real="真實 D6331（2 重複均值）",
    mock_ylabel="Pearson r（genus 豐度）", mock_title="模擬 → 真實：全面崩落（n=11 屬，CI 重疊、無明顯排名）",
    mock_panel_title="真實 D6331（2 重複）關鍵數字",
    mock_rows=[("樣本", "ZymoBIOMICS D6331，SRR33710519/518，各 3M reads，135bp"),
               ("In-set 屬", "11/21（81.5% reads），其餘含 Escherichia(14%) 不在訓練集內"),
               ("Kraken2 分類率", "真實僅 ~44%（模擬時 ~70%）"),
               ("偵測敏感度", "Bracken 55.6%  ·  NT-v2 83.3%  ·  MT13-mer 77.8%"),
               ("Bray-Curtis", "Kraken2+Bracken 最差 0.434（Veillonella 高估）"),
               ("共同失敗模式", "Roseburia 兩神經模型都嚴重低估；Clostridium 吸收多餘質量")],
    s_mock_caption="§6 收尾訊息：neural 模型在 in-database 豐度估計上並不比配置完整的 Kraken2+Bracken 更強；真正的價值在 read-level 準確度與 out-of-database 穩健性",

    s_reco_eyebrow="§7 · 實務建議", s_reco_title="協定、決策指南與建模者建議",
    s_reco_bullets=[
        ("Box 1：8 步驟可重用評估協定（每步都附上 rationale/pitfall）", INK, 0, True),
        ("先固定 read 長度 → read+genome 雙重切分 → 算唯一 reads → 一次只改一軸 → 三層級分開評估 → 同時用 raw 與修正後 baseline → 統一 RC-TTA 協定 → 用真實數據驗證並公開一切", MUTED, 1, False),
        ("Box 2 決策指南：低豐度偵測用 Kraken2；closed-set 讀分類選長 k 模型；只有 6-mer GFM 別指望資料量救得了；真實群落豐度目前沒有穩定贏家", INK, 0, False),
        ("給建模者的建議：tokenization 是槓桿最大的設計選擇 — 未來 microbial FM 應以 overlapping 12–13-mer 預訓練；工程挑戰是長k詞彙表意味著數十億參數的 embedding table", ORANGE, 0, False),
    ],

    s_limits_eyebrow="§8 · 主動揭露", s_limits_title="限制與未來方向",
    s_limits_bullets=[
        ("P0：Closed-set 評估 — 98.7%/r≈1 都是同一目錄；已證明是「學到組成」而非查表，但 out-of-genome generalization 仍是最重要待辦", CRITICAL, 0, True),
        ("模擬 reads 的侷限：ART 模擬可能對 alignment/k-mer 方法有利；真實 mock community 已顯示落差，但只有 1 個標準品/11 個屬", INK, 0, False),
        ("Baseline 廣度：只比較 Kraken2（raw + Bracken），建議未來加入 Centrifuge 等", INK, 0, False),
        ("單一 read 長度：模擬用 150bp、mock 用 135bp，其他長度未測試；RC-TTA 只用在 NT-v2（已揭露、不影響定性結論）", INK, 0, False),
    ],

    s_supp_eyebrow="Supplementary + 待辦", s_supp_title="附錄內容與投稿前待補清單",
    s_supp_bullets=[
        ("Supplementary（S1–S8）：資料切分 · 訓練設定 · 其他 backbone 比較 · 完整結果 · train-fit 上限 · 查表分析 · 計算成本 · 真實 mock 細節", INK, 0, True),
        ("✅ 已填：Conflicts of interest、Code availability（GitHub repo）、AI 使用揭露", GOOD, 0, False),
        ("⚠️ 待補：Funding 金額/grant number、Data DOI、第二作者 + ORCID + CRediT initials、Acknowledgments 名單", WARNING, 0, False),
        ("⚠️ 待補：Supplementary S2 的 optimizer/scheduler 名稱需再次確認", WARNING, 0, False),
    ],
    s_supp_caption="科學內容已完整一致；行政資訊是目前唯一的「非科學」缺口",

    s_summary_eyebrow="總結", s_summary_title="一張圖看完整論文",
    s_summary_bullets=[
        ("核心貢獻", "把「data / pre-training / tokenization」三個常被混在一起的因子，用嚴謹的 one-factor-at-a-time 設計拆開，找出 tokenization 才是決定性變數", BLUE),
        ("最誠實的部分", "主動加入 Bracken baseline 撤回過度宣稱的豐度優勢，又主動加入真實 mock community 揭露模擬-真實的落差", CRITICAL),
        ("最大待辦", "out-of-genome generalization 實驗（P0）+ 作者/funding/DOI 行政資訊", WARNING),
    ],
    s_summary_footer="下一步：先做 P0 實驗，再補齊行政資訊，準備投稿",

    footer_title="GFM Metagenomic Benchmark · BIB Problem-solving Protocol",
),

"en": dict(
    cover_eyebrow="BRIEFINGS IN BIOINFORMATICS · PROBLEM-SOLVING PROTOCOL",
    cover_title="Full Manuscript Walkthrough",
    cover_sub="How should genomic foundation models be evaluated for\nmetagenomic taxonomic classification?",
    cover_sub2="A controlled benchmark of pre-training, tokenization, and data scaling",
    cover_thesis="Tokenization — not backbone pre-training — sets the performance ceiling for 150bp read classification",
    cover_footer="Ming-Ju Yang · GIBEB, National Taiwan University · 2026-07-14",

    div_abstract_num="Abstract", div_abstract_title="The Abstract", div_abstract_sub="Background → Methods → Results → Conclusion",
    div_walk_num="Body", div_walk_title="Section-by-Section Walkthrough", div_walk_sub="§1 framing → §8 honest limitations",
    div_close_num="Closing", div_close_title="Recommendations, Limitations & Summary", div_close_sub="§7–§8 + Supplementary + Summary",

    s_position_eyebrow="Manuscript status", s_position_title="Submission Target & Article Type",
    s_position_bullets=[
        ("Journal: Briefings in Bioinformatics (BIB)", BLUE, 0),
        ("Article type: Problem-solving Protocol (2,000–5,000 words), not a traditional original-research paper", MUTED, 1),
        ("Framing: not \"we propose a new method\" but \"we answer a methodological question — how should GFMs be evaluated for metagenomic taxonomic classification\"", VIOLET, 0),
        ("Corresponding author: Ming-Ju Yang; affiliation: GIBEB, NTU; second author / ORCID / CRediT still pending", MUTED, 1),
        ("Current completeness: the science (abstract through limitations) is complete and internally consistent; administrative info (authors, funding, data DOI) is still outstanding", ORANGE, 0),
    ],

    s_bg_eyebrow="ABSTRACT · BACKGROUND", s_bg_title="Motivation",
    s_bg_bullets=[
        "Genomic foundation models (GFMs) show promise across many genomic tasks",
        "But their value for 150bp, signal-sparse metagenomic taxonomic classification is poorly characterized",
        "It's unclear whether pre-training, tokenization, or data scale actually governs performance",
        "This is the central question the manuscript systematically decomposes",
    ],

    s_methods_eyebrow="ABSTRACT · METHODS", s_methods_title="Approach",
    s_methods_bullets=[
        "Reads simulated from a unified human-gut genome catalogue: 120 genera, 1,535 species",
        "Four axes varied one at a time: data scale (500K–250M reads), pre-training (498M pretrained vs. random-init), tokenization (6/12/13-mer), evaluation level",
        "Three evaluation levels: read-level accuracy, sample-level abundance, binary detection",
        "Benchmarked against the alignment-free classifier Kraken2 (raw and Bracken-corrected)",
    ],

    s_results_eyebrow="ABSTRACT · RESULTS", s_results_title="Headline Findings",
    s_results_bullets=[
        "Overlapping 13-mer tokenization sets the ceiling: 87.4% from scratch vs. 67.1% for a pretrained 6-mer model",
        "The 6-mer model nearly saturates past 50M reads (+0.22pp); the 13-mer model keeps climbing to 98.7%",
        "With tokenization fixed, pre-training still contributes +13.2pp",
        "Abundance: a 67%-accurate model reaches r=0.99; raw Kraken2 is r=0.82, Bracken-corrected reaches r=0.997",
        "On a real mock community, every method collapses to r≈0.4–0.6 — no clear winner",
    ],

    s_keypoints_eyebrow="KEY POINTS", s_keypoints_title="Five Headline Findings",
    s_keypoints=[
        ("① Data scaling benefit is tokenizer-dependent", "50M→250M yields ~0 gain for 6-mer (+0.22pp) but +11pp for 13-mer", BLUE),
        ("② Tokenization, not pre-training, sets the ceiling", "13-mer from scratch (87.4%) beats a 498M-parameter pretrained 6-mer model (67.1%)", VIOLET),
        ("③ Pre-training is still an independent, real factor", "Contributes +13.2pp at fixed tokenization — the two factors are separable and both matter", ORANGE),
        ("④ Read-level and sample-level utility diverge", "67%-accurate model reaches r=0.99 simulated, but that ranking doesn't transfer to a real mock community", RED),
        ("⑤ Practical, use-case-specific recommendations", "Future microbial foundation models should be pretrained with overlapping long-k tokenization", AQUA),
    ],

    s_roadmap_eyebrow="Manuscript structure", s_roadmap_title="How Nine Sections Build One Argument",
    roadmap_labels=["§1\nIntroduction", "§2\nWhy GFMs\nStruggle", "§3\nBenchmark\nDesign", "§4\nData Scaling\n& Saturation",
                    "§5\nPretrain vs\nTokenization", "§6\nRead vs Sample\n+ Real Data", "§7\nRecommen-\ndations", "§8\nLimitations\n& Future Work"],
    roadmap_caption="The through-line: tokenization (overlapping 13-mer) — not backbone pre-training — sets the ceiling for 150bp read classification",

    s_challenge_eyebrow="§1–2 · Framing", s_challenge_title="Why Short-Read Classification Is an Uphill Battle for GFMs",
    s_challenge_bullets=[
        ("Alignment-free tools (Kraken2/Centrifuge) are fast and accurate in-database but fail on novel organisms; GFMs excel elsewhere — can they help here?", MUTED, 0),
        ("① Short reads make long-range context modeling — GFMs' main strength — irrelevant", BLUE, 1),
        ("② Tokenization determines signal retention: NT-v2's non-overlapping 6-mers split motifs vs. Kraken2's overlapping long k-mers", BLUE, 1),
        ("③ A 120–1,535-way, long-tailed label space — the metric you pick changes who \"wins\"", BLUE, 1),
    ],

    s_design_eyebrow="§3 · Benchmark Design", s_design_title="Design and Evaluation Protocol",
    design_labels=["Gut genome\ncatalogue\n120 genera / 1,535 species", "Leakage-controlled\nsimulated reads\n(ART, 150bp)\nread + genome split",
                   "Vary one axis\nat a time\nData scale / Pre-training /\nTokenization", "Three evaluation\nlevels\nRead-level / Sample-level /\nDetection"],
    design_kraken="Benchmarked vs.\nKraken2 (+Bracken)",
    design_caption="Key detail: the 258.67M-read source pool has only 42.64M unique reads (6.07× redundancy); three test pools: natural-distribution / harder leftover / Kraken2 coverage-matched",

    s_scaling_eyebrow="§4 · Data Scale", s_scaling_title="Data Scaling and the Saturation Effect",
    scaling_legend6="NT-v2 non-overlapping 6-mer (pretrained)", scaling_legend13="MT overlapping 13-mer (from scratch)",
    scaling_xlabel="Training reads (log scale)", scaling_ylabel="Genus Top-1 accuracy (%)",
    scaling_ann6="50M→250M: only +0.22pp\n(saturated)", scaling_ann13="50M→250M: +11.3pp\n(keeps climbing)",
    scaling_caption="Further evidence: the 6-mer model can't even fit its own training set past ~66% (train≈val) — a representational, not optimization, ceiling",

    s_decomp_eyebrow="§5 · Causal Decomposition", s_decomp_title="Pre-training vs. Tokenization",
    decomp_labels=["Random-init\n6-mer", "Pretrained\nNT-v2 6-mer", "13-mer\nscratch"],
    decomp_title_l="At fixed 50M reads: two independent contributions",
    decomp_ann_pre="+13.19pp\nPre-training", decomp_ann_tok="+20.35pp\nTokenization",
    decomp_labels2=["Unique-key\nlookup", "Majority\nvote", "Naive\nBayes", "Neural\n13-mer"],
    decomp_title_r="\"Is 13-mer's 98.7% just lookup?\" — No",
    decomp_caption="The test chain: exact-match lookup gets only 0.72% → the neural net adds +13pp over naive Bayes → conclusion: it learns composition, not memorization",

    s_kraken_eyebrow="§6a · Sample-Level Utility", s_kraken_title="The Trade-off vs. Kraken2",
    kraken_xlabel="Abundance Pearson r (in-database, coverage-matched)", kraken_ylabel="Detection sensitivity @95% specificity (%)",
    kraken_title="Abundance vs. detection trade-off in-database (with Bracken correction)",
    kraken_ann="Bracken correction moves\nKraken2 to best-at-both\n(top right)",

    s_mock_eyebrow="§6b · Real-World Trial", s_mock_title="ZymoBIOMICS Mock Community (D6331)",
    mock_legend_sim="Simulated (in-DB)", mock_legend_real="Real D6331 (mean of 2 replicates)",
    mock_ylabel="Pearson r (genus abundance)", mock_title="Simulated → real: total collapse (n=11 genera, overlapping CIs, no clear ranking)",
    mock_panel_title="Key numbers on real D6331 (2 replicates)",
    mock_rows=[("Sample", "ZymoBIOMICS D6331, SRR33710519/518, 3M reads each, 135bp"),
               ("In-set genera", "11/21 (81.5% of reads); rest incl. Escherichia (14%) absent from training genera"),
               ("Kraken2 classification rate", "only ~44% real (vs. ~70% simulated)"),
               ("Detection sensitivity", "Bracken 55.6%  ·  NT-v2 83.3%  ·  MT 13-mer 77.8%"),
               ("Bray-Curtis", "worst for Kraken2+Bracken at 0.434 (Veillonella over-estimated)"),
               ("Shared failure modes", "both neural models badly under-call Roseburia; Clostridium absorbs excess mass")],
    s_mock_caption="§6's closing message: neural models hold no sample-level abundance advantage over a properly configured classical pipeline in-database; their real value is read-level accuracy and out-of-database robustness",

    s_reco_eyebrow="§7 · Practical Recommendations", s_reco_title="Protocol, Decision Guide & Advice for Model Builders",
    s_reco_bullets=[
        ("Box 1: an 8-step reusable evaluation protocol (each step carries its own rationale/pitfall)", INK, 0, True),
        ("Fix read length first → split by read AND genome → count unique reads → change one axis at a time → evaluate 3 levels separately → use both raw and corrected classical baselines → unify the RC-TTA protocol → validate on real data and release everything", MUTED, 1, False),
        ("Box 2 decision guide: use Kraken2 for low-abundance detection; favor a long-k model for closed-set read classification; don't expect data scale to save a 6-mer-only GFM; no stable winner yet for real-community abundance", INK, 0, False),
        ("For model builders: tokenization is the highest-leverage design choice — future microbial FMs should pretrain with overlapping 12–13-mers; the engineering challenge is that a long-k vocabulary implies a multi-billion-parameter embedding table", ORANGE, 0, False),
    ],

    s_limits_eyebrow="§8 · Disclosed Proactively", s_limits_title="Limitations and Future Directions",
    s_limits_bullets=[
        ("P0: closed-set evaluation — the 98.7%/r≈1 headline numbers are all in-catalogue; shown to be \"learned composition,\" not lookup, but out-of-genome generalization remains the most important open item", CRITICAL, 0, True),
        ("Simulated-read limits: ART simulation may favor alignment/k-mer methods; the real mock-community probe already shows the gap, but it's still one standard / 11 genera", INK, 0, False),
        ("Baseline breadth: only Kraken2 (raw + Bracken) compared; recommend adding Centrifuge and others", INK, 0, False),
        ("Single read length: simulation uses 150bp, mock uses 135bp, other lengths untested; RC-TTA applied only to NT-v2 (disclosed, doesn't affect qualitative conclusions)", INK, 0, False),
    ],

    s_supp_eyebrow="Supplementary + To-do", s_supp_title="Supplementary Contents & Pre-Submission Checklist",
    s_supp_bullets=[
        ("Supplementary (S1–S8): dataset/splits · training configs · other backbone comparisons · full results · train-fit ceiling · lookup analysis · compute cost · real mock-community detail", INK, 0, True),
        ("✅ Filled in: conflicts of interest, code availability (GitHub repo), AI-use disclosure", GOOD, 0, False),
        ("⚠️ Still pending: funding amount/grant number, data DOI, second author + ORCID + CRediT initials, acknowledgments list", WARNING, 0, False),
        ("⚠️ Still pending: re-confirm the optimizer/scheduler names in Supplementary S2", WARNING, 0, False),
    ],
    s_supp_caption="The science is complete and consistent; administrative info is the only remaining \"non-scientific\" gap",

    s_summary_eyebrow="Summary", s_summary_title="The Whole Paper on One Slide",
    s_summary_bullets=[
        ("Core contribution", "Separates three usually-entangled factors — data, pre-training, tokenization — with a rigorous one-factor-at-a-time design, and finds tokenization is the decisive variable", BLUE),
        ("Most honest part", "Proactively added the Bracken baseline to retract an over-claimed abundance advantage, and added a real mock community that exposes the simulated-to-real gap", CRITICAL),
        ("Biggest to-do", "The out-of-genome generalization experiment (P0) + author/funding/DOI administrative info", WARNING),
    ],
    s_summary_footer="Next: run the P0 experiment, then complete the administrative info, then submit",

    footer_title="GFM Metagenomic Benchmark · BIB Problem-solving Protocol",
),
}

# shared numeric data (language-independent)
SCALING_X6, SCALING_Y6 = [0.5, 5, 50, 250], [55.29, 63.05, 67.07, 67.29]
SCALING_X13, SCALING_Y13 = [50, 250], [87.42, 98.7]
DECOMP_VALS = [53.88, 67.07, 87.42]
DECOMP_VALS2 = [0.72, 35.3, 74.9, 87.42]
KRAKEN_POINTS = [("Kraken2 (raw)", 0.823, 93.5, C_KRAKEN_RAW), ("Kraken2+Bracken", 0.997, 93.5, C_KRAKEN_BRACK),
                 ("NT-v2 6-mer", 0.992, 17.5, C_NTV2), ("MT 13-mer", 0.999, 66.5, C_MT13)]
MOCK_METHODS = ["Kraken2+Bracken", "MT 13-mer", "NT-v2 6-mer"]
MOCK_SIM = [0.997, 0.999, 0.992]
MOCK_REAL = [0.578, 0.538, 0.412]
MOCK_COLORS = [C_KRAKEN_BRACK, C_MT13, C_NTV2]


# ═══════════════════════════════════════════════════════════════════════
# FIGURES
# ═══════════════════════════════════════════════════════════════════════

def fig_roadmap(lang, out):
    L = T[lang]
    fig, ax = plt.subplots(figsize=(13, 5.3))
    ax.set_xlim(0, 13); ax.set_ylim(0, 6); ax.axis("off")
    colors = [SEC_COLOR["12"], SEC_COLOR["12"], SEC_COLOR["3"], SEC_COLOR["4"],
              SEC_COLOR["5"], SEC_COLOR["6"], SEC_COLOR["7"], SEC_COLOR["8"]]
    labels = L["roadmap_labels"]
    n = len(labels)
    xs = np.linspace(0.95, 12.05, n)
    y = 3.3; w, h = 1.3, 1.65
    for x, label, color in zip(xs, labels, colors):
        box = FancyBboxPatch((x - w/2, y - h/2), w, h, boxstyle="round,pad=0.08",
                              facecolor=color, edgecolor="none", alpha=0.94)
        ax.add_patch(box)
        ax.text(x, y, label, ha="center", va="center", fontsize=10, color="white",
                fontproperties=fprop(lang, bold=True), linespacing=1.4)
    for i in range(n - 1):
        ax.annotate("", xy=(xs[i+1] - w/2 - 0.03, y), xytext=(xs[i] + w/2 + 0.03, y),
                    arrowprops=dict(arrowstyle="-|>", color=MUTED, lw=1.4))
    ax.text(6.5, 0.55, L["roadmap_caption"], ha="center", fontsize=11.5, color=INK_2,
            fontproperties=fprop(lang))
    plt.savefig(out / "roadmap.png"); plt.close()


def fig_design_pipeline(lang, out):
    L = T[lang]
    fig, ax = plt.subplots(figsize=(13, 4.7))
    ax.set_xlim(0, 13); ax.set_ylim(0, 6); ax.axis("off")
    colors = [BLUE, VIOLET, ORANGE, AQUA]
    xs = [1.5, 4.7, 7.9, 11.1]
    w, h = 2.3, 2.3
    y = 3.6
    for x, label, color in zip(xs, L["design_labels"], colors):
        box = FancyBboxPatch((x - w/2, y - h/2), w, h, boxstyle="round,pad=0.1",
                              facecolor=color, edgecolor="none", alpha=0.94)
        ax.add_patch(box)
        ax.text(x, y, label, ha="center", va="center", fontsize=10, color="white",
                fontproperties=fprop(lang, bold=True), linespacing=1.4)
    for i in range(len(xs) - 1):
        x0 = xs[i] + w/2 + 0.05; x1 = xs[i+1] - w/2 - 0.05
        ax.annotate("", xy=(x1, y), xytext=(x0, y), arrowprops=dict(arrowstyle="-|>", color=MUTED, lw=1.6))
    box2 = FancyBboxPatch((9.7, 0.45), 2.6, 1.2, boxstyle="round,pad=0.1",
                          facecolor=RED, edgecolor="none", alpha=0.94)
    ax.add_patch(box2)
    ax.text(11.0, 1.05, L["design_kraken"], ha="center", va="center", fontsize=10, color="white",
            fontproperties=fprop(lang, bold=True), linespacing=1.4)
    ax.annotate("", xy=(11.0, 1.7), xytext=(11.0, 2.45), arrowprops=dict(arrowstyle="-|>", color=MUTED, lw=1.4))
    ax.text(6.5, 0.15, L["design_caption"], ha="center", fontsize=10, color=INK_2, fontproperties=fprop(lang))
    plt.savefig(out / "design_pipeline.png"); plt.close()


def fig_scaling(lang, out):
    L = T[lang]
    fig, ax = plt.subplots(figsize=(12.2, 5.2))
    ax.plot(SCALING_X6, SCALING_Y6, "o-", color=C_NTV2, lw=2.5, ms=9, label=L["scaling_legend6"])
    ax.plot(SCALING_X13, SCALING_Y13, "o-", color=C_MT13, lw=2.5, ms=9, label=L["scaling_legend13"])
    ax.set_xscale("log")
    ax.set_xticks([0.5, 5, 50, 250]); ax.set_xticklabels(["500K", "5M", "50M", "250M"], fontsize=11)
    ax.set_xlabel(L["scaling_xlabel"], fontsize=12, fontproperties=fprop(lang))
    ax.set_ylabel(L["scaling_ylabel"], fontsize=12, fontproperties=fprop(lang))
    ax.set_ylim(45, 105)
    ax.grid(axis="y", color=HAIRLINE, lw=1, zorder=0)
    ax.legend(fontsize=10.5, prop=fprop(lang), loc="upper left", frameon=False)
    ax.annotate(L["scaling_ann6"], xy=(250, 67.29), xytext=(60, 50), fontsize=10.5, color=C_NTV2,
                fontweight="bold", fontproperties=fprop(lang, bold=True),
                arrowprops=dict(arrowstyle="->", color=C_NTV2, lw=1.4))
    ax.annotate(L["scaling_ann13"], xy=(250, 98.7), xytext=(90, 90), fontsize=10.5, color=C_MT13,
                fontweight="bold", fontproperties=fprop(lang, bold=True),
                arrowprops=dict(arrowstyle="->", color=C_MT13, lw=1.4))
    plt.savefig(out / "scaling.png"); plt.close()


def fig_decomposition(lang, out):
    L = T[lang]
    fig, axes = plt.subplots(1, 2, figsize=(13, 5.0))

    ax = axes[0]
    colors = [C_BASELINE, C_NTV2, C_MT13]
    bars = ax.bar(L["decomp_labels"], DECOMP_VALS, color=colors, width=0.55, zorder=3)
    for b, v in zip(bars, DECOMP_VALS):
        ax.text(b.get_x() + b.get_width()/2, v + 3.5, f"{v:.2f}", ha="center", fontsize=11.5,
                fontweight="bold", color=INK, zorder=5,
                bbox=dict(facecolor=SURFACE, edgecolor="none", pad=1.5, alpha=0.9))
    ax.annotate("", xy=(1, 67.07 - 8), xytext=(0, 53.88 + 8), arrowprops=dict(arrowstyle="->", color=GOOD, lw=2))
    ax.text(0.5, 57.5, L["decomp_ann_pre"], ha="center", fontsize=10, color=GOOD, fontweight="bold", fontproperties=fprop(lang, bold=True))
    ax.annotate("", xy=(2, 87.42 - 8), xytext=(1, 67.07 + 8), arrowprops=dict(arrowstyle="->", color=CRITICAL, lw=2))
    ax.text(1.5, 74, L["decomp_ann_tok"], ha="center", fontsize=10, color=CRITICAL, fontweight="bold", fontproperties=fprop(lang, bold=True))
    ax.set_ylim(0, 110); ax.set_ylabel("Genus Top-1 (%)", fontsize=11)
    ax.set_title(L["decomp_title_l"], fontsize=12.5, fontweight="bold", color=INK, fontproperties=fprop(lang, bold=True), pad=12)
    ax.grid(axis="y", color=HAIRLINE, lw=1, zorder=0)
    plt.setp(ax.get_xticklabels(), fontproperties=fprop(lang))

    ax2 = axes[1]
    colors2 = [RED, ORANGE, YELLOW, VIOLET]
    bars2 = ax2.barh(L["decomp_labels2"], DECOMP_VALS2, color=colors2, height=0.55, zorder=3)
    for b, v in zip(bars2, DECOMP_VALS2):
        ax2.text(v + 1.5, b.get_y() + b.get_height()/2, f"{v:.1f}%", va="center", fontsize=11, fontweight="bold", color=INK)
    ax2.set_xlim(0, 105); ax2.set_xlabel("Genus Top-1 (%)", fontsize=11)
    ax2.set_title(L["decomp_title_r"], fontsize=12.5, fontweight="bold", color=INK, fontproperties=fprop(lang, bold=True), pad=12)
    ax2.grid(axis="x", color=HAIRLINE, lw=1, zorder=0)
    ax2.invert_yaxis()
    plt.setp(ax2.get_yticklabels(), fontproperties=fprop(lang))

    plt.savefig(out / "decomposition.png"); plt.close()


def fig_kraken(lang, out):
    L = T[lang]
    fig, ax = plt.subplots(figsize=(11, 5.5))
    for name, r, sens, color in KRAKEN_POINTS:
        ax.scatter(r, sens, s=800, color=color, alpha=0.9, edgecolor=SURFACE, linewidth=2.5, zorder=3)
        ax.annotate(name, (r, sens), xytext=(0, 22), textcoords="offset points",
                    ha="center", fontsize=11, fontweight="bold", color=color, fontproperties=fprop(lang, bold=True))
    ax.set_xlabel(L["kraken_xlabel"], fontsize=12, fontproperties=fprop(lang))
    ax.set_ylabel(L["kraken_ylabel"], fontsize=12, fontproperties=fprop(lang))
    ax.set_xlim(0.78, 1.03); ax.set_ylim(0, 105)
    ax.grid(color=HAIRLINE, lw=1, zorder=0)
    ax.annotate(L["kraken_ann"], xy=(0.997, 93.5), xytext=(0.87, 72), fontsize=10.5, color=GOOD,
                fontweight="bold", fontproperties=fprop(lang, bold=True),
                arrowprops=dict(arrowstyle="->", color=GOOD, lw=1.4))
    ax.set_title(L["kraken_title"], fontsize=13.5, fontweight="bold", color=INK,
                 fontproperties=fprop(lang, bold=True), pad=14)
    plt.savefig(out / "kraken.png"); plt.close()


def fig_mock(lang, out):
    L = T[lang]
    fig, axes = plt.subplots(1, 2, figsize=(13, 5.0))

    ax = axes[0]
    x = np.arange(len(MOCK_METHODS)); w = 0.32
    ax.bar(x - w/2, MOCK_SIM, width=w, color=[BLUE]*3, alpha=0.35, label=L["mock_legend_sim"], zorder=3)
    ax.bar(x + w/2, MOCK_REAL, width=w, color=MOCK_COLORS, label=L["mock_legend_real"], zorder=3)
    for xi, sv, rv in zip(x, MOCK_SIM, MOCK_REAL):
        ax.text(xi - w/2, sv + 0.03, f"{sv:.3f}", ha="center", fontsize=9, color=INK_2)
        ax.text(xi + w/2, rv + 0.03, f"{rv:.3f}", ha="center", fontsize=10, fontweight="bold", color=INK)
    ax.set_xticks(x); ax.set_xticklabels(MOCK_METHODS, fontsize=10.5, fontproperties=fprop(lang))
    ax.set_ylabel(L["mock_ylabel"], fontsize=11, fontproperties=fprop(lang))
    ax.set_ylim(0, 1.3)
    ax.legend(fontsize=10, prop=fprop(lang), loc="upper center", bbox_to_anchor=(0.5, 1.16), ncol=2, frameon=False)
    ax.set_title(L["mock_title"], fontsize=12, fontweight="bold", color=INK,
                 fontproperties=fprop(lang, bold=True), pad=40)
    ax.grid(axis="y", color=HAIRLINE, lw=1, zorder=0)

    ax2 = axes[1]
    ax2.axis("off"); ax2.set_xlim(0, 10); ax2.set_ylim(0, 10)
    ax2.text(5, 9.4, L["mock_panel_title"], fontsize=13, fontweight="bold", color=INK, ha="center",
              fontproperties=fprop(lang, bold=True))
    y0 = 8.2
    for name, val in L["mock_rows"]:
        ax2.text(0.2, y0, name, fontsize=11, fontweight="bold", color=BLUE, fontproperties=fprop(lang, bold=True))
        ax2.text(0.2, y0 - 0.5, f"  {val}", fontsize=9.5, color=INK_2, fontproperties=fprop(lang))
        y0 -= 1.35
    plt.savefig(out / "mock.png"); plt.close()


# ═══════════════════════════════════════════════════════════════════════
# PPTX HELPERS — design system
# ═══════════════════════════════════════════════════════════════════════

SW, SH = Inches(13.333), Inches(7.5)


def set_no_line(shape):
    shape.line.fill.background()


def fill_solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(color)


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    fill_solid(bg, color); set_no_line(bg)
    bg.shadow.inherit = False
    return bg


def txt(slide, text, left, top, width, height, size=14, color=INK, bold=False,
        align=PP_ALIGN.LEFT, font=FONT_ZH, italic=False, anchor=None, line_spacing=None,
        letter_spacing=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = rgb(color); r.font.name = font
    return tb


def eyebrow(slide, text, left, top, color, font):
    tb = txt(slide, text.upper(), left, top, 10, 0.35, size=12, color=color, bold=True, font=font)
    return tb


def accent_rule(slide, left, top, width=1.1, color=BLUE, thickness=0.045):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(thickness))
    fill_solid(r, color); set_no_line(r)
    r.shadow.inherit = False
    return r


def footer(slide, lang, page_no, total, color=MUTED):
    L = T[lang]
    txt(slide, L["footer_title"], 0.6, 7.12, 8, 0.3, size=9, color=color, font=fontname(lang))
    txt(slide, f"{page_no:02d} / {total:02d}", 12.2, 7.12, 0.6, 0.3, size=9, color=color,
        align=PP_ALIGN.RIGHT, font=fontname(lang))


def content_header(slide, lang, eyebrow_text, title_text, accent):
    add_bg(slide, PAGE)
    eyebrow(slide, eyebrow_text, 0.6, 0.55, accent, fontname(lang))
    txt(slide, title_text, 0.6, 0.86, 12.0, 0.85, size=27, color=INK, bold=True, font=fontname(lang))
    accent_rule(slide, 0.62, 1.62, width=1.0, color=accent)


def bullet_card(slide, items, left, top, width, height, lang, marker_default=BLUE, gap=0.16, base_size=15):
    """items: list of (text, color, level) or (text, color, level, bold)"""
    y = top
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    fill_solid(card, SURFACE)
    card.line.color.rgb = rgb(HAIRLINE); card.line.width = Pt(1)
    card.shadow.inherit = False
    pad = 0.28
    yy = top + pad
    for item in items:
        if len(item) == 4:
            text, color, level, bold = item
        else:
            text, color, level = item; bold = (level == 0)
        indent = 0.32 + level * 0.35
        marker_w = 0.1
        if level == 0:
            m = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.28), Inches(yy + 0.09),
                                        Inches(marker_w), Inches(marker_w))
            fill_solid(m, color); set_no_line(m)
            m.shadow.inherit = False
        size = base_size if level == 0 else base_size - 2
        tcolor = INK if (level == 0 and color != MUTED) else (color if level > 0 else INK)
        tb = txt(slide, text, left + indent, yy, width - indent - 0.3, 1.4,
                 size=size, color=(INK if level == 0 else INK_2), bold=(bold and level == 0),
                 font=fontname(lang), line_spacing=1.15)
        # estimate consumed height by rough char-per-line heuristic
        # CJK glyphs render roughly 2x the width of Latin glyphs at the same point size
        chars_per_inch_at_15pt = 4.3 if lang == "zh" else 9.5
        chars_per_line = max(8, int((width - indent - 0.3) * chars_per_inch_at_15pt * (15 / size)))
        n_lines = max(1, -(-len(text) // chars_per_line))
        line_h = (size / 72) * 1.28
        used = n_lines * line_h + 0.22
        yy += used + gap
    if yy > top + height:
        print(f"    !! bullet_card overflow: needs {yy - top:.2f}in, allocated {height:.2f}in "
              f"(first item: {items[0][0][:24]!r})")
    return card


def key_point_row(slide, num_text, title, body, color, left, top, width, height, lang):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    fill_solid(card, SURFACE)
    card.line.color.rgb = rgb(HAIRLINE); card.line.width = Pt(1)
    card.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.09), Inches(height))
    fill_solid(bar, color); set_no_line(bar); bar.shadow.inherit = False
    txt(slide, title, left + 0.35, top + 0.1, width - 0.6, 0.5, size=15.5, color=INK, bold=True, font=fontname(lang))
    txt(slide, body, left + 0.35, top + 0.55, width - 0.6, height - 0.65, size=12.5, color=INK_2, font=fontname(lang), line_spacing=1.1)


def add_image(slide, path, left, top, width, height):
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))


def divider_slide(prs, lang, num_label, title, sub, color):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, DARK_BG)
    accent_rule(sl, 0.9, 2.55, width=0.55, color=color, thickness=0.06)
    txt(sl, num_label.upper(), 0.9, 2.15, 8, 0.4, size=14, color=color, bold=True, font=fontname(lang))
    txt(sl, title, 0.85, 2.75, 11.5, 1.2, size=40, color="#FFFFFF", bold=True, font=fontname(lang))
    txt(sl, sub, 0.9, 3.85, 10.5, 0.6, size=15, color=DARK_INK2, font=fontname(lang))
    return sl


# ═══════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════

def build(lang):
    L = T[lang]
    out = DOCS / f"paper_designed_{lang}_figs"
    out.mkdir(exist_ok=True)
    print(f"[{lang}] generating figures...")
    fig_roadmap(lang, out)
    fig_design_pipeline(lang, out)
    fig_scaling(lang, out)
    fig_decomposition(lang, out)
    fig_kraken(lang, out)
    fig_mock(lang, out)

    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    TOTAL = 19
    pg = [0]

    def new_content(eyebrow_text, title_text, accent):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        content_header(sl, lang, eyebrow_text, title_text, accent)
        pg[0] += 1
        footer(sl, lang, pg[0], TOTAL)
        return sl

    # 1. Cover
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, DARK_BG)
    accent_rule(sl, 1.0, 1.55, width=0.7, color=BLUE, thickness=0.05)
    txt(sl, L["cover_eyebrow"], 1.0, 1.75, 11, 0.4, size=13, color=BLUE, bold=True, font=fontname(lang))
    txt(sl, L["cover_title"], 1.0, 2.2, 11.3, 1.3, size=42, color="#FFFFFF", bold=True, font=fontname(lang))
    txt(sl, L["cover_sub"], 1.0, 3.35, 11.3, 1.0, size=16, color="#C7D2E0", font=fontname(lang), line_spacing=1.2)
    txt(sl, L["cover_sub2"], 1.0, 4.15, 11.3, 0.5, size=12.5, color=DARK_INK2, italic=True, font=fontname(lang))
    thesis_card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.85), Inches(11.3), Inches(0.85))
    fill_solid(thesis_card, "#132A40"); thesis_card.line.color.rgb = rgb(BLUE); thesis_card.line.width = Pt(1)
    thesis_card.shadow.inherit = False
    txt(sl, L["cover_thesis"], 1.3, 5.0, 10.7, 0.6, size=14.5, color="#FFFFFF", bold=True, font=fontname(lang),
        anchor=MSO_ANCHOR.MIDDLE)
    txt(sl, L["cover_footer"], 1.0, 6.7, 11.3, 0.4, size=11, color=DARK_INK2, font=fontname(lang))
    pg[0] += 1
    footer(sl, lang, pg[0], TOTAL, color=DARK_INK2)

    # 2. Divider: Abstract
    divider_slide(prs, lang, L["div_abstract_num"], L["div_abstract_title"], L["div_abstract_sub"], BLUE)
    pg[0] += 1

    # 3-5. Abstract B/M/R
    sl = new_content(L["s_bg_eyebrow"], L["s_bg_title"], BLUE)
    bullet_card(sl, [(b, BLUE, 0) for b in L["s_bg_bullets"]], 0.6, 1.35, 12.1, 4.9, lang, base_size=17)

    sl = new_content(L["s_methods_eyebrow"], L["s_methods_title"], VIOLET)
    bullet_card(sl, [(b, VIOLET, 0) for b in L["s_methods_bullets"]], 0.6, 1.35, 12.1, 4.9, lang, base_size=16)

    sl = new_content(L["s_results_eyebrow"], L["s_results_title"], AQUA)
    bullet_card(sl, [(b, AQUA, 0) for b in L["s_results_bullets"]], 0.6, 1.35, 12.1, 4.9, lang, base_size=15.5)

    # 6. Key points
    sl = new_content(L["s_keypoints_eyebrow"], L["s_keypoints_title"], ORANGE)
    y = 1.3; hh = 0.98
    for title, body, color in L["s_keypoints"]:
        key_point_row(sl, "", title, body, color, 0.6, y, 12.1, hh, lang)
        y += hh + 0.1

    # 7. Roadmap divider-ish content slide
    divider_slide(prs, lang, L["div_walk_num"], L["div_walk_title"], L["div_walk_sub"], VIOLET)
    pg[0] += 1

    sl = new_content(L["s_roadmap_eyebrow"], L["s_roadmap_title"], INK)
    add_image(sl, out / "roadmap.png", 0.15, 1.35, 13.0, 5.15)

    # Challenge
    sl = new_content(L["s_challenge_eyebrow"], L["s_challenge_title"], BLUE)
    bullet_card(sl, L["s_challenge_bullets"], 0.6, 1.35, 12.1, 4.9, lang, base_size=15.5)

    # Design
    sl = new_content(L["s_design_eyebrow"], L["s_design_title"], VIOLET)
    add_image(sl, out / "design_pipeline.png", 0.15, 1.3, 13.0, 4.65)
    txt(sl, L["design_caption"], 0.6, 6.15, 12.1, 0.7, size=11.5, color=INK_2, font=fontname(lang))

    # Scaling
    sl = new_content(L["s_scaling_eyebrow"], L["s_scaling_title"], ORANGE)
    add_image(sl, out / "scaling.png", 1.0, 1.3, 11.3, 4.85)
    txt(sl, L["scaling_caption"], 1.0, 6.2, 11.3, 0.65, size=11.5, color=INK_2, font=fontname(lang))

    # Decomposition
    sl = new_content(L["s_decomp_eyebrow"], L["s_decomp_title"], RED)
    add_image(sl, out / "decomposition.png", 0.15, 1.3, 13.0, 4.65)
    txt(sl, L["decomp_caption"], 0.6, 6.05, 12.1, 0.75, size=11.5, color=INK_2, font=fontname(lang))

    # Kraken/Bracken
    sl = new_content(L["s_kraken_eyebrow"], L["s_kraken_title"], AQUA)
    add_image(sl, out / "kraken.png", 1.4, 1.3, 10.5, 5.4)

    # Mock
    sl = new_content(L["s_mock_eyebrow"], L["s_mock_title"], AQUA)
    add_image(sl, out / "mock.png", 0.15, 1.3, 13.0, 4.85)
    txt(sl, L["s_mock_caption"], 0.6, 6.2, 12.1, 0.65, size=11.5, color=VIOLET, bold=True, font=fontname(lang))

    # Closing divider
    divider_slide(prs, lang, L["div_close_num"], L["div_close_title"], L["div_close_sub"], GREEN)
    pg[0] += 1

    # Recommendations
    sl = new_content(L["s_reco_eyebrow"], L["s_reco_title"], GREEN)
    bullet_card(sl, L["s_reco_bullets"], 0.6, 1.35, 12.1, 4.9, lang, base_size=14.5)

    # Limitations
    sl = new_content(L["s_limits_eyebrow"], L["s_limits_title"], CRITICAL)
    bullet_card(sl, L["s_limits_bullets"], 0.6, 1.35, 12.1, 4.9, lang, base_size=14.5)

    # Supplementary
    sl = new_content(L["s_supp_eyebrow"], L["s_supp_title"], MUTED)
    bullet_card(sl, L["s_supp_bullets"], 0.6, 1.35, 12.1, 4.3, lang, base_size=14.5)
    txt(sl, L["s_supp_caption"], 0.6, 5.85, 12.1, 0.5, size=13, color=INK_2, bold=True, font=fontname(lang))

    # Summary
    sl = new_content(L["s_summary_eyebrow"], L["s_summary_title"], BLUE)
    y = 1.4; hh = 1.35
    for title, body, color in L["s_summary_bullets"]:
        key_point_row(sl, "", title, body, color, 0.6, y, 12.1, hh, lang)
        y += hh + 0.15
    footer_card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y + 0.05), Inches(12.1), Inches(0.6))
    fill_solid(footer_card, DARK_BG); set_no_line(footer_card); footer_card.shadow.inherit = False
    txt(sl, L["s_summary_footer"], 0.6, y + 0.05, 12.1, 0.6, size=14, color="#FFFFFF", bold=True,
        align=PP_ALIGN.CENTER, font=fontname(lang), anchor=MSO_ANCHOR.MIDDLE)

    out_pptx = DOCS / f"paper_designed_{lang}.pptx"
    prs.save(out_pptx)
    print(f"[{lang}] saved: {out_pptx} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    for lang in ["zh", "en"]:
        build(lang)

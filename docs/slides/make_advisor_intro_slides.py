#!/usr/bin/env python3
"""Generate 'Co-advisor Project Introduction' slide deck (full project overview).

Source script: docs/advisor_intro_script.md (35 slides).
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

TEMPLATE = Path("/work/ymj1123ntu/token_level_gfm_classifier/docs/0424.pptx")
OUT      = Path("/work/ymj1123ntu/token_level_gfm_classifier/docs/advisor_intro.pptx")

BLUE   = RGBColor(0x00, 0x70, 0xC0)
BLACK  = RGBColor(0x00, 0x00, 0x00)
GRAY   = RGBColor(0x40, 0x40, 0x40)
GREEN  = RGBColor(0x00, 0x70, 0x00)
RED    = RGBColor(0xC0, 0x00, 0x00)
ORANGE = RGBColor(0xE0, 0x70, 0x00)

prs = Presentation(TEMPLATE)
xml_slides = prs.slides._sldIdLst
for sld in list(xml_slides):
    xml_slides.remove(sld)

LAY_TITLE   = prs.slide_layouts[0]
LAY_CONTENT = prs.slide_layouts[1]
LAY_SECTION = prs.slide_layouts[19]


def add_section(text):
    slide = prs.slides.add_slide(LAY_SECTION)
    txBox = slide.shapes.add_textbox(
        Emu(590550), Emu(1530900), Emu(7962900), Emu(2081700))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(44)
    run.font.color.rgb = BLUE
    run.font.name = "Calibri"


def add_content(title, bullets):
    slide = prs.slides.add_slide(LAY_CONTENT)
    title_ph = slide.placeholders[0]
    title_ph.text = title
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(24)
    body_ph = slide.placeholders[1]
    tf = body_ph.text_frame
    tf.clear()
    first = True
    for level, text, bold, color in bullets:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.level = level
        run = para.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(17) if level == 0 else Pt(14)
        if color:
            run.font.color.rgb = color
    return slide


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(LAY_TITLE)
    slide.placeholders[0].text = title
    for para in slide.placeholders[0].text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.bold = True
    sub_ph = slide.placeholders[1]
    sub_ph.text = subtitle
    for para in sub_ph.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(18)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ──────────────────────────────────────────────────────────────────────────────
add_title_slide(
    "Token-Level GFM Classifier for Metagenomic Reads",
    "Co-advisor Project Introduction   |   楊明儒 (資工所碩二)   |   2026-04-23"
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — 大綱
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Outline",
    [
        (0, "1.  Project Background & Problem Definition", True, BLUE),
        (0, "2.  Background (GFM, Tokenization, LoRA, RC Symmetry)", True, BLUE),
        (0, "3.  System Architecture & Design Decisions", True, BLUE),
        (0, "4.  Experiment Timeline (Genus & Species)", True, BLUE),
        (0, "5.  Key Findings (Data, RC TTA, Hierarchical, Pre-training vs Tokenization)", True, BLUE),
        (0, "6.  Ongoing Experiments (DNABERT / DNABERT-2 Tokenization Comparison)", True, BLUE),
        (0, "7.  Challenges & Open Questions for Discussion", True, BLUE),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Section: 專案背景
# ──────────────────────────────────────────────────────────────────────────────
add_section("1.  Project Background & Problem Definition")

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — 問題定義
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Problem Definition",
    [
        (0, "Task:  Metagenomic taxonomic classification at read level", True, BLUE),
        (0, "Input:   Single 150 bp DNA short read  (simulated Illumina paired-end)", False, BLACK),
        (0, "Output:  Genus  (120 classes)  or  Species  (1,535 classes)", False, BLACK),
        (0, "Data:    HGR-UMGS  —  2,505 human-gut bacterial genomes", False, BLACK),
        (1, "ART Illumina simulator generates ~258 M reads", False, GRAY),
        (1, "Balanced subsamples: 5 M / 50 M", False, GRAY),
        (0, "Core difficulty:  classify a 150 bp fragment  (≈ 1 / 1,000,000 of a genome)", True, RED),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — 為什麼這個問題難
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Why This Task is Hard",
    [
        (0, "Short reads — low information density", True, BLUE),
        (1, "150 bp → 25 NT-v2 tokens vs 3 Mbp bacterial genome  (1 : 100,000)", False, BLACK),
        (0, "Very large class space", True, BLUE),
        (1, "1,535 species classes;  random baseline = 0.065 %", False, BLACK),
        (0, "Long-tail class distribution", True, BLUE),
        (1, "Head classes dominate loss; rare species have <100 reads even after balancing", False, BLACK),
        (0, "Phylogenetic proximity", True, BLUE),
        (1, "Closely related species share >100 bp identical stretches", False, BLACK),
        (0, "Compared to standard DNA benchmarks (promoter, splice site):  fine-grained over huge class space with very short context", True, ORANGE),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Section: 背景知識
# ──────────────────────────────────────────────────────────────────────────────
add_section("2.  Background")

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — GFMs
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Genomic Foundation Models (GFMs)",
    [
        (0, "Analogous to BERT / GPT in NLP — self-supervised pre-training on raw genomes", False, BLACK),
        (0, "Goal:  learn the 'language' of DNA  (motifs, codon bias, regulatory elements)", False, BLACK),
        (0, "Three representative models (k-mer based):", True, BLUE),
        (1, "DNABERT  (2021)        —  BERT-base,  86 M params,  overlapping 6-mer (stride=1)", False, BLACK),
        (1, "Nucleotide Transformer v2  (2024)  —  500 M params,  non-overlapping 6-mer (stride=6)", False, BLACK),
        (1, "DNABERT-2  (2023)      —  MosaicBERT, 117 M params,  BPE (learned tokenization)", False, BLACK),
        (0, "Downstream tasks historically focus on human regulatory elements", False, GRAY),
        (0, "Metagenomic classification is a relatively under-explored application domain", True, ORANGE),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — NT-v2
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Nucleotide Transformer v2  (our main backbone)",
    [
        (0, "Why we chose NT-v2:", True, BLUE),
        (1, "Multi-species pre-training:  850 species spanning bacteria / archaea / eukaryote", False, BLACK),
        (1, "Largest publicly available GFM  (500 M params;  others ~100 M)", False, BLACK),
        (1, "Actively maintained by InstaDeepAI / BioNTech", False, BLACK),
        (0, "Architecture:", True, BLUE),
        (1, "29 layers,  hidden dim = 1024,  498 M parameters", False, BLACK),
        (1, "Tokenization:  non-overlapping 6-mer,  stride = 6  →  150 bp = 25 tokens", False, BLACK),
        (1, "Vocab:  4,104 tokens  (4^6 = 4096 k-mers + special)", False, BLACK),
        (0, "Key design choice: fixed 6-mer is baked into pre-training — this becomes a discussion point later", True, RED),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — LoRA
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "LoRA — Low-Rank Adaptation",
    [
        (0, "Full fine-tune of 500 M params infeasible on a single GPU", False, BLACK),
        (0, "LoRA:  add low-rank delta matrices to attention Q / K / V  (r = 16)", True, BLUE),
        (1, "Backbone frozen;  only LoRA adapters + classification head update", False, BLACK),
        (1, "Trainable params:  5.54 M  ≈  1.11 %  of total", True, GREEN),
        (0, "Pros:  memory-efficient, fast training, avoids catastrophic forgetting", False, BLACK),
        (0, "Cons:  expressivity bounded by low-rank constraint", False, BLACK),
        (1, "If task / pre-train domain gap is large, full fine-tune may be needed", False, GRAY),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — RC symmetry
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Reverse-Complement (RC) Symmetry",
    [
        (0, "DNA double helix: ACGT → RC  (complement + reverse)", False, BLACK),
        (1, "Sequencer randomly reads either strand → both should be the SAME biological unit", False, BLACK),
        (0, "Transformers are NOT RC-equivariant — model sees s and RC(s) as distinct inputs", True, RED),
        (0, "Our three mitigations:", True, BLUE),
        (1, "RC augmentation  (train)  —  50 % flip probability", False, BLACK),
        (1, "RC TTA  (inference)  —  ŷ = argmax[ logits(s) + logits(RC(s)) ]", False, BLACK),
        (1, "RC consistency loss  (optional)  —  λ · KL( fwd ∥ RC )", False, BLACK),
        (0, "RC TTA is 'free lunch':  no retraining,  +1 ~ 2 pp at inference", True, GREEN),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Section: 系統架構
# ──────────────────────────────────────────────────────────────────────────────
add_section("3.  System Architecture")

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Pipeline
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "End-to-End Pipeline",
    [
        (0, "150 bp DNA read", True, BLACK),
        (1, "↓   6-mer tokenization (stride = 6)", False, GRAY),
        (0, "25 token IDs", True, BLACK),
        (1, "↓   NT-v2 backbone  (29 layers,  LoRA on Q / K / V)", False, GRAY),
        (0, "25 × 1024 token embeddings   ← FULL sequence preserved", True, BLUE),
        (1, "↓   4-head AttentionPool classification head", False, GRAY),
        (0, "1024-dim pooled representation", True, BLACK),
        (1, "↓   MLP  (2 layers,  512 hidden,  dropout 0.15)", False, GRAY),
        (0, "120 logits (genus)  /  1,535 logits (species)", True, BLACK),
        (0, "Key decision:  NO mean pooling before classifier  —  attention head learns its own weights", True, ORANGE),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — 為什麼不 mean pool
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Why not mean pooling?",
    [
        (0, "Phase-1 baseline  (frozen embedding):   backbone → mean pool → 1024-dim → MLP", False, BLACK),
        (1, "Problem:  averaging across 25 tokens washes out local patterns", False, GRAY),
        (0, "Token-level  AttentionPool head:", True, BLUE),
        (1, "Preserves full 25 × 1024 sequence", False, BLACK),
        (1, "4-head attention  ⇒  model learns which token is informative", False, BLACK),
        (1, "Similar spirit to CLS token but not tied to a specific position", False, BLACK),
        (0, "Result:  token-level beats mean-pool by ~ 8 pp genus accuracy  (same backbone, same training)", True, GREEN),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — 兩階段訓練
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Two-Phase Training with Differential LR",
    [
        (0, "Phase 1  (optional):  freeze backbone, train head only", True, BLUE),
        (1, "Avoids large gradients destabilizing pre-trained weights early on", False, BLACK),
        (0, "Phase 2:  unfreeze LoRA  +  head, train jointly", True, BLUE),
        (1, "Differential learning rate", True, BLACK),
        (1, "Head LR  =  5 e-4   (random-init  ⇒  needs to converge fast)", False, BLACK),
        (1, "Backbone LR (LoRA)  =  3 e-5   (pre-trained  ⇒  small adjustments)", False, BLACK),
        (0, "LR schedule:  cosine with 5 % warmup,  early-stopping patience = 5", False, GRAY),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — 資料設計
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Data Design",
    [
        (0, "Source:  258 M reads  from  2,505 genomes  (HGR-UMGS, human gut)", False, BLACK),
        (0, "Simulation:  ART Illumina simulator,  150 bp paired-end", False, BLACK),
        (0, "Balanced subsamples  (prevents head-class domination of loss):", True, BLUE),
        (1, "5 M reads   (v8, DNABERT / DNABERT-2 5M experiments)", False, BLACK),
        (1, "50 M reads  (v9, sp_v4 — thesis main scale)", False, BLACK),
        (0, "Train / Val split:  90 / 10,  seed fixed for reproducibility", False, GRAY),
        (0, "Training-time augmentation:  RC flip with 50 % probability", False, GRAY),
        (0, "120 genera,  1,535 species  (filtered to ≥ N reads per species)", False, GRAY),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — Section: 實驗歷程
# ──────────────────────────────────────────────────────────────────────────────
add_section("4.  Experiment Timeline")

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 17 — Genus 實驗演進
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Genus Experiment Progression  (v4 → v11)",
    [
        (0, "Version   Data      Key change                        Val Acc (RC TTA)", True, BLUE),
        (0, "v4        500 K     Basic architecture                 55 %", False, BLACK),
        (0, "v5        500 K     + Logit Adjustment (τ = 0.3)       ~ 55.5 %", False, BLACK),
        (0, "v6        500 K     + RC Consistency (λ = 0.1)         ~ 55.2 %", False, BLACK),
        (0, "v7        500 K     RC Consistency only                ~ 55.3 %", False, BLACK),
        (0, "v8        5 M       Basic, 10× data                    63.05 %", True, GREEN),
        (0, "v9        50 M      Basic, 100× data                   67.07 %", True, GREEN),
        (0, "v11       50 M      Shallow Transformer (ablation)     ~ 36 %", False, GRAY),
        (0, "Any single training trick:  ± 0.5 pp    |    500 K → 50 M:  +12 pp", True, RED),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 18 — 核心發現一：資料量 >> tricks
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Key Finding 1 — Data Scale ≫ Training Tricks",
    [
        (0, "Explored many tricks:  logit adjustment,  RC consistency,  weighted sampler,", False, BLACK),
        (0, "    class weighting,  label smoothing,  focal loss, …", False, BLACK),
        (1, "No single trick contributed more than +1 pp", True, RED),
        (0, "But data scaling alone:", True, BLUE),
        (1, "500 K → 5 M    :   + 8 pp", True, GREEN),
        (1, "5 M → 50 M     :   + 4 pp", True, GREEN),
        (1, "500 K → 50 M   :   + 12 pp", True, GREEN),
        (0, "Consistent with recent LLM scaling-law observations", False, GRAY),
        (0, "Practical implication:  ROI of data scaling ≫ loss / sampler / regularizer tuning", True, ORANGE),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 19 — RC TTA Free Lunch
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Key Finding 2 — RC TTA: Free Lunch at Inference",
    [
        (0, "Method:   ŷ  =  argmax [ logits(s)  +  logits(RC(s)) ]", False, BLACK),
        (0, "Cost:     inference time  ×  2   (two forward passes)", False, GRAY),
        (0, "v9 (50 M genus) results:", True, BLUE),
        (0, "Metric               Forward only    + RC TTA     Δ", True, BLACK),
        (0, "Top-1 accuracy        66.29 %         67.07 %     +0.78 pp", False, BLACK),
        (0, "Top-3                 83.85 %         84.40 %     +0.56 pp", False, BLACK),
        (0, "Top-10                95.02 %         95.25 %     +0.23 pp", False, BLACK),
        (0, "Why it works:  model's errors on forward vs RC are partially independent → strand ensemble", True, ORANGE),
        (0, "Theoretical justification in thesis §2.7 (strand-symmetry derivation)", False, GRAY),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 20 — Species 任務
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Species Task  (sp_v4, 1,535 Classes)",
    [
        (0, "sp_v4 (50 M, NT-v2 + LoRA):  val_acc ≈ 17.5 %", False, BLACK),
        (1, "Looks low but random baseline = 0.065 %  →  relative lift ≈ 270 ×", False, GRAY),
        (0, "Six hierarchical evaluation modes explored:", True, BLUE),
        (0, "Mode                              Top-1    Top-5    F1-macro    F1=0 classes", True, BLACK),
        (0, "Flat classifier                   16.4 %   46.7 %   0.137       804", False, BLACK),
        (0, "Hard Top-5 genus routing          16.5 %   46.1 %   0.139       ~", False, BLACK),
        (0, "Soft genus routing                16.2 %   46.2 %   0.135       ~", False, BLACK),
        (0, "Oracle genus (upper bound)        27.9 %   65.3 %   0.256       ~", True, GREEN),
        (0, "Per-genus  (predicted genus)      11.2 %   32.9 %   0.109       401", False, BLACK),
        (0, "Per-genus  (oracle genus)         21.8 %   60.2 %   0.204       362", False, BLACK),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 21 — 核心發現二：Genus bottleneck
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Key Finding 3 — Genus is the Species Bottleneck",
    [
        (0, "Oracle (27.9 %) − Flat (16.4 %)  =  + 11.5 pp", True, GREEN),
        (0, "If genus label were always correct, species accuracy gains 11.5 pp", False, BLACK),
        (1, "Approximately half of species errors are carried in from genus errors", False, GRAY),
        (0, "BUT:  Top-k and soft genus routing deliver ≈ 0 improvement  (±0.3 pp)", True, RED),
        (0, "Why is oracle strong yet real routing flat?", True, BLUE),
        (1, "Genus and species model predictions conflict at boundary cases", False, BLACK),
        (1, "Naive combine strategies (top-k mask / prob multiplication) cannot resolve", False, BLACK),
        (0, "Open question for advisor:  better combine strategies?  learned routing?", True, ORANGE),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 22 — 核心發現三：Per-Genus 幫助長尾
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Key Finding 4 — Per-Genus Helps the Long Tail",
    [
        (0, "Per-genus:  one dedicated species classifier per genus (81 models + 39 single-species genera)", False, BLACK),
        (0, "Architecture:  NT-v2 frozen backbone  +  attention head only", False, BLACK),
        (1, "5 M reads divided by genus  →  ~ 3,260 reads / species  (10× less than flat)", False, GRAY),
        (0, "Results:  worse on overall Top-1  (21.8 % oracle vs 27.9 % flat oracle)", False, BLACK),
        (1, "Caused by data being split into small per-genus slices", False, GRAY),
        (0, "BUT:  F1 = 0 species counts", True, BLUE),
        (1, "Flat classifier:           804 / 1,535   ( 52 % )", False, BLACK),
        (1, "Per-genus (predicted):     401 / 1,535   ( 26 % )", True, GREEN),
        (1, "Per-genus (oracle):        362 / 1,535   ( 24 % )", True, GREEN),
        (0, "Per-genus is a legitimate long-tail improvement  —  scale mismatch hides it in aggregate metrics", True, ORANGE),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 23 — 核心發現四：NT-v2 wins at matched tokenization
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Key Finding 5 — Pre-training Wins at Matched Tokenization",
    [
        (0, "MetaTransformer (MT):  shallow Transformer trained from scratch, no pre-training", False, GRAY),
        (0, "At identical 6-mer tokenization, our NT-v2 + LoRA wins decisively:", True, BLUE),
        (0, "Model                           Tokenization                    Genus Acc    Δ", True, BLACK),
        (0, "NT-v2 + LoRA  (ours, RC TTA)   6-mer non-overlap (stride = 6)   67.07 %     —", True, GREEN),
        (0, "MetaTransformer                 6-mer non-overlap (stride = 6)   ~ 36 %      − 31 pp", True, RED),
        (0, "MetaTransformer                 6-mer overlap (stride = 1)       48.87 %     − 18 pp", True, RED),
        (0, "Three conclusions:", True, BLUE),
        (1, "Pre-training representation quality is real  (+31 pp at matched k)", False, BLACK),
        (1, "6× more tokens (overlap) still lags NT-v2 by 18 pp  —  density ≠ quality", False, BLACK),
        (1, "Our pipeline (NT-v2 + LoRA + AttentionPool + RC TTA) is SOTA at sensible token budget", True, GREEN),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 24 — MT 高 k setting vocab cost
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "MetaTransformer's High-k: Vocabulary-for-Accuracy Trade-off",
    [
        (0, "MT surpasses NT-v2 only at k = 12 / 13 overlap, with exploding vocab:", True, BLUE),
        (0, "k     stride            Theoretical vocab (4^k)      Genus Acc", True, BLACK),
        (0, "6     6 (non-overlap)   4,096    (matches NT-v2)      ~ 36 %", False, BLACK),
        (0, "6     1 (overlap)       4,096                         48.87 %", False, BLACK),
        (0, "12    1 (overlap)       16.7 M                        78.83 %", False, BLACK),
        (0, "13    1 (overlap)       67 M                          87.42 %", False, RED),
        (0, "13    13 (non-overlap)  67 M                          27.30 %", False, GRAY),
        (0, "Interpretation:", True, BLUE),
        (1, "NT-v2 vocab:  4,104  →  embedding ≈ 4 M params", False, BLACK),
        (1, "MT k = 13:  theoretical 67 M  →  embedding table ≫ entire NT-v2 backbone", True, RED),
        (1, "High-k MT effectively MEMORIZES the k-mer → label map in embeddings", True, RED),
        (0, "Future-work direction:  GFM pre-trained with high-k overlapping tokenization", True, ORANGE),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 25 — Section: 目前進行中
# ──────────────────────────────────────────────────────────────────────────────
add_section("6.  Ongoing Experiments")

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 26 — 新一輪 tokenization comparison
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "New Tokenization Comparison  (jobs running now)",
    [
        (0, "Addressing advisor's prompt:  can we break past NT-v2's fixed 6-mer tokenization?", False, GRAY),
        (0, "Two new backbones submitted  (5 M genus task):", True, BLUE),
        (0, "DNABERT  (2021)", True, BLACK),
        (1, "zhihan1996/DNA_bert_6  —  86 M params,  BERT-base", False, BLACK),
        (1, "Overlapping 6-mer (stride = 1)  —  k matched to NT-v2, only stride differs", False, BLACK),
        (1, "Goal:  isolate the stride effect under the same k = 6", False, GRAY),
        (0, "DNABERT-2  (2023)", True, BLACK),
        (1, "zhihan1996/DNABERT-2-117M  —  MosaicBERT arch", False, BLACK),
        (1, "BPE (learned) tokenization", False, BLACK),
        (1, "Goal:  test whether learned tokenization beats fixed k-mer", False, GRAY),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 27 — Controlled comparison limitation
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Controlled-Comparison Caveats",
    [
        (0, "Not fully fair — parameter counts differ:", True, BLUE),
        (0, "Model        Params     Tokenization                Hidden dim", True, BLACK),
        (0, "NT-v2        498 M      non-overlap 6-mer           1024", False, BLACK),
        (0, "DNABERT      86 M       overlap 6-mer               768", False, BLACK),
        (0, "DNABERT-2    117 M      BPE                         768", False, BLACK),
        (0, "5.8 × parameter-count gap  —  but outcomes still interpretable:", True, ORANGE),
        (1, "If DNABERT (86 M) > NT-v2 (498 M):  stride effect dominates model scale", False, BLACK),
        (1, "If DNABERT-2 (BPE) > DNABERT (6-mer):  learned > fixed k-mer", False, BLACK),
        (0, "Thesis will label parameter-count differences explicitly", False, GRAY),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 28 — 其他已完成實驗
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Other Completed Experiments in the Thesis",
    [
        (0, "Shallow Transformer ablation (v11)", True, BLUE),
        (1, "Confirms pre-training value  (v9 − v11 gap ≈ 31 pp at matched k = 6)", False, BLACK),
        (0, "Logit Adjustment / RC Consistency sweeps", True, BLUE),
        (1, "Negative results  —  included to show rigorous ablation, supports 'data ≫ tricks'", False, BLACK),
        (0, "Per-genus pipeline end-to-end", True, BLUE),
        (1, "81 dedicated species models trained, evaluated, analyzed for long-tail", False, BLACK),
        (0, "RC TTA theoretical justification", True, BLUE),
        (1, "Thesis §2.7  —  strand-symmetry derivation explaining the +1 ~ 2 pp", False, BLACK),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 29 — Section: 困難
# ──────────────────────────────────────────────────────────────────────────────
add_section("7.  Challenges & Open Questions")

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 30 — 技術性挑戰
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Technical Challenges (and how we solved them)",
    [
        (0, "PyTorch 2.6 breaking change", True, BLUE),
        (1, "weights_only=True became default;  numpy scalars in ckpt → crash", False, BLACK),
        (1, "Fix:  all torch.load calls now pass weights_only=False", False, GRAY),
        (0, "TWCC SLURM 48h timeout", True, BLUE),
        (1, "v9 / sp_v4 both hit the wall  →  auto-resume script checks last.pt", False, BLACK),
        (0, "Eval OOM at 128 G memory", True, BLUE),
        (1, "5 M val × 1,535 species logits ≈ 30 GB  +  buffers  →  bumped to 200 G (TWCC cap)", False, BLACK),
        (0, "DNABERT trust_remote_code conflict", True, BLUE),
        (1, "Custom BertConfig collided with stock  →  made trust_remote_code a config field", False, BLACK),
        (0, "Per-genus sequential training used 5 % GPU", True, BLUE),
        (1, "Parallelized  (N = 4 concurrent)  →  3h  →  30 min", False, BLACK),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 31 — 方法論挑戰
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Methodological Challenges",
    [
        (0, "Fair comparison under 5.8× parameter-count gap", True, BLUE),
        (1, "How to make a tokenization claim when model sizes differ?", False, BLACK),
        (0, "Hierarchical routing strategy", True, BLUE),
        (1, "Oracle gap is 11.5 pp but real top-k routing delivers ~ 0  —  what combines?", False, BLACK),
        (0, "Long-tail classification trade-off", True, BLUE),
        (1, "Per-genus boosts tail F1 but drops aggregate Top-1  —  which is thesis-worthy?", False, BLACK),
        (0, "Ablation design under time budget", True, BLUE),
        (1, "Orthogonal single-variable tests vs combinatorial matrix — which gives cleaner story?", False, BLACK),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 32 — Section: 開放問題
# ──────────────────────────────────────────────────────────────────────────────
add_section("Open Questions — For Discussion")

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 33 — 開放問題清單
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Open Questions for Today's Discussion",
    [
        (0, "1.  High-k + pre-training GFM — worth pursuing?", True, BLUE),
        (1, "Matched-k:  NT-v2 beats MT by +31 pp.  High-k MT embedding-memorization still hits 87 %.", False, BLACK),
        (1, "Combining both remains an open, unexplored direction in the literature", False, GRAY),
        (0, "2.  DNABERT / DNABERT-2 5M results  (running)", True, BLUE),
        (1, "Will decide whether the 'stride' or 'BPE' variable delivers real gains", False, BLACK),
        (0, "3.  Per-genus vs top-k routing — which direction to commit to?", True, BLUE),
        (1, "Oracle says top-k;  implementation says per-genus (long-tail)", False, BLACK),
        (0, "4.  RC-equivariant architectures  (Caduceus)  — worth trying?", True, BLUE),
        (1, "Would eliminate the 2× inference cost of RC TTA", False, BLACK),
        (0, "5.  METAGENE-1 (7 B metagenomic)  — feasible as an upper-bound reference?", True, BLUE),
        (0, "6.  Thesis narrative framing:", True, RED),
        (1, "'Pre-trained GFM is SOTA at matched setting,  and tokenization is an under-appreciated axis'", True, BLACK),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 34 — 目前專案狀態
# ──────────────────────────────────────────────────────────────────────────────
add_content(
    "Current Project Status",
    [
        (0, "Code", True, BLUE),
        (1, "/work/ymj1123ntu/token_level_gfm_classifier    (pending git init)", False, BLACK),
        (0, "Thesis", True, BLUE),
        (1, "/work/ymj1123ntu/thesis    (Chapter 2 / 4 / 5 ≈ 75 % complete)", False, BLACK),
        (0, "Running jobs", True, BLUE),
        (1, "174794  —  sp_v4 species full eval       (~ 1 – 2 h remaining)", False, BLACK),
        (1, "175206  —  DNABERT 5M genus training     (queued)", False, BLACK),
        (1, "175207  —  DNABERT-2 5M genus training   (queued)", False, BLACK),
        (0, "Next milestone", True, BLUE),
        (1, "All pending experiments complete in ~ 10 – 12 h", False, BLACK),
        (1, "Thesis Chapter 4 update with tokenization comparison + final species numbers", False, BLACK),
    ]
)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 35 — Thank you
# ──────────────────────────────────────────────────────────────────────────────
add_section("Thank you   |   Q & A")


prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides.__iter__.__doc__) if False else ''})")
print(f"Total slides: {len(list(prs.slides))}")

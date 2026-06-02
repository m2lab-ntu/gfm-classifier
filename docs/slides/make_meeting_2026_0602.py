#!/usr/bin/env python3
"""
Weekly meeting deck · 2026-05-31
Covers progress 2026-05-28 advisor meeting → 2026-05-30.
"""

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Rectangle
import numpy as np
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DOCS = Path("/work/ymj1123ntu/gfm-classifier/docs/slides")
TMP = DOCS / "meeting_0602_figs"
TMP.mkdir(exist_ok=True)
OUT = DOCS / "meeting_2026_0602.pptx"

NAVY = "#0A2540"
TEAL = "#1E6091"
GREEN = "#2E7D32"
RED = "#C62828"
ORANGE = "#E65100"
GOLD = "#B8860B"
GRAY = "#5A6068"
LIGHT = "#F4F6F8"
WHITE = "#FFFFFF"

def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

plt.rcParams.update({"font.family": "DejaVu Sans", "font.size": 11,
                     "axes.spines.top": False, "axes.spines.right": False,
                     "figure.facecolor": "white", "savefig.dpi": 200,
                     "savefig.bbox": "tight", "savefig.pad_inches": 0.15})


# ═══════════════════════════════════════════════════════════════════════════════
# FIGURES
# ═══════════════════════════════════════════════════════════════════════════════

def fig_progress_timeline():
    fig, ax = plt.subplots(figsize=(12, 3.5))
    ax.set_xlim(0, 12); ax.set_ylim(0, 4); ax.axis("off")

    # Date axis
    ax.add_patch(Rectangle((0.5, 1.9), 11, 0.12, facecolor=NAVY))
    dates = [(0.7,  "5/28",  "Advisor\nmeeting"),
             (2.3,  "5/29",  "MT aligned\n+ inversion\nfound"),
             (3.8,  "5/30",  "Thesis →\n139 pp"),
             (5.2,  "5/30",  "TWCC\nbudget\nexhausted"),
             (6.6,  "6/1",   "GitHub\nrepo +\nExp F (55.93%)"),
             (8.2,  "6/2",   "Nano4 H200\nenv setup"),
             (9.7,  "6/2",   "Nano4 sanity\ncheck FAIL\n(transformers)"),
             (11.0, "now",   "Unblock\nNano4\n→ resume P0")]
    colors = [TEAL, GREEN, GREEN, RED, GREEN, TEAL, ORANGE, GOLD]
    for (x, date, label), color in zip(dates, colors):
        ax.scatter(x, 1.96, s=170, color=color, zorder=4, edgecolor="white", linewidth=2)
        ax.text(x, 2.6, date, ha="center", fontsize=10.5, fontweight="bold", color=color)
        ax.text(x, 1.35, label, ha="center", fontsize=8.5, color=NAVY, va="top")

    ax.text(6, 3.6, "Progress Timeline · 5/28 advisor meeting → 6/2 (today)",
            ha="center", fontsize=13, fontweight="bold", color=NAVY)
    plt.savefig(TMP / "timeline.png"); plt.close()


def fig_mt_alignment_fix():
    """Before/After bar chart: what the alignment fix changed."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4))

    # Left: species level full 100K (replacing old values)
    models = ["MT 13-mer\nflat", "MT 6-mer\nflat", "MT 6-mer\nhier."]
    old = [49.7, 9.2, 6.4]
    new = [53.7, 9.0, 6.4]
    x = np.arange(len(models)); w = 0.36
    ax1.bar(x - w/2, old, w, color=GRAY, label="OLD (on MT val pool)",
            edgecolor="white", linewidth=1.5)
    ax1.bar(x + w/2, new, w, color=GREEN, label="NEW (on TWCC 100K test, aligned)",
            edgecolor="white", linewidth=1.5)
    for i, (o, n) in enumerate(zip(old, new)):
        ax1.text(i - w/2, o + 1, f"{o:.1f}%", ha="center", fontsize=9.5, color=NAVY)
        ax1.text(i + w/2, n + 1, f"{n:.1f}%", ha="center", fontsize=9.5, fontweight="bold", color=GREEN)
    ax1.set_xticks(x); ax1.set_xticklabels(models, fontsize=10)
    ax1.set_ylim(0, 65)
    ax1.set_ylabel("Species read accuracy (%)", fontsize=11)
    ax1.set_title("Species — MT 13-mer flat +4.0 pp on canonical test set",
                  fontsize=11, fontweight="bold", color=NAVY)
    ax1.legend(loc="upper right", fontsize=9)
    ax1.grid(axis="y", linestyle="--", alpha=0.3)

    # Right: genus level (different evaluation conditions)
    models_g = ["MT 13-mer\ngenus", "MT 6-mer\ngenus"]
    old_g = [87.4, 48.8]
    new_g = [94.25, 48.76]
    x = np.arange(len(models_g))
    ax2.bar(x - w/2, old_g, w, color=GRAY, label="OLD (5M val pool, 50K rps)",
            edgecolor="white", linewidth=1.5)
    ax2.bar(x + w/2, new_g, w, color=GREEN, label="NEW (100K test, 1K rps)",
            edgecolor="white", linewidth=1.5)
    for i, (o, n) in enumerate(zip(old_g, new_g)):
        ax2.text(i - w/2, o + 1.5, f"{o:.1f}%", ha="center", fontsize=10, color=NAVY)
        ax2.text(i + w/2, n + 1.5, f"{n:.2f}%", ha="center", fontsize=10, fontweight="bold", color=GREEN)
    ax2.set_xticks(x); ax2.set_xticklabels(models_g, fontsize=10)
    ax2.set_ylim(0, 110)
    ax2.set_ylabel("Genus read accuracy (%)", fontsize=11)
    ax2.set_title("Genus — MT 13-mer +6.85 pp on different test set",
                  fontsize=11, fontweight="bold", color=NAVY)
    ax2.legend(loc="upper right", fontsize=9)
    ax2.grid(axis="y", linestyle="--", alpha=0.3)

    plt.tight_layout()
    plt.savefig(TMP / "mt_alignment.png"); plt.close()


def fig_inversion():
    """THE key new finding: species-vs-genus performance inversion."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.5))

    # Species level (in-DB) — Kraken2 dominates
    models_s = ["NT-Species\nflat", "NT-v2\nper-gen\n(oracle)",
                "MT 13-mer\nflat", "Kraken2"]
    sp_acc = [17.12, 27.69, 52.64, 77.18]
    sp_r   = [0.142, 0.156, 0.514, 0.847]
    colors_s = [TEAL, ORANGE, GREEN, GOLD]
    bars = ax1.bar(models_s, sp_acc, color=colors_s, edgecolor="white", linewidth=2)
    for b, a, r in zip(bars, sp_acc, sp_r):
        ax1.text(b.get_x() + b.get_width()/2, a + 1.5, f"{a:.1f}%",
                 ha="center", fontsize=10.5, fontweight="bold", color=NAVY)
        ax1.text(b.get_x() + b.get_width()/2, a/2, f"r={r:.2f}",
                 ha="center", fontsize=9, color="white")
    ax1.set_ylim(0, 95)
    ax1.set_ylabel("Read Top-1 accuracy (%)", fontsize=11)
    ax1.set_title("Species level · 1,535 classes\nKraken2 dominates ✓",
                  fontsize=12, fontweight="bold", color=GOLD, pad=10)
    ax1.grid(axis="y", linestyle="--", alpha=0.3)
    ax1.text(0.5, 85, "Kraken2 lead: +24 pp acc, +0.33 r",
             ha="center", transform=ax1.transAxes,
             fontsize=10, color=GOLD, fontweight="bold",
             bbox=dict(boxstyle="round,pad=0.4", facecolor="#FFF8E1", edgecolor=GOLD))

    # Genus level (in-DB) — MT 13-mer dominates
    models_g = ["NT-Genus\nv9", "NT-v2\nper-gen\n(oracle)*",
                "MT 13-mer\ngenus", "Kraken2"]
    gn_acc = [68.88, 100.0, 94.89, 77.68]
    gn_r   = [0.992, 1.000, 0.9993, 0.823]
    colors_g = [TEAL, ORANGE, GREEN, GOLD]
    bars = ax2.bar(models_g, gn_acc, color=colors_g, edgecolor="white", linewidth=2)
    for b, a, r in zip(bars, gn_acc, gn_r):
        ax2.text(b.get_x() + b.get_width()/2, a + 1.5, f"{a:.1f}%",
                 ha="center", fontsize=10.5, fontweight="bold", color=NAVY)
        ax2.text(b.get_x() + b.get_width()/2, a/2, f"r={r:.3f}",
                 ha="center", fontsize=9, color="white")
    ax2.set_ylim(0, 115)
    ax2.set_ylabel("Read Top-1 accuracy (%)", fontsize=11)
    ax2.set_title("Genus level · 120 classes\nMT 13-mer surpasses Kraken2 ⚡",
                  fontsize=12, fontweight="bold", color=GREEN, pad=10)
    ax2.grid(axis="y", linestyle="--", alpha=0.3)
    ax2.text(0.5, 102, "MT 13-mer lead: +17 pp acc, +0.18 r",
             ha="center", transform=ax2.transAxes,
             fontsize=10, color=GREEN, fontweight="bold",
             bbox=dict(boxstyle="round,pad=0.4", facecolor="#E8F5E9", edgecolor=GREEN))
    ax2.text(0.5, 0.02, "*oracle = uses true genus label (trivially 100%)",
             ha="center", transform=ax2.transAxes,
             fontsize=8, color=GRAY, style="italic")

    plt.tight_layout()
    plt.savefig(TMP / "inversion.png"); plt.close()


def fig_inversion_explainer():
    """Why the inversion happens: commit-rate × cardinality."""
    fig, ax = plt.subplots(figsize=(11, 4))
    ax.set_xlim(0, 12); ax.set_ylim(0, 6); ax.axis("off")

    # Two boxes side by side
    # Left: species level mechanism
    ax.add_patch(FancyBboxPatch((0.3, 0.5), 5.4, 5.0, boxstyle="round,pad=0.1",
                                 linewidth=2.5, edgecolor=GOLD, facecolor="#FFF8E1"))
    ax.text(3.0, 5.0, "Species level (1,535 classes)", ha="center",
            fontsize=12, fontweight="bold", color=GOLD)
    ax.text(0.6, 4.0, "Kraken2 commits to 70% of reads", fontsize=10.5, color=NAVY)
    ax.text(0.6, 3.5, "• Per-species expected count: 0.65 reads/sample", fontsize=9.5, color=NAVY)
    ax.text(0.6, 3.0, "• Lost 30% spreads across 1,535 species", fontsize=9.5, color=NAVY)
    ax.text(0.6, 2.5, "  → per-species lost signal ≈ 0.2 reads", fontsize=9.5, color=NAVY)
    ax.text(0.6, 2.0, "• On committed reads: 99.33% accurate", fontsize=9.5, color=NAVY)
    ax.text(0.6, 1.0, "Result: high precision on commits wins,",
            fontsize=10.5, fontweight="bold", color=GOLD)
    ax.text(0.6, 0.7, "distributed 30% loss is acceptable", fontsize=10.5, fontweight="bold", color=GOLD)

    # Right: genus level mechanism
    ax.add_patch(FancyBboxPatch((6.3, 0.5), 5.4, 5.0, boxstyle="round,pad=0.1",
                                 linewidth=2.5, edgecolor=GREEN, facecolor="#E8F5E9"))
    ax.text(9.0, 5.0, "Genus level (120 classes)", ha="center",
            fontsize=12, fontweight="bold", color=GREEN)
    ax.text(6.6, 4.0, "Kraken2 still commits to 70%", fontsize=10.5, color=NAVY)
    ax.text(6.6, 3.5, "• Per-genus expected count: 8.3 reads/sample", fontsize=9.5, color=NAVY)
    ax.text(6.6, 3.0, "• Lost 30% spreads across only 120 genera", fontsize=9.5, color=NAVY)
    ax.text(6.6, 2.5, "  → per-genus lost signal ≈ 2.5 reads", fontsize=9.5, color=NAVY)
    ax.text(6.6, 2.0, "• MT 13-mer commits 100%, ~95% accuracy", fontsize=9.5, color=NAVY)
    ax.text(6.6, 1.0, "Result: 30% loss now 30% of signal,",
            fontsize=10.5, fontweight="bold", color=GREEN)
    ax.text(6.6, 0.7, "MT 13-mer's full coverage wins", fontsize=10.5, fontweight="bold", color=GREEN)

    plt.savefig(TMP / "inversion_explainer.png"); plt.close()


def fig_decisions_matrix():
    """6 decision points laid out as cards."""
    fig, ax = plt.subplots(figsize=(12, 6))
    ax.set_xlim(0, 12); ax.set_ylim(0, 7); ax.axis("off")

    decisions = [
        ("Q1", "DNABERT-2 50M resubmit",
         "Move to Nano4 H200 (free) instead of TWCC\n→ STATUS: ready once Nano4 sanity-check unblocks",
         GREEN),
        ("Q2", "DNABERT-1 50M continue?",
         "Cost: 13 days   |   Value: marginal\n→ STATUS: recommended CANCEL (5M + cost-arg suffices)",
         RED),
        ("Q3", "MT 13-mer hier retrain (Taiwana-2)?",
         "Cost: 1-2 days V100   |   Value: fill Table 4.30 row\n→ STATUS: pending advisor decision",
         ORANGE),
        ("Q4", "Speed/memory unified benchmark",
         "Move MT models to Nano4 → H200 unified throughput\n→ STATUS: depends on Nano4 unblock + MT migration",
         GREEN),
        ("Q5", "258M training?",
         "Cost: 18 days wall clock   |   Value: +4.2 pp pred.\n→ STATUS: recommended SKIP (extrapolation in thesis)",
         RED),
        ("Q6", "HMP real-dataset validation",
         "Cost: 3 hr H200 on Nano4   |   Value: must for paper\n→ STATUS: queue after sanity check passes",
         GREEN),
    ]

    positions = [(0.2, 4.5), (4.2, 4.5), (8.2, 4.5),
                 (0.2, 1.0), (4.2, 1.0), (8.2, 1.0)]
    for (qn, title, body, color), (x, y) in zip(decisions, positions):
        ax.add_patch(FancyBboxPatch((x, y), 3.6, 2.5, boxstyle="round,pad=0.1",
                                     linewidth=2.5, edgecolor=color, facecolor=WHITE))
        ax.add_patch(FancyBboxPatch((x, y + 2.0), 3.6, 0.5, boxstyle="round,pad=0.1",
                                     linewidth=0, facecolor=color))
        ax.text(x + 0.3, y + 2.25, qn, fontsize=14, fontweight="bold", color="white", va="center")
        ax.text(x + 1.0, y + 2.25, title, fontsize=11, fontweight="bold",
                color="white", va="center")
        ax.text(x + 0.2, y + 1.0, body, fontsize=9.5, color=NAVY, va="center")

    plt.savefig(TMP / "decisions.png"); plt.close()


def fig_budget_status():
    """TWCC budget out — but Nano4 H200 free 1-month unblocks everything."""
    fig, ax = plt.subplots(figsize=(11.5, 4.2))
    ax.set_xlim(0, 12); ax.set_ylim(0, 5); ax.axis("off")

    # Top: dual-status banner
    ax.add_patch(FancyBboxPatch((0.3, 3.2), 5.6, 1.3, boxstyle="round,pad=0.1",
                                 linewidth=2.5, edgecolor=RED, facecolor="#FFEBEE"))
    ax.text(3.1, 4.15, "⚠ TWCC (Nano5) wallet: -802.5 pts",
            ha="center", fontsize=12, fontweight="bold", color=RED)
    ax.text(3.1, 3.65, "Cannot sbatch new jobs", ha="center", fontsize=10, color=RED)

    ax.add_patch(FancyBboxPatch((6.1, 3.2), 5.6, 1.3, boxstyle="round,pad=0.1",
                                 linewidth=2.5, edgecolor=GREEN, facecolor="#E8F5E9"))
    ax.text(8.9, 4.15, "✓ Nano4 H200: FREE 1 month",
            ha="center", fontsize=12, fontweight="bold", color=GREEN)
    ax.text(8.9, 3.65, "Repo cloned, env up, sanity check blocked on deps",
            ha="center", fontsize=9.5, color=GREEN)

    # Three columns: done / blocked / needs unblock
    cols = [
        ("✓ Done (no GPU needed)", GREEN,
         "• GitHub repo published\n• 4 envs STATUS.md filled\n• Thesis 139 pp w/ Table 4.30\n• Per-genus 13-mer (55.93%)"),
        ("→ In progress on Nano4", TEAL,
         "• Env setup (✓ conda, deps, data)\n• Sanity check (✗ transformers ver.)\n• Need: downgrade transformers\n• Need: sync DNABERT-2 last.pt"),
        ("→ Once Nano4 unblocked", GOLD,
         "• DNABERT-2 50M resume (36 hr)\n• Speed/memory benchmark (3 hr)\n• HMP real-dataset (3 hr)\n• Per-genus 13-mer realign"),
    ]
    for i, (title, color, body) in enumerate(cols):
        x = 0.3 + i * 4.0
        ax.add_patch(FancyBboxPatch((x, 0.2), 3.6, 2.6, boxstyle="round,pad=0.1",
                                     linewidth=2, edgecolor=color, facecolor=WHITE))
        ax.text(x + 1.8, 2.55, title, ha="center", fontsize=11, fontweight="bold", color=color)
        ax.text(x + 0.2, 1.4, body, fontsize=9.5, color=NAVY, va="center")

    plt.savefig(TMP / "budget.png"); plt.close()


# Generate figures
print("Generating figures...")
fig_progress_timeline()
fig_mt_alignment_fix()
fig_inversion()
fig_inversion_explainer()
fig_decisions_matrix()
fig_budget_status()
print(f"  → {TMP}/")


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX
# ═══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW = prs.slide_width
SH = prs.slide_height


def add_header(slide, title, accent=NAVY, page=None, total=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb(accent)
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), SW - Inches(1), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; run = p.add_run(); run.text = title
    run.font.size = Pt(24); run.font.bold = True; run.font.color.rgb = rgb(NAVY)
    if page is not None:
        pn = slide.shapes.add_textbox(SW - Inches(1.2), Inches(7.05),
                                       Inches(1.0), Inches(0.4))
        ptf = pn.text_frame; ptf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        prun = ptf.paragraphs[0].add_run(); prun.text = f"{page} / {total}"
        prun.font.size = Pt(10); prun.font.color.rgb = rgb(GRAY)


def add_textbox(slide, text, left, top, w, h,
                font_size=14, bold=False, color=NAVY, align="left",
                bg=None, border=None):
    if bg or border:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        if bg:
            shp.fill.solid(); shp.fill.fore_color.rgb = rgb(bg)
        else:
            shp.fill.background()
        if border:
            shp.line.color.rgb = rgb(border); shp.line.width = Pt(1.5)
        else:
            shp.line.fill.background()
        tf = shp.text_frame
    else:
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run(); run.text = line
        run.font.size = Pt(font_size); run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_picture(slide, fp, left, top, w=None, h=None):
    if w is not None:
        return slide.shapes.add_picture(str(fp), left, top, width=w)
    elif h is not None:
        return slide.shapes.add_picture(str(fp), left, top, height=h)
    else:
        return slide.shapes.add_picture(str(fp), left, top)


def add_takeaway(slide, text, color=GREEN):
    add_textbox(slide, text, Inches(0.5), SH - Inches(1.1),
                SW - Inches(1), Inches(0.7),
                font_size=14, bold=True, color="#FFFFFF", align="center", bg=color)


def make_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = rgb(WHITE); bg.line.fill.background()
    return s


TOTAL = 10

# ─── SLIDE 1: Title ───────────────────────────────────────────────────────
s = make_slide()
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.8))
bar.fill.solid(); bar.fill.fore_color.rgb = rgb(NAVY); bar.line.fill.background()
add_textbox(s, "Weekly Progress · Updates Since 5/28",
            Inches(0.6), Inches(2.4), Inches(12), Inches(1.0),
            font_size=34, bold=True, color=NAVY, align="left")
add_textbox(s, "Per-genus 13-mer Exp F · GitHub migration · Nano4 H200 (free) brought online",
            Inches(0.6), Inches(3.5), Inches(12), Inches(0.5),
            font_size=20, color=TEAL, align="left")
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.15),
                         Inches(4), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = rgb(TEAL); ln.line.fill.background()
add_textbox(s, "楊明儒 · 2026-06-02",
            Inches(0.6), Inches(4.5), Inches(12), Inches(0.5),
            font_size=16, color=GRAY, align="left")
add_textbox(s,
    "Since 5/28 advisor meeting:\n"
    "  • Earlier this week: MT alignment fix · in-DB Table 4.30 · genus-level inversion · thesis 139 pp · TWCC budget exhausted\n"
    "  • This week (5/31–6/2): Per-genus 13-mer Exp F result · GitHub repo published · Nano4 H200 onboarding\n\n"
    "Today's deck: progress update, blocker status, decision tracker.",
    Inches(0.6), Inches(5.5), Inches(12), Inches(1.5),
    font_size=12.5, color=NAVY, align="left", bg="#F4F6F8", border=TEAL)


# ─── SLIDE 2: Progress Timeline ───────────────────────────────────────────
s = make_slide()
add_header(s, "Progress Timeline · 5/28 → 6/2", accent=TEAL, page=2, total=TOTAL)
add_picture(s, TMP / "timeline.png", Inches(0.5), Inches(1.0), w=Inches(12.3))
add_textbox(s,
    "Week 1 (5/28-30): MT aligned (5/6), in-DB Table 4.30, genus inversion finding, thesis 139pp · TWCC wallet exhausted\n"
    "Week 2 (5/31-6/2): GitHub repo (4 envs synced) · Per-genus 13-mer (55.93%, 3rd router data point) · Nano4 onboarded\n\n"
    "Current blocker: Nano4 sanity check fails on transformers version (NT-v2 ESM model uses removed API).\n"
    "Fix: pip install transformers==4.35.2 in Nano4 gfm env. ETA ~1 hr to unblock.",
    Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
    font_size=13, color=NAVY, align="left", bg="#F4F6F8", border=NAVY)


# ─── SLIDE 3: MT Alignment Fix ────────────────────────────────────────────
s = make_slide()
add_header(s, "MT Predictions · Read-Order Alignment Fix",
           accent=GREEN, page=3, total=TOTAL)
add_picture(s, TMP / "mt_alignment.png", Inches(0.3), Inches(1.1), w=Inches(12.7))
add_textbox(s,
    "Why old numbers were wrong  ·  Previous MT predictions came from MT's val_dir "
    "(multi-file FASTA, sorted by filename). The read order did NOT match TWCC's canonical "
    "reads_100K.fa. Re-ran on Taiwana-2 with reads in single-file order → labels now "
    "1:1-aligned to NT-v2 / Kraken2 references.",
    Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.4),
    font_size=12, color=NAVY, align="left", bg="#F4F6F8", border=GREEN)
add_takeaway(s,
    "All MT in-DB metrics now correctly computed on the canonical 100K test set",
    color=GREEN)


# ─── SLIDE 4: The Inversion (key finding) ─────────────────────────────────
s = make_slide()
add_header(s, "New Finding · Species/Genus Performance Inversion",
           accent=ORANGE, page=4, total=TOTAL)
add_picture(s, TMP / "inversion.png", Inches(0.2), Inches(1.0), w=Inches(12.9))
add_takeaway(s,
    "Same Kraken2 DB, same test set, same protocol — Kraken2 wins species; MT 13-mer wins genus",
    color=ORANGE)


# ─── SLIDE 5: Why the Inversion ──────────────────────────────────────────
s = make_slide()
add_header(s, "Why the Inversion · Commit-Rate × Task Cardinality",
           accent=ORANGE, page=5, total=TOTAL)
add_picture(s, TMP / "inversion_explainer.png", Inches(0.3), Inches(1.0), w=Inches(12.7))
add_textbox(s,
    "Theoretical implication  ·  Kraken2's 30% abstention is a precision-vs-recall trade-off. "
    "At high task granularity (1,535 species), the precision wins. At low granularity (120 genera), "
    "the recall loss dominates. This adds nuance to the thesis's complementary positioning: "
    "they are not just \"in-DB vs OOD\" — they trade off across task cardinality.",
    Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.0),
    font_size=11.5, color=NAVY, align="left", bg="#FFF3E0", border=ORANGE)


# ─── SLIDE 6: Thesis updates summary ──────────────────────────────────────
s = make_slide()
add_header(s, "Thesis Updates Summary",
           accent=TEAL, page=6, total=TOTAL)
# Two-column layout: what changed | what was added
add_textbox(s, "What changed", Inches(0.5), Inches(1.0), Inches(6.0), Inches(0.5),
            font_size=16, bold=True, color=TEAL)
add_textbox(s,
    "• Table 4.29 MT rows updated to aligned values\n"
    "    — MT 13-mer flat: 49.7% → 53.7% / r 0.46 → 0.50\n"
    "    — MT 6-mer flat/hier: minor (9.0% / 6.4%)\n"
    "    — MT 13-mer hier: N/A (checkpoint corrupted)\n\n"
    "• Hierarchical-masking paragraph: removed broken MT 13-mer hier reference;\n"
    "    revised to 2 routers (NT-Genus 66%, MT 6-mer 49%) + reference to read-level 3 points\n\n"
    "• OOD framing fixed: was \"neural retains OOD signal\";\n"
    "    now \"deployment coverage asymmetry\" (the 219 species are in-distribution)\n\n"
    "• Thesis page count: 137 → 139",
    Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2),
    font_size=11, color=NAVY, align="left", bg="#F4F6F8", border=TEAL)

add_textbox(s, "What was added", Inches(6.8), Inches(1.0), Inches(6.0), Inches(0.5),
            font_size=16, bold=True, color=GREEN)
add_textbox(s,
    "• New Table 4.30: Fair-restricted comparison\n"
    "    on in-DB 85,819 reads (species + genus panels)\n\n"
    "• New paragraph: \"Coverage-restricted fair comparison\"\n"
    "    explains why restrict to 85,819 reads + interpretation\n\n"
    "• New paragraph: \"Sharp inversion between species\n"
    "    and genus levels\" — describes the new finding\n"
    "    with commit-rate × cardinality mechanism\n\n"
    "• New paragraph: \"Kraken2 dependence on database\n"
    "    coverage\" — corrected deployment vs OOD framing",
    Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.2),
    font_size=11, color=NAVY, align="left", bg="#E8F5E9", border=GREEN)


# ─── SLIDE 7: Budget situation ────────────────────────────────────────────
s = make_slide()
add_header(s, "Blocker · TWCC Budget Exhausted on 5/30",
           accent=RED, page=7, total=TOTAL)
add_picture(s, TMP / "budget.png", Inches(0.3), Inches(1.1), w=Inches(12.7))
add_takeaway(s,
    "Need decision: budget top-up, or accept partial DNABERT-2 + defer remaining tasks",
    color=RED)


# ─── SLIDE 8: 6 Decision Points status ────────────────────────────────────
s = make_slide()
add_header(s, "6 Open Decisions · Status Update",
           accent=GOLD, page=8, total=TOTAL)
add_picture(s, TMP / "decisions.png", Inches(0.3), Inches(1.0), w=Inches(12.7))


# ─── SLIDE 9: Recommended priorities ──────────────────────────────────────
s = make_slide()
add_header(s, "Recommended Priorities · Next 2 Weeks",
           accent=GREEN, page=9, total=TOTAL)

add_textbox(s,
    "Immediate (CPU / no GPU needed)",
    Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
    font_size=16, bold=True, color=GREEN)
add_textbox(s,
    "• Realign Per-genus 13-mer (Exp F) on Taiwana-2 → integrate in-DB row of Table 4.30\n"
    "• Thesis figure regeneration with new aligned MT numbers (53.7%, 94.25%)\n"
    "• Thesis Section 4.13 proofreading + format consistency check\n"
    "• Optional: MT 13-mer hier retraining on Taiwana-2 (Q3)",
    Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.6),
    font_size=12, color=NAVY, align="left", bg="#E8F5E9", border=GREEN)

add_textbox(s,
    "After Nano4 sanity check unblocks (~45 GPU-hr on H200, FREE)",
    Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
    font_size=16, bold=True, color=GOLD)
add_textbox(s,
    "• Sync DNABERT-2 50M last.pt → resume training (36 hr H200 from epoch 17)\n"
    "• Migrate MT model checkpoints from Taiwana-2 → Nano4\n"
    "• Run speed/memory unified benchmark on H200 (3 hr) — must-do, advisor ask\n"
    "• HMP mock community real-dataset inference (3 hr) — paper submission prep",
    Inches(0.5), Inches(3.8), Inches(12.3), Inches(1.6),
    font_size=12, color=NAVY, align="left", bg="#FFF8E1", border=GOLD)

add_textbox(s,
    "Recommended NOT to do (advisor confirm if disagree)",
    Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
    font_size=16, bold=True, color=RED)
add_textbox(s,
    "• DNABERT-1 50M continue: ~300 GPU-hr / 13 days — 5M evidence + 11.7 hr/epoch cost argument suffices\n"
    "• 258M training: ~290 GPU-hr / 18 days — log-fit predicts +4.2 pp, does not cross tokenisation gap",
    Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.1),
    font_size=12, color=NAVY, align="left", bg="#FFEBEE", border=RED)


# ─── SLIDE 10: Ask + Questions ────────────────────────────────────────────
s = make_slide()
add_header(s, "Asks for Today's Meeting",
           accent=GOLD, page=10, total=TOTAL)
add_textbox(s,
    "Three decisions I need from you today:",
    Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.6),
    font_size=18, bold=True, color=NAVY)

qs = [
    ("1", "Migration plan: Nano4 takes over",
     "Confirm Nano4 (H200, free 1-month) replaces TWCC as primary GPU. "
     "DNABERT-2 resumes on Nano4. TWCC dormant. Speed benchmark on H200. OK?",
     GREEN),
    ("2", "DNABERT-1 50M — final call",
     "Cancel. Use 5M acc (61.78%) + 11.7 hr/epoch cost vs NT-v2's 2.8 hr "
     "as deployment-limitation evidence in thesis. Agreed?",
     RED),
    ("3", "Per-genus 13-mer Exp F integration",
     "Once Taiwana-2 sends realigned predictions, add 4th row to Table 4.30 "
     "showing 55.93% (3rd router-quality data point validating monotonic theorem). OK?",
     GOLD),
]
for i, (n, title, body, color) in enumerate(qs):
    y = Inches(1.8 + i * 1.7)
    add_textbox(s, n, Inches(0.5), y, Inches(0.7), Inches(0.7),
                font_size=24, bold=True, color="#FFFFFF", align="center",
                bg=color)
    add_textbox(s, title, Inches(1.3), y + Inches(0.05), Inches(3.5), Inches(0.6),
                font_size=14, bold=True, color=color)
    add_textbox(s, body, Inches(4.9), y, Inches(8.0), Inches(1.4),
                font_size=11.5, color=NAVY, align="left", bg="#F4F6F8", border=color)

prs.save(str(OUT))
print(f"\n✅ Saved: {OUT}")
print(f"   Slides: {len(prs.slides)}")
